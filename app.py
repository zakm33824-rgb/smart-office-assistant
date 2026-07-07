"""轻量化智能办公助手。"""

from __future__ import annotations

import io
import os
import base64
import subprocess
import sys
import json
import math
import re
import shutil
import tempfile
import uuid
from collections import Counter
from contextlib import nullcontext
from datetime import date
from pathlib import Path
from typing import Any, Iterable

import pandas as pd
import streamlit as st

from ppt_template_library import app_ui as ppt_app_ui
from ppt_template_library import catalog as ppt_catalog
from ppt_template_library import manager as ppt_manager
from ppt_template_library import page_engine as ppt_page_engine
from ppt_template_library import page_library as ppt_page_library
from ppt_template_library import seeds as ppt_seeds
from ppt_template_library import smart_planner as ppt_smart_planner


def _maybe_bootstrap_space_port() -> None:
    if os.environ.get("DEPLOY_PORT_SHIM_CHILD") == "1":
        return
    if os.environ.get("STREAMLIT_SERVER_PORT") == "8501":
        return
    if not any(os.environ.get(name) for name in ("SPACE_ID", "SPACE_HOST", "HF_SPACE_ID")):
        return
    if os.environ.get("_OFFICE_SPACE_CHILD_STARTED") == "1":
        return

    child_env = os.environ.copy()
    child_env["DEPLOY_PORT_SHIM_CHILD"] = "1"
    child_env["_OFFICE_SPACE_CHILD_STARTED"] = "1"
    child_env["STREAMLIT_SERVER_ADDRESS"] = "0.0.0.0"
    child_env["STREAMLIT_SERVER_PORT"] = "8501"
    child_env["STREAMLIT_SERVER_HEADLESS"] = "true"
    child_env["STREAMLIT_BROWSER_GATHER_USAGE_STATS"] = "false"

    subprocess.Popen(
        [
            sys.executable,
            "-m",
            "streamlit",
            "run",
            str(Path(__file__).resolve()),
            "--server.address=0.0.0.0",
            "--server.port=8501",
            "--server.headless=true",
        ],
        env=child_env,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        close_fds=True,
    )
    st.stop()


_maybe_bootstrap_space_port()

from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from ppt_template_library.app_ui import render_template_library_page
from ppt_template_library.page_engine import build_library_overview, build_page_level_outline

APP_TITLE = '轻量化智能办公助手'
DEFAULT_PPT_PROMPT = '生成一份Q3产品运营复盘PPT，包含业绩概况、用户增长、问题分析、下月规划4个章节'
DEFAULT_TABLE_PROMPT = '生成2024年各门店营收表，包含门店名、月度营收、同比增速、排名4列，填充合理示例数据'

STYLE_UI_TO_CODE = {
    '\u7b80\u7ea6\u5546\u52a1': 'business',
    '\u5b66\u672f\u6c47\u62a5': 'academic',
    '\u521b\u610f\u6f14\u793a': 'creative',
    '\u79d1\u6280\u84dd': 'technology',
    '\u91d1\u878d\u7eff': 'finance',
    '\u653f\u52a1\u7ea2': 'government',
    '\u6781\u7b80\u767d': 'minimal',
    '\u9ed1\u91d1\u9ad8\u7aef': 'dark',
    '\u53ef\u6301\u7eed\u62a5\u544a': 'eco_report',
    '\u62fc\u63a5\u4e2a\u4eba\u9875': 'split_portfolio',
    '\u9ed1\u767d\u6587\u6848': 'editorial_copy',
}
STYLE_CODE_TO_UI = {
    'business': '\u7b80\u7ea6\u5546\u52a1',
    'academic': '\u5b66\u672f\u6c47\u62a5',
    'creative': '\u521b\u610f\u6f14\u793a',
    'technology': '\u79d1\u6280\u84dd',
    'finance': '\u91d1\u878d\u7eff',
    'government': '\u653f\u52a1\u7ea2',
    'minimal': '\u6781\u7b80\u767d',
    'dark': '\u9ed1\u91d1\u9ad8\u7aef',
    'eco_report': '\u53ef\u6301\u7eed\u62a5\u544a',
    'split_portfolio': '\u62fc\u63a5\u4e2a\u4eba\u9875',
    'editorial_copy': '\u9ed1\u767d\u6587\u6848',
}

PPT_STYLES: dict[str, dict[str, RGBColor]] = {
    '\u7b80\u7ea6\u5546\u52a1': {'bg': RGBColor(247, 249, 252), 'panel': RGBColor(255, 255, 255), 'primary': RGBColor(32, 83, 117), 'accent': RGBColor(0, 152, 138), 'text': RGBColor(38, 45, 52), 'muted': RGBColor(99, 110, 123), 'light': RGBColor(228, 238, 244)},
    '\u5b66\u672f\u6c47\u62a5': {'bg': RGBColor(248, 248, 245), 'panel': RGBColor(255, 255, 255), 'primary': RGBColor(74, 66, 124), 'accent': RGBColor(153, 111, 44), 'text': RGBColor(42, 42, 46), 'muted': RGBColor(105, 103, 98), 'light': RGBColor(235, 232, 244)},
    '\u521b\u610f\u6f14\u793a': {'bg': RGBColor(250, 250, 255), 'panel': RGBColor(255, 255, 255), 'primary': RGBColor(31, 76, 153), 'accent': RGBColor(228, 98, 74), 'text': RGBColor(32, 39, 52), 'muted': RGBColor(94, 103, 120), 'light': RGBColor(235, 241, 255)},
    '\u79d1\u6280\u84dd': {'bg': RGBColor(244, 248, 255), 'panel': RGBColor(255, 255, 255), 'primary': RGBColor(18, 60, 122), 'accent': RGBColor(47, 128, 237), 'text': RGBColor(21, 28, 43), 'muted': RGBColor(102, 112, 133), 'light': RGBColor(225, 235, 250)},
    '\u91d1\u878d\u7eff': {'bg': RGBColor(245, 251, 247), 'panel': RGBColor(255, 255, 255), 'primary': RGBColor(28, 108, 79), 'accent': RGBColor(59, 179, 113), 'text': RGBColor(27, 40, 35), 'muted': RGBColor(92, 123, 111), 'light': RGBColor(226, 244, 233)},
    '\u653f\u52a1\u7ea2': {'bg': RGBColor(253, 247, 247), 'panel': RGBColor(255, 255, 255), 'primary': RGBColor(139, 30, 45), 'accent': RGBColor(201, 54, 62), 'text': RGBColor(42, 29, 29), 'muted': RGBColor(139, 110, 110), 'light': RGBColor(248, 228, 230)},
    '\u6781\u7b80\u767d': {'bg': RGBColor(249, 250, 251), 'panel': RGBColor(255, 255, 255), 'primary': RGBColor(17, 24, 39), 'accent': RGBColor(59, 130, 246), 'text': RGBColor(17, 24, 39), 'muted': RGBColor(107, 114, 128), 'light': RGBColor(229, 231, 235)},
    '\u9ed1\u91d1\u9ad8\u7aef': {'bg': RGBColor(15, 17, 21), 'panel': RGBColor(23, 26, 33), 'primary': RGBColor(229, 195, 106), 'accent': RGBColor(245, 166, 35), 'text': RGBColor(245, 247, 250), 'muted': RGBColor(167, 176, 191), 'light': RGBColor(35, 38, 47)},
    '\u53ef\u6301\u7eed\u62a5\u544a': {'bg': RGBColor(250, 245, 226), 'panel': RGBColor(255, 255, 252), 'primary': RGBColor(62, 138, 148), 'accent': RGBColor(245, 190, 66), 'text': RGBColor(44, 66, 70), 'muted': RGBColor(98, 121, 121), 'light': RGBColor(220, 242, 238)},
    '\u62fc\u63a5\u4e2a\u4eba\u9875': {'bg': RGBColor(4, 64, 92), 'panel': RGBColor(255, 255, 255), 'primary': RGBColor(255, 247, 238), 'accent': RGBColor(255, 94, 58), 'text': RGBColor(245, 248, 250), 'muted': RGBColor(198, 210, 218), 'light': RGBColor(12, 84, 117)},
    '\u9ed1\u767d\u6587\u6848': {'bg': RGBColor(230, 227, 214), 'panel': RGBColor(244, 241, 231), 'primary': RGBColor(44, 38, 82), 'accent': RGBColor(61, 55, 122), 'text': RGBColor(28, 28, 31), 'muted': RGBColor(75, 73, 84), 'light': RGBColor(236, 233, 225)},
}
HEX_STYLES: dict[str, dict[str, str]] = {
    '\u7b80\u7ea6\u5546\u52a1': {'bg': 'F7F9FC', 'panel': 'FFFFFF', 'primary': '205375', 'accent': '00988A', 'text': '262D34', 'muted': '636E7B', 'light': 'E4EEF4'},
    '\u5b66\u672f\u6c47\u62a5': {'bg': 'F8F8F5', 'panel': 'FFFFFF', 'primary': '4A427C', 'accent': '996F2C', 'text': '2A2A2E', 'muted': '69675E', 'light': 'EBE8F4'},
    '\u521b\u610f\u6f14\u793a': {'bg': 'FAFAFF', 'panel': 'FFFFFF', 'primary': '1F4C99', 'accent': 'E4624A', 'text': '202734', 'muted': '5E6778', 'light': 'EBF1FF'},
    '\u79d1\u6280\u84dd': {'bg': 'F4F8FF', 'panel': 'FFFFFF', 'primary': '123C7A', 'accent': '2F80ED', 'text': '151C2B', 'muted': '667085', 'light': 'E1EBFA'},
    '\u91d1\u878d\u7eff': {'bg': 'F5FBF7', 'panel': 'FFFFFF', 'primary': '1C6C4F', 'accent': '3BB371', 'text': '1B2823', 'muted': '5C7B6F', 'light': 'E2F4EA'},
    '\u653f\u52a1\u7ea2': {'bg': 'FDF7F7', 'panel': 'FFFFFF', 'primary': '8B1E2D', 'accent': 'C9363E', 'text': '2A1D1D', 'muted': '8B6E6E', 'light': 'F8E4E6'},
    '\u6781\u7b80\u767d': {'bg': 'F9FAFB', 'panel': 'FFFFFF', 'primary': '111827', 'accent': '3B82F6', 'text': '111827', 'muted': '6B7280', 'light': 'E5E7EB'},
    '\u9ed1\u91d1\u9ad8\u7aef': {'bg': '0F1115', 'panel': '171A21', 'primary': 'E5C36A', 'accent': 'F5A623', 'text': 'F5F7FA', 'muted': 'A7B0BF', 'light': '23262F'},
    '\u53ef\u6301\u7eed\u62a5\u544a': {'bg': 'FAF5E2', 'panel': 'FFFEFC', 'primary': '3E8A94', 'accent': 'F5BE42', 'text': '2C4246', 'muted': '627979', 'light': 'DCF2EE'},
    '\u62fc\u63a5\u4e2a\u4eba\u9875': {'bg': '04405C', 'panel': 'FFFFFF', 'primary': 'FFF7EE', 'accent': 'FF5E3A', 'text': 'F5F8FA', 'muted': 'C6D2DA', 'light': '0C5475'},
    '\u9ed1\u767d\u6587\u6848': {'bg': 'E6E3D6', 'panel': 'F4F1E7', 'primary': '2C2652', 'accent': '3D377A', 'text': '1C1C1F', 'muted': '4B4954', 'light': 'ECE9E1'},
}
STORE_NAMES = ['星河旗舰店', '晨光中心店', '云栖体验店', '远山旗舰店', '海韵门店', '启航门店', '未来门店', '明日门店', '极光门店', '城市会客厅']
REGIONS = ['华东', '华北', '华南', '西南', '西北', '华中', '东北', '海外']
CITIES = ['上海', '北京', '广州', '深圳', '杭州', '成都', '武汉', '西安', '南京', '重庆']
PRODUCTS = ['A款', 'B款', 'C款', 'D款', 'E款', 'F款']
DEPARTMENTS = ['销售部', '市场部', '运营部', '研发部', '财务部', '人力资源部']
PEOPLE = ['张伟', '李娜', '王强', '赵敏', '刘洋', '陈晨', '周婷', '黄磊']
MONTHS = [f'{month}月' for month in range(1, 13)]


def setup_page() -> None:
    st.set_page_config(page_title=APP_TITLE, page_icon='📊', layout='wide')
    st.markdown(
        '''
        <style>
        html, body, [data-testid="stAppViewContainer"] {height: 100%; background: transparent !important;}
        [data-testid="stAppViewContainer"] > .main {background: transparent !important;}
        [data-testid="stHeader"], [data-testid="stToolbar"], .stDeployButton, #MainMenu, footer {display: none !important;}
        .stApp {background: transparent !important;}
        .block-container {padding-top: 1.2rem; padding-bottom: 2rem;}
        [data-testid="stSidebar"] {background: rgba(246, 248, 251, 0.92);}
        h1, h2, h3 {letter-spacing: 0 !important;}
        .office-card {border: 1px solid #e6ebf2; border-radius: 8px; padding: 16px 18px; background: rgba(255, 255, 255, 0.94); margin-bottom: 12px;}
        .small-muted {color: #667085; font-size: 0.92rem;}
        .metric-pill {display: inline-block; padding: 4px 10px; border-radius: 999px; background: #eef5ff; color: #245a8d; font-size: 0.85rem; margin-right: 6px; margin-bottom: 6px;}
        .office-video-bg {
            position: fixed;
            inset: 0;
            width: 100vw;
            height: 100vh;
            object-fit: cover;
            z-index: -2;
            pointer-events: none;
            filter: saturate(0.95) contrast(1.02) brightness(0.78);
        }
        .office-video-overlay {
            position: fixed;
            inset: 0;
            z-index: -1;
            pointer-events: none;
            background: linear-gradient(180deg, rgba(7, 10, 18, 0.42) 0%, rgba(10, 14, 24, 0.32) 100%);
        }
        </style>
        ''',
        unsafe_allow_html=True,
    )
    background_path = Path(__file__).resolve().parent / 'assets' / 'background.mp4'
    if background_path.exists():
        video_b64 = base64.b64encode(background_path.read_bytes()).decode('ascii')
        st.markdown(
            f'''
            <video class="office-video-bg" autoplay loop muted playsinline>
                <source src="data:video/mp4;base64,{video_b64}" type="video/mp4">
            </video>
            <div class="office-video-overlay"></div>
            ''',
            unsafe_allow_html=True,
        )


def normalize_text(text: str) -> str:
    return re.sub(r'\s+', ' ', (text or '').strip())


def safe_filename(name: str, suffix: str) -> str:
    cleaned = re.sub(r'[\\/:*?"<>|\s]+', '_', name.strip())[:48].strip('_')
    return f"{cleaned or '智能办公结果'}{suffix}"


def resolve_style(style_name: str) -> str:
    if style_name in PPT_STYLES:
        return style_name
    return STYLE_CODE_TO_UI.get(style_name, '简约商务')


def style_code(style_name: str) -> str:
    if style_name in STYLE_UI_TO_CODE:
        return STYLE_UI_TO_CODE[style_name]
    if style_name in STYLE_CODE_TO_UI:
        return style_name
    return 'business'


def theme(style_name: str) -> dict[str, RGBColor]:
    return PPT_STYLES[resolve_style(style_name)]


def theme_hex(style_name: str) -> dict[str, str]:
    return HEX_STYLES[resolve_style(style_name)]


def add_background(slide: Any, palette: dict[str, RGBColor]) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = palette['bg']


def set_shape_fill(shape: Any, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_shape_line(shape: Any, color: RGBColor | None = None, transparency: int = 100000) -> None:
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = color
    shape.line.transparency = transparency


def add_text(slide: Any, text: str, left: float, top: float, width: float, height: float, font_size: int = 20, color: RGBColor | None = None, bold: bool = False, align: PP_ALIGN = PP_ALIGN.LEFT, font_name: str = 'Microsoft YaHei') -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def add_bullets(slide: Any, bullets: list[str], left: float, top: float, width: float, height: float, color: RGBColor, font_size: int = 18) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f'• {bullet}'
        p.level = 0
        p.space_after = Pt(8)
        p.line_spacing = 1.1
        run = p.runs[0]
        run.font.name = 'Microsoft YaHei'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def split_items(raw_text: str) -> list[str]:
    cleaned = normalize_text(raw_text)
    cleaned = re.sub(r'\d+\s*个章节', '', cleaned)
    cleaned = re.sub(r'\d+\s*列', '', cleaned)
    cleaned = cleaned.replace('以及', '、').replace('和', '、').replace('与', '、')
    parts = re.split(r'[、,，;；/|]+', cleaned)
    return [item.strip(' .、,；;/|') for item in parts if item.strip(' .、,；;/|')]


def extract_explicit_columns(prompt: str) -> list[str]:
    text = normalize_text(prompt)
    patterns = [
        r'(?:包含|包括|字段|列名|列有|含有)(.+?)(?:\d+\s*列|列|。|；|;|,|，|$)',
        r'(?:列为|列是|需要)(.+?)(?:\d+\s*列|列|。|；|;|,|，|$)',
    ]
    for pattern in patterns:
        match = re.search(pattern, text)
        if match:
            columns = split_items(match.group(1))
            columns = [col for col in columns if len(col) <= 24]
            if columns:
                return columns
    return []


def infer_row_count(prompt: str, default: int = 8) -> int:
    text = normalize_text(prompt)
    match = re.search(r'(\d+)\s*(?:行|条|家|项|个样本|条数据)', text)
    if match:
        return max(3, min(30, int(match.group(1))))
    if '前十' in text or '10' in text and '排名' in text:
        return 10
    return default


def infer_table_title(prompt: str) -> str:
    text = re.split(r'(?:包含|包括|字段|列|并|，|。|；|;)', normalize_text(prompt), maxsplit=1)[0]
    text = re.sub(r'^(请|帮我|麻烦|生成|制作|输出|做一份|做个|帮忙)?', '', text).strip(' ，,。:：')
    return text or '智能表格'


def infer_table_columns(prompt: str) -> list[str]:
    explicit = extract_explicit_columns(prompt)
    if explicit:
        return explicit
    text = normalize_text(prompt)
    if any(word in text for word in ['门店', '店铺', '营收', '销售额', '收入']):
        return ['门店名', '月度营收', '同比增速', '排名']
    if any(word in text for word in ['月份', '月度', '季度', '趋势', '增长']):
        return ['月份', '销售额', '环比增速', '累计销售额']
    if any(word in text for word in ['项目', '进度', '计划']):
        return ['项目', '负责人', '进度', '状态', '截止日期']
    if any(word in text for word in ['员工', '人员', '人力', '团队']):
        return ['姓名', '部门', '岗位', '绩效评分']
    if any(word in text for word in ['学校', '班级', '学生', '成绩']):
        return ['姓名', '班级', '成绩', '排名']
    return ['名称', '指标A', '指标B', '指标C']


def infer_prompt_topic(prompt: str) -> str:
    text = normalize_text(prompt)
    title = re.split(r'(?:包含|包括|围绕|结合|实现|并|，|。|；|;)', text, maxsplit=1)[0]
    title = re.sub(r'^(请|帮我|麻烦|生成|制作|输出|做一份|做个|帮忙)?', '', title).strip(' ，,。:：')
    return title or '智能PPT'


def style_of_text(text: str) -> str:
    lowered = text.lower()
    if any(key in lowered for key in ['科技', 'ai', '互联网', '云', '大数据', '软件', '网络安全']):
        return 'technology'
    if any(key in text for key in ['学术', '教育', '科研', '答辩']):
        return 'academic'
    if any(key in text for key in ['创意', '活泼', '插画', '卡通']):
        return 'creative'
    return 'business'


def _add_tag_row(slide: Any, tags: Iterable[str], left: float, top: float, width: float, color: RGBColor, bg: RGBColor) -> None:
    x = left
    for tag in list(tags)[:8]:
        tag_w = max(0.9, min(1.7, 0.12 * len(str(tag)) + 0.7))
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(tag_w), Inches(0.34))
        set_shape_fill(pill, bg)
        set_shape_line(pill, bg)
        add_text(slide, str(tag), x + 0.06, top + 0.02, tag_w - 0.12, 0.25, font_size=11, color=color, bold=True)
        x += tag_w + 0.12
        if x > left + width - 0.6:
            break


def _presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs



from copy import deepcopy


GENERATED_LIBRARY_DIRS = (
    'templates',
    'components',
    'preview',
    'metadata',
    'database',
    'commercial_allowed',
    'personal_use_only',
    'license_uncertain',
    'premium_quality',
    'industry',
    'scenario',
    'style',
    'color',
    'layout',
    'logs',
    'source_records',
)


def _clear_generated_library_assets(root: str | Path = 'ppt_template_library') -> None:
    root_path = Path(root)
    for name in GENERATED_LIBRARY_DIRS:
        target = root_path / name
        try:
            if target.is_dir():
                shutil.rmtree(target, ignore_errors=True)
            elif target.exists():
                target.unlink()
        except Exception:
            pass


def _normalize_prompt(text: str) -> str:
    return re.sub(r'\s+', ' ', (text or '').strip())


def _prompt_title(prompt: str) -> str:
    text = _normalize_prompt(prompt)
    for prefix in ('\u8bf7', '\u5e2e\u6211', '\u9ebb\u70e6', '\u751f\u6210', '\u5236\u4f5c', '\u8f93\u51fa', '\u505a\u4e00\u4efd', '\u505a\u4e2a', '\u5e2e\u5fd9'):
        if text.startswith(prefix):
            text = text[len(prefix):].lstrip()
            break
    if text.startswith('\u4e00\u4efd'):
        text = text[2:].lstrip()
    for marker in ('\u5305\u542b', '\u5305\u62ec', '\u6db5\u76d6', '\u56f4\u7ed5', '\u7ed3\u5408', '\u5b9e\u73b0', '\u5e76\u4e14', '\u5e76', '\u4ee5\u53ca', '\u3001', '\uff0c', '\u3002', '\uff1b', ';', ':', '\uff1a'):
        if marker in text:
            text = text.split(marker, 1)[0]
            break
    text = text.replace('PPTX', '').replace('pptx', '').replace('PPT', '').replace('ppt', '')
    text = text.replace('\u5e7b\u706f\u7247', '').replace('\u6f14\u793a\u6587\u7a3f', '').replace('\u6c47\u62a5', '').replace('\u62a5\u544a', '')
    return text.strip(' ,.:-_\u3001\u3002\uff1a\uff0c') or '\u667a\u80fdPPT'


def _prompt_sections(prompt: str) -> list[str]:
    text = _normalize_prompt(prompt)
    fallback_map = [
        (['\u590d\u76d8', '\u603b\u7ed3', '\u5e74\u62a5', '\u7ecf\u8425', '\u4e1a\u7ee9', '\u6c47\u62a5'], ['\u4e1a\u7ee9\u6982\u51b5', '\u589e\u957f\u5206\u6790', '\u95ee\u9898\u5206\u6790', '\u884c\u52a8\u8ba1\u5212']),
        (['\u9879\u76ee', '\u8ba1\u5212', '\u8def\u7ebf\u56fe', '\u91cc\u7a0b\u7891', '\u63a8\u8fdb'], ['\u9879\u76ee\u80cc\u666f', '\u63a8\u8fdb\u8fdb\u5ea6', '\u98ce\u9669\u4e0e\u8d44\u6e90', '\u4e0b\u4e00\u6b65\u8ba1\u5212']),
        (['\u4ea7\u54c1', '\u53d1\u5e03', '\u529f\u80fd', '\u4ecb\u7ecd'], ['\u4ea7\u54c1\u80cc\u666f', '\u6838\u5fc3\u4eae\u70b9', '\u5e94\u7528\u573a\u666f', '\u53d1\u5e03\u8ba1\u5212']),
        (['\u5b66\u672f', '\u7814\u7a76', '\u8bba\u6587', '\u7b54\u8fa9'], ['\u7814\u7a76\u80cc\u666f', '\u65b9\u6cd5\u4e0e\u5b9e\u9a8c', '\u7ed3\u679c\u5206\u6790', '\u7ed3\u8bba\u5c55\u671b']),
        (['\u5e02\u573a', '\u8425\u9500', '\u589e\u957f', '\u7528\u6237'], ['\u5e02\u573a\u80cc\u666f', '\u6e20\u9053\u8868\u73b0', '\u7528\u6237\u53cd\u9988', '\u8425\u9500\u8ba1\u5212']),
        (['\u8d22\u52a1', '\u9500\u552e', '\u8fd0\u8425', '\u6570\u636e', '\u6307\u6807'], ['\u6838\u5fc3\u6307\u6807', '\u8d8b\u52bf\u5206\u6790', '\u7ed3\u6784\u62c6\u89e3', '\u7ba1\u7406\u5efa\u8bae']),
    ]
    fallback = ['\u80cc\u666f\u4e0e\u76ee\u6807', '\u6838\u5fc3\u8fdb\u5c55', '\u95ee\u9898\u5206\u6790', '\u884c\u52a8\u8ba1\u5212']
    for keywords, sections in fallback_map:
        if any(keyword in text for keyword in keywords):
            fallback = sections
            break
    extracted = ''
    for marker in ('\u5305\u542b', '\u5305\u62ec', '\u6db5\u76d6', '\u56f4\u7ed5', '\u5206\u4e3a', '\u5206\u6210', '\u5217\u51fa', '\u7ed9\u51fa'):
        if marker in text:
            extracted = text.split(marker, 1)[1]
            break
    if not extracted and ':' in text:
        extracted = text.split(':', 1)[1]
    if not extracted and '\uff1a' in text:
        extracted = text.split('\uff1a', 1)[1]
    for terminator in ('\u4e2a\u7ae0\u8282', '\u7ae0\u8282', '\u90e8\u5206', '\u6a21\u5757', '\u677f\u5757', '\u9875'):
        if terminator in extracted:
            extracted = extracted.split(terminator, 1)[0]
            break
    items: list[str] = []
    if extracted:
        for part in re.split(r'[,\uff0c\u3001;\uff1b\\/|&]+', extracted):
            item = re.sub(r'^\d+[.\u3001\)\]]\s*', '', part).strip(' ,.:-_\u3001\u3002\uff1a\uff0c')
            item = re.sub(r'\d+$', '', item).strip()
            if item and len(item) <= 24:
                items.append(item)
    return list(dict.fromkeys(items)) or fallback


def _seedless_bootstrap_library(*args: Any, **kwargs: Any) -> dict[str, Any]:
    _clear_generated_library_assets()
    return _ORIGINAL_BOOTSTRAP_LIBRARY(*args, **kwargs)


def _prompt_page_type_for_section(section: str) -> tuple[str, str]:
    text = _normalize_prompt(section).lower()
    groups = [
        (("问题", "痛点", "风险", "挑战", "瓶颈", "异常", "problem", "risk", "challenge", "issue", "原因", "改进"), ('comparison', 'before_after')),
        (("计划", "规划", "下月", "下一步", "行动", "路线图", "里程碑", "plan", "roadmap", "milestone"), ('timeline', 'roadmap_timeline')),
        (("用户", "引流", "统计", "增长", "趋势", "流量", "留存", "转化", "growth", "traffic", "retention", "conversion"), ('chart', 'trend_chart')),
        (("业绩", "收入", "销售", "财务", "经营", "指标", "年报", "数据", "kpi", "performance", "revenue", "sales", "financial", "metric"), ('data_analysis', 'kpi_dashboard')),
        (("表格", "清单", "明细", "table", "list"), ('table', 'ranking_table')),
        (("流程", "步骤", "机制", "执行", "process", "workflow", "flow"), ('process', 'workflow_process')),
        (("地图", "区域", "分布", "map", "region", "distribution"), ('map', 'business_map')),
        (("团队", "组织", "成员", "角色", "team", "organization", "people", "profile"), ('team', 'team_profile')),
        (("对比", "比较", "竞品", "前后", "comparison", "before", "after"), ('comparison', 'before_after')),
        (("swot", "pest", "策略", "strategy", "结论", "总结"), ('strategy', 'swot_matrix')),
    ]
    for terms, result in groups:
        if any(term.lower() in text for term in terms):
            return result
    return 'content', 'text_image_layout'


def _prompt_layout_for_section(section: str, data_profile: Any | None = None) -> str:
    text = _normalize_prompt(section).lower()
    groups = [
        (("问题", "痛点", "风险", "挑战", "瓶颈", "异常", "problem", "risk", "challenge", "issue", "原因", "改进"), 'comparison_slide'),
        (("计划", "规划", "下月", "下一步", "行动", "路线图", "里程碑", "plan", "roadmap", "milestone"), 'timeline_slide'),
        (("用户", "引流", "统计", "增长", "趋势", "流量", "留存", "转化", "growth", "traffic", "retention", "conversion"), 'chart_slide'),
        (("业绩", "收入", "销售", "财务", "经营", "指标", "年报", "数据", "kpi", "performance", "revenue", "sales", "financial", "metric"), 'data_analysis_slide'),
        (("表格", "清单", "明细", "table", "list"), 'table_slide'),
        (("流程", "步骤", "机制", "执行", "process", "workflow", "flow"), 'process_slide'),
        (("地图", "区域", "分布", "map", "region", "distribution"), 'map_slide'),
        (("团队", "组织", "成员", "角色", "team", "organization", "people", "profile"), 'people_slide'),
        (("对比", "比较", "竞品", "前后", "comparison", "before", "after"), 'comparison_slide'),
        (("swot", "pest", "策略", "strategy", "结论", "总结"), 'strategy_slide'),
    ]
    for terms, layout in groups:
        if any(term.lower() in text for term in terms):
            return layout
    if data_profile is not None:
        try:
            if getattr(data_profile, 'has_map', False):
                return 'map_slide'
            if getattr(data_profile, 'has_trend', False):
                return 'chart_slide'
            if getattr(data_profile, 'metric_columns', None):
                return 'data_analysis_slide'
        except Exception:
            pass
    return 'content_slide'


def _prompt_bullets_for_section(section: str, prompt: str, style_label: str, index: int, total: int) -> list[str]:
    page_type, _ = _prompt_page_type_for_section(section)
    clean_section = _normalize_prompt(section) or section
    base_map = {
        'data_analysis': [f'拆解{clean_section}的核心指标', '对比同比、环比和目标差异', '提炼主要驱动因素与结论'],
        'kpi_dashboard': [f'围绕{clean_section}建立关键看板', '突出最重要的业务指标', '补充异常和机会点'],
        'chart': [f'呈现{clean_section}的变化趋势', '标记关键拐点和阶段差异', '说明趋势对后续决策的影响'],
        'trend_chart': [f'呈现{clean_section}的变化趋势', '标记关键拐点和阶段差异', '说明趋势对后续决策的影响'],
        'comparison': [f'梳理{clean_section}的现状', '定位关键问题或差异', '给出对应优化动作'],
        'before_after': [f'对比{clean_section}的前后变化', '突出改进效果或差距', '明确下一步建议'],
        'timeline': [f'拆解{clean_section}的阶段目标', '排列关键里程碑与时间点', '说明执行路径与责任人'],
        'roadmap_timeline': [f'拆解{clean_section}的阶段目标', '排列关键里程碑与时间点', '说明执行路径与责任人'],
        'table': [f'列出{clean_section}的结构化信息', '统一字段口径与示例', '补充汇总和排序规则'],
        'ranking_table': [f'列出{clean_section}的结构化信息', '统一字段口径与示例', '补充汇总和排序规则'],
        'process': [f'拆分{clean_section}的关键步骤', '说明各步骤之间的衔接', '强调执行规范与责任分工'],
        'workflow_process': [f'拆分{clean_section}的关键步骤', '说明各步骤之间的衔接', '强调执行规范与责任分工'],
        'strategy': [f'形成{clean_section}的核心判断', '明确优先级与取舍', '落到可执行动作'],
        'swot_matrix': [f'形成{clean_section}的核心判断', '明确优先级与取舍', '落到可执行动作'],
        'map': [f'展示{clean_section}的空间分布', '标出重点区域与差异', '提炼区域策略建议'],
        'business_map': [f'展示{clean_section}的空间分布', '标出重点区域与差异', '提炼区域策略建议'],
        'team': [f'介绍{clean_section}的人员构成', '说明职责分工', '突出关键成员与协作方式'],
        'team_profile': [f'介绍{clean_section}的人员构成', '说明职责分工', '突出关键成员与协作方式'],
        'content': [f'围绕{clean_section}展开说明', '补足背景、要点和结论', '确保表达简洁清晰'],
        'text_image_layout': [f'围绕{clean_section}展开说明', '补足背景、要点和结论', '确保表达简洁清晰'],
    }
    bullets = base_map.get(page_type, [f'围绕{clean_section}展开说明', '补足背景、要点和结论', '确保表达简洁清晰'])
    if bullets and total > 1:
        bullets = bullets.copy()
        bullets[0] = f'{bullets[0]}（第{index}/{total}页，{style_label}）'
    return bullets[:3]


def _prompt_content_titles(prompt: str, sections: list[str], style_label: str) -> list[str]:
    text = _normalize_prompt(prompt).lower()
    titles: list[str] = []
    if any(keyword in text for keyword in ('业绩', '收入', '销售', '财务', '经营', '指标', '增长', '趋势', '用户', '留存', '转化', '数据')):
        titles.extend(['关键数据', '趋势洞察'])
    if any(keyword in text for keyword in ('问题', '风险', '挑战', '瓶颈', '痛点', '异常', '复盘')):
        titles.append('问题与对策')
    if any(keyword in text for keyword in ('计划', '规划', '下一步', '行动', '路线图', '里程碑', '推进')):
        titles.append('推进计划')
    if not titles:
        titles.extend(['核心内容', '总结与行动'])
    return list(dict.fromkeys(titles))[: (2 if len(sections) >= 3 else 1)]


def _prompt_page_sequence(prompt: str, request: Any, page_count: int) -> list[dict[str, Any]]:
    sections = list(dict.fromkeys(_prompt_sections(prompt)))
    slots: list[dict[str, Any]] = [
        {'page_type': 'cover', 'page_subtype': request.style, 'title': request.title},
        {'page_type': 'agenda', 'page_subtype': 'number_agenda', 'title': '目录'},
    ]
    for section in sections[:max(0, page_count - 3)]:
        page_type, page_subtype = _prompt_page_type_for_section(section)
        slots.append({'page_type': page_type, 'page_subtype': page_subtype, 'title': section})
    for extra_title in _prompt_content_titles(prompt, sections, request.style):
        if len(slots) >= page_count - 1:
            break
        page_type, page_subtype = _prompt_page_type_for_section(extra_title)
        slots.append({'page_type': page_type, 'page_subtype': page_subtype, 'title': extra_title})
    slots.append({'page_type': 'ending', 'page_subtype': 'thanks_end', 'title': '总结与行动'})
    while len(slots) < page_count:
        slots.insert(-1, {'page_type': 'content', 'page_subtype': 'text_image_layout', 'title': '核心内容'})
    return slots[:page_count]


def _pick_recommended_slide_ids(plan: dict[str, Any], section: str, layout_type: str) -> list[str]:
    ids: list[str] = []
    for page in plan.get('pages', []):
        if normalize_text(str(page.get('section', ''))) == normalize_text(section):
            ids.extend([str(item) for item in page.get('recommended_slide_ids', [])[:3] if str(item).strip()])
            if normalize_text(str(page.get('layout_type', ''))) == normalize_text(layout_type):
                ids = [str(item) for item in page.get('recommended_slide_ids', [])[:3] if str(item).strip()] + ids
    return list(dict.fromkeys(ids))[:3]


def _refine_prompt_plan(plan: dict[str, Any], prompt: str, style_label: str, page_range: tuple[int, int]) -> dict[str, Any]:
    sections = list(dict.fromkeys(plan.get('sections') or _prompt_sections(prompt)))
    if not sections:
        sections = ['核心内容', '总结与行动']
    min_pages, max_pages = page_range
    content_limit = max(1, max_pages - 3)
    content_sections = sections[:content_limit]
    data_profile = plan.get('data_profile')
    title = normalize_text(str(plan.get('title', ''))) or _prompt_title(prompt)
    subtitle = normalize_text(str(plan.get('subtitle', ''))) or f'{style_label} · {date.today().strftime("%Y-%m-%d")}'

    pages: list[dict[str, Any]] = [{
        'page_no': 1,
        'layout_type': 'cover_slide',
        'section': '封面',
        'title': title,
        'subtitle': subtitle,
        'bullets': [],
        'metrics': [],
        'tags': ['prompt_driven', 'cover_slide', style_label],
        'recommended_slide_ids': _pick_recommended_slide_ids(plan, '封面', 'cover_slide'),
    }, {
        'page_no': 2,
        'layout_type': 'agenda_slide',
        'section': '目录',
        'title': '目录',
        'subtitle': '',
        'bullets': content_sections,
        'metrics': [],
        'tags': ['prompt_driven', 'agenda_slide', style_label],
        'recommended_slide_ids': _pick_recommended_slide_ids(plan, '目录', 'agenda_slide'),
    }]

    for idx, section in enumerate(content_sections, start=1):
        page_type, page_subtype = _prompt_page_type_for_section(section)
        layout_type = _prompt_layout_for_section(section, data_profile)
        pages.append({
            'page_no': len(pages) + 1,
            'layout_type': layout_type,
            'section': section,
            'title': section,
            'subtitle': f'{section} · {style_label}',
            'bullets': _prompt_bullets_for_section(section, prompt, style_label, idx, len(content_sections)),
            'metrics': [],
            'tags': ['prompt_driven', page_type, page_subtype, layout_type, style_label],
            'recommended_slide_ids': _pick_recommended_slide_ids(plan, section, layout_type),
        })

    extras = _prompt_content_titles(prompt, content_sections, style_label)
    spare_slots = max(0, max_pages - (len(content_sections) + 3))
    for extra_title in extras[:spare_slots]:
        page_type, page_subtype = _prompt_page_type_for_section(extra_title)
        layout_type = _prompt_layout_for_section(extra_title, data_profile)
        pages.append({
            'page_no': len(pages) + 1,
            'layout_type': layout_type,
            'section': extra_title,
            'title': extra_title,
            'subtitle': f'{extra_title} · {style_label}',
            'bullets': _prompt_bullets_for_section(extra_title, prompt, style_label, len(pages), len(content_sections) + len(extras) + 1),
            'metrics': [],
            'tags': ['prompt_driven', page_type, page_subtype, layout_type, style_label],
            'recommended_slide_ids': _pick_recommended_slide_ids(plan, extra_title, layout_type),
        })

    pages.append({
        'page_no': len(pages) + 1,
        'layout_type': 'ending_slide',
        'section': '结束',
        'title': '感谢聆听',
        'subtitle': '谢谢',
        'bullets': ['欢迎继续补充信息', '可根据提示词再次生成新版本'],
        'metrics': [],
        'tags': ['prompt_driven', 'ending_slide', style_label],
        'recommended_slide_ids': _pick_recommended_slide_ids(plan, '结束', 'ending_slide'),
    })

    refined = deepcopy(plan)
    refined['title'] = title
    refined['subtitle'] = subtitle
    refined['sections'] = content_sections
    refined['pages'] = pages
    refined['total_pages'] = len(pages)
    refined['summary'] = ' · '.join(content_sections[:4]) if content_sections else prompt
    return refined
def _apply_prompt_runtime_patches() -> None:
    _clear_generated_library_assets()
    ppt_seeds.build_seed_sources = lambda *args, **kwargs: []
    ppt_seeds.build_template_catalog = lambda *args, **kwargs: []
    ppt_catalog.build_source_catalog = lambda *args, **kwargs: []
    ppt_catalog.build_template_catalog = lambda *args, **kwargs: []
    ppt_page_library.build_seed_slide_catalog = lambda *args, **kwargs: []
    ppt_page_library.build_seed_component_catalog = lambda *args, **kwargs: []
    ppt_manager.bootstrap_library = _seedless_bootstrap_library
    ppt_page_engine._infer_title = _prompt_title
    ppt_page_engine._infer_sections = _prompt_sections
    ppt_page_engine._sequence = _prompt_page_sequence
    ppt_smart_planner.infer_title = _prompt_title
    ppt_smart_planner.infer_sections = _prompt_sections
    ppt_smart_planner._pick_layout_for_section = _prompt_layout_for_section

    def _no_seed_demo_page_library(root: str | Path = 'ppt_template_library', refresh: bool = False) -> dict[str, Any]:
        return {'slide_count': 0, 'component_count': 0, 'template_count': 0, 'page_count': 0, 'source_count': 0}

    ppt_page_engine.seed_demo_page_library = _no_seed_demo_page_library
    ppt_app_ui.seed_demo_page_library = _no_seed_demo_page_library


_ORIGINAL_BOOTSTRAP_LIBRARY = ppt_manager.bootstrap_library
bootstrap_library = _seedless_bootstrap_library
_apply_prompt_runtime_patches()

from ppt_template_library.manager import load_library_preview_counts
from ppt_template_library.slide_storage import load_component_catalog, load_slide_catalog, load_slide_summary
from ppt_template_library.smart_planner import build_deck_plan, build_excel_aware_plan
from ppt_template_library.smart_renderer import render_smart_deck

BOOTSTRAP_KEY = '_office_library_ready'
PLAN_KEY = 'ppt_plan'
PLAN_SOURCE_KEY = 'ppt_plan_source'
PPT_BYTES_KEY = 'ppt_bytes'
PPT_PATH_KEY = 'ppt_path'
SELECTED_CANDIDATES_KEY = 'ppt_selected_candidates'
TABLE_DF_KEY = 'table_result_df'
TABLE_SUMMARY_KEY = 'table_summary_df'
UPLOADED_DF_KEY = 'uploaded_df'
UPLOADED_PROFILE_KEY = 'uploaded_profile'


def _has_streamlit_context() -> bool:
    try:
        from streamlit.runtime.scriptrunner import get_script_run_ctx
        return get_script_run_ctx() is not None
    except Exception:
        return False


def ensure_library_ready(force_refresh: bool = False) -> dict[str, Any]:
    has_context = _has_streamlit_context()
    if has_context and not force_refresh:
        cached_manifest = st.session_state.get('_office_library_manifest')
        if st.session_state.get(BOOTSTRAP_KEY) and cached_manifest:
            return cached_manifest
    if not force_refresh:
        try:
            slide_df = load_slide_catalog()
            component_df = load_component_catalog()
            manifest = load_library_preview_counts()
            if not slide_df.empty or not component_df.empty:
                if has_context:
                    st.session_state[BOOTSTRAP_KEY] = True
                    st.session_state['_office_library_manifest'] = manifest
                return manifest
        except Exception:
            pass
    spinner = st.spinner('正在初始化模板资源库...') if has_context else nullcontext()
    with spinner:
        manifest = bootstrap_library(generate_previews=False)
    if has_context:
        st.session_state[BOOTSTRAP_KEY] = True
        st.session_state['_office_library_manifest'] = manifest
    return manifest


def load_rich_slide_catalog() -> pd.DataFrame:
    # 延迟初始化模板库，避免首页冷启动时阻塞。
    return load_slide_catalog()


def _read_dataframe_from_upload(uploaded_file: Any) -> pd.DataFrame:
    name = getattr(uploaded_file, 'name', '').lower()
    raw = io.BytesIO(uploaded_file.getvalue())
    if name.endswith('.csv'):
        for encoding in ('utf-8-sig', 'gbk', 'utf-8'):
            try:
                raw.seek(0)
                df = pd.read_csv(raw, encoding=encoding)
                break
            except Exception:
                df = pd.DataFrame()
        else:
            df = pd.DataFrame()
    elif name.endswith(('.xlsx', '.xls')):
        df = pd.read_excel(raw)
    else:
        raise ValueError('仅支持 CSV / XLSX 文件。')
    df.columns = [str(col).strip() for col in df.columns]
    return df


def _normalize_columns(df: pd.DataFrame) -> pd.DataFrame:
    work = df.copy()
    work.columns = [str(col).strip() for col in work.columns]
    return work


def _infer_table_rows_from_prompt(prompt: str, default: int) -> int:
    match = re.search(r'(\d+)\s*(?:行|条|家|项|个样本|条数据)', normalize_text(prompt))
    if match:
        return max(3, min(30, int(match.group(1))))
    return default


def _table_category_values(column_name: str, row_count: int) -> list[str]:
    if any(key in column_name for key in ['门店', '店铺']):
        return [STORE_NAMES[i % len(STORE_NAMES)] for i in range(row_count)]
    if any(key in column_name for key in ['地区', '区域']):
        return [REGIONS[i % len(REGIONS)] for i in range(row_count)]
    if any(key in column_name for key in ['城市', '城市名']):
        return [CITIES[i % len(CITIES)] for i in range(row_count)]
    if any(key in column_name for key in ['产品', '品类']):
        return [f'产品{PRODUCTS[i % len(PRODUCTS)]}' for i in range(row_count)]
    if any(key in column_name for key in ['部门']):
        return [DEPARTMENTS[i % len(DEPARTMENTS)] for i in range(row_count)]
    if any(key in column_name for key in ['姓名', '成员', '负责人']):
        return [PEOPLE[i % len(PEOPLE)] for i in range(row_count)]
    if any(key in column_name for key in ['月份', '月度']):
        return [MONTHS[i % len(MONTHS)] for i in range(row_count)]
    if any(key in column_name for key in ['季度']):
        return [f'Q{(i % 4) + 1}' for i in range(row_count)]
    if any(key in column_name for key in ['日期', '时间']):
        return [f'2024-{(i % 12) + 1:02d}-01' for i in range(row_count)]
    if any(key in column_name for key in ['状态']):
        states = ['进行中', '已完成', '待开始', '有风险']
        return [states[i % len(states)] for i in range(row_count)]
    if any(key in column_name for key in ['排名', '序号', '名次']):
        return [str(i + 1) for i in range(row_count)]
    return [f'样本{i + 1}' for i in range(row_count)]


def _table_numeric_values(column_name: str, row_count: int) -> list[str]:
    series: list[str] = []
    for idx in range(row_count):
        if any(key in column_name for key in ['同比', '增速', '增长率', '环比', '占比', '转化率', '完成率', '%', '率']):
            value = 8 + idx * 1.7 + (math.sin(idx + 1) * 1.2)
            series.append(f'{value:.1f}%')
        elif any(key in column_name for key in ['营收', '销售额', '收入', '金额', '成本', '利润', '预算', '支出', '成交额']):
            base = 80 + idx * 8 + (math.sin(idx + 1) * 6)
            series.append(f'{base:.1f}万')
        elif any(key in column_name for key in ['用户', '订单', '访问', '数量', '销量', '人数', '流量', '产量']):
            series.append(str(120 + idx * 11 + int(abs(math.sin(idx + 2)) * 9)))
        elif any(key in column_name for key in ['排名', '序号', '名次']):
            series.append(str(idx + 1))
        elif any(key in column_name for key in ['分', '评分', '绩效']):
            series.append(f'{78 + (idx % 12)}')
        else:
            series.append(f'{100 + idx * 7}')
    return series


def generate_table_dataframe(prompt: str, row_count: int = 8) -> pd.DataFrame:
    columns = infer_table_columns(prompt)
    row_count = _infer_table_rows_from_prompt(prompt, row_count)
    rows: list[dict[str, Any]] = []
    for idx in range(row_count):
        row: dict[str, Any] = {}
        for column in columns:
            column_name = str(column)
            if any(key in column_name for key in ['营收', '销售额', '同比', '增速', '利润', '成本', '转化率', '排名', '序号', '名次', '得分', '评分', '数量', '用户', '订单', '访问', '产量', '%', '率']):
                values = _table_numeric_values(column_name, row_count)
                row[column_name] = values[idx]
            else:
                values = _table_category_values(column_name, row_count)
                row[column_name] = values[idx]
        rows.append(row)
    df = pd.DataFrame(rows)
    return _normalize_columns(df)


def analyze_uploaded_dataframe(df: pd.DataFrame) -> dict[str, Any]:
    df = _normalize_columns(df)
    profile = {
        'rows': int(len(df)),
        'columns': int(len(df.columns)),
        'numeric_columns': [str(col) for col in df.select_dtypes(include='number').columns],
        'text_columns': [str(col) for col in df.columns if col not in df.select_dtypes(include='number').columns],
        'date_columns': [str(col) for col in df.columns if any(keyword in str(col).lower() for keyword in ['date', 'time', 'month', 'year', '日期', '时间', '月份', '季度'])],
        'has_percent': any(re.search(r'%|率|同比|环比|增速|完成率|转化率', str(col), re.IGNORECASE) for col in df.columns),
        'has_map': any(re.search(r'地区|区域|省|市|国家|region|area|country', str(col), re.IGNORECASE) for col in df.columns),
        'has_trend': any(re.search(r'月|季度|时间|日期|year|month|time', str(col), re.IGNORECASE) for col in df.columns),
        'has_table': len(df.columns) >= 4,
        'suggested_visuals': [],
        'summary': '',
    }
    suggested: list[str] = []
    if profile['numeric_columns']:
        suggested.extend(['kpi_slide', 'dashboard_slide'])
    if profile['has_trend']:
        suggested.extend(['chart_slide', 'timeline_slide'])
    if profile['has_map']:
        suggested.append('map_slide')
    if profile['has_table']:
        suggested.append('table_slide')
    profile['suggested_visuals'] = list(dict.fromkeys(suggested))
    profile['summary'] = f"{profile['rows']} 行 × {profile['columns']} 列，数值列 {len(profile['numeric_columns'])} 个"
    return profile


def _apply_total_row(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df
    total: dict[str, Any] = {col: '' for col in df.columns}
    total[df.columns[0]] = '合计'
    for column in df.select_dtypes(include='number').columns:
        total[column] = df[column].sum()
    return pd.concat([df, pd.DataFrame([total])], ignore_index=True)


def process_dataframe(
    df: pd.DataFrame,
    dedupe: bool = True,
    fill_missing: bool = True,
    sort_column: str = '',
    ascending: bool = False,
    group_column: str = '',
    agg_column: str = '',
    agg_func: str = 'sum',
    add_total_row: bool = False,
) -> tuple[pd.DataFrame, pd.DataFrame | None]:
    work = _normalize_columns(df)
    if dedupe:
        work = work.drop_duplicates()
    if fill_missing:
        for column in work.columns:
            if pd.api.types.is_numeric_dtype(work[column]):
                fill_value = float(work[column].median()) if not work[column].dropna().empty else 0
                work[column] = work[column].fillna(fill_value)
            else:
                work[column] = work[column].fillna('未知')
    if sort_column and sort_column in work.columns:
        work = work.sort_values(sort_column, ascending=ascending)
    summary_df: pd.DataFrame | None = None
    if group_column and group_column in work.columns and agg_column in work.columns:
        agg_map = {
            'sum': 'sum',
            'mean': 'mean',
            'max': 'max',
            'min': 'min',
            'count': 'count',
        }
        op = agg_map.get(agg_func, 'sum')
        summary_df = work.groupby(group_column, dropna=False)[agg_column].agg(op).reset_index()
        summary_df.columns = [group_column, f'{agg_column}_{agg_func}']
    if add_total_row:
        work = _apply_total_row(work)
    return work, summary_df


def _style_worksheet(ws: Any, palette: dict[str, str], header_row: int = 1) -> None:
    header_fill = PatternFill(fill_type='solid', fgColor=palette['primary'])
    header_font = Font(color='FFFFFF', bold=True)
    body_font = Font(color='1F2937')
    thin = Side(style='thin', color='D8E0EA')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    for cell in ws[header_row]:
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal='center', vertical='center')
        cell.border = border
    for row in ws.iter_rows(min_row=header_row + 1):
        for cell in row:
            cell.font = body_font
            cell.border = border
            cell.alignment = Alignment(vertical='top')
    for column_cells in ws.columns:
        column_letter = column_cells[0].column_letter
        max_length = 0
        for cell in column_cells:
            value = '' if cell.value is None else str(cell.value)
            max_length = max(max_length, len(value))
        ws.column_dimensions[column_letter].width = min(max_length + 4, 36)
    ws.freeze_panes = 'A2'
    ws.auto_filter.ref = ws.dimensions


def dataframe_to_bytes(df: pd.DataFrame, file_format: str = 'xlsx', sheet_name: str = '数据表', summary_df: pd.DataFrame | None = None, theme_name: str = '简约商务') -> bytes:
    file_format = file_format.lower()
    if file_format == 'csv':
        return df.to_csv(index=False).encode('utf-8-sig')
    palette = theme_hex(theme_name)
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name=sheet_name[:31])
        _style_worksheet(writer.book[sheet_name[:31]], palette)
        if summary_df is not None and not summary_df.empty:
            summary_sheet = '汇总'
            summary_df.to_excel(writer, index=False, sheet_name=summary_sheet)
            _style_worksheet(writer.book[summary_sheet], palette)
    return buffer.getvalue()


def _resolve_plan_style(style_label: str) -> str:
    return style_code(style_label)


def build_ppt_plan(prompt: str, style_label: str, page_range: tuple[int, int], excel_df: pd.DataFrame | None = None) -> dict[str, Any]:
    page_catalog = load_rich_slide_catalog()
    style = _resolve_plan_style(style_label)
    if excel_df is not None and not excel_df.empty:
        return build_excel_aware_plan(prompt, style, excel_df, page_catalog=page_catalog, page_range=page_range)
    base_plan = build_deck_plan(prompt, style, page_range=page_range, excel_df=None, page_catalog=page_catalog)
    return _refine_prompt_plan(base_plan, prompt, style_label, page_range)
def _slide_catalog_lookup() -> pd.DataFrame:
    df = load_rich_slide_catalog()
    if df.empty:
        return df
    return df.set_index('slide_id', drop=False)


def _recommendation_rows(plan_page: dict[str, Any], slide_lookup: pd.DataFrame) -> pd.DataFrame:
    ids = list(plan_page.get('recommended_slide_ids', [])[:3])
    if slide_lookup.empty or not ids:
        return pd.DataFrame()
    rows = []
    for slide_id in ids:
        if slide_id in slide_lookup.index:
            row = slide_lookup.loc[slide_id]
            if isinstance(row, pd.DataFrame):
                row = row.iloc[0]
            rows.append(row.to_dict())
    if not rows:
        return pd.DataFrame()
    result = pd.DataFrame(rows)
    result['recommendation_order'] = pd.Categorical(result['slide_id'], categories=ids, ordered=True)
    return result.sort_values('recommendation_order')


def _page_editor_df(plan: dict[str, Any]) -> pd.DataFrame:
    rows = []
    for page in plan.get('pages', []):
        rows.append({
            'page_no': page.get('page_no'),
            'layout_type': page.get('layout_type', ''),
            'title': page.get('title', ''),
            'section': page.get('section', ''),
            'subtitle': page.get('subtitle', ''),
            'bullets': '\n'.join(page.get('bullets', [])[:4]),
            'candidate_count': len(page.get('recommended_slide_ids', [])),
        })
    return pd.DataFrame(rows)


def _apply_editor_changes(plan: dict[str, Any], edited_df: pd.DataFrame) -> dict[str, Any]:
    updated = deepcopy(plan)
    pages: list[dict[str, Any]] = []
    for _, row in edited_df.iterrows():
        page_no = int(row.get('page_no', 0))
        page = deepcopy(plan['pages'][page_no - 1]) if 0 < page_no <= len(plan.get('pages', [])) else {}
        page['title'] = normalize_text(str(row.get('title', page.get('title', '')))) or page.get('title', '')
        page['section'] = normalize_text(str(row.get('section', page.get('section', '')))) or page.get('section', '')
        page['subtitle'] = normalize_text(str(row.get('subtitle', page.get('subtitle', '')))) or page.get('subtitle', '')
        layout_type = normalize_text(str(row.get('layout_type', page.get('layout_type', ''))))
        if layout_type:
            page['layout_type'] = layout_type
        bullets_text = str(row.get('bullets', '')).strip()
        if bullets_text:
            bullets = [item.strip() for item in re.split(r'[\n；;|,，]+', bullets_text) if item.strip()]
            page['bullets'] = bullets[:6]
        pages.append(page)
    updated['pages'] = pages
    updated['total_pages'] = len(pages)
    return updated


def _apply_candidate_selection(plan: dict[str, Any], selection_map: dict[int, str], slide_lookup: pd.DataFrame) -> dict[str, Any]:
    updated = deepcopy(plan)
    selected_first: dict[str, Any] | None = None
    for page in updated.get('pages', []):
        page_no = int(page.get('page_no', 0))
        selected_id = selection_map.get(page_no)
        if not selected_id or slide_lookup.empty or selected_id not in slide_lookup.index:
            continue
        row = slide_lookup.loc[selected_id]
        if isinstance(row, pd.DataFrame):
            row = row.iloc[0]
        row_dict = row.to_dict()
        page['selected_slide_id'] = selected_id
        page['selected_candidate'] = row_dict
        page['layout_type'] = row_dict.get('layout_type', page.get('layout_type', 'content_slide'))
        page['selected_style'] = row_dict.get('style', page.get('selected_style', ''))
        if page_no == 1:
            selected_first = row_dict
    if selected_first:
        design = deepcopy(updated.get('design_system', {}))
        for key in ['primary_color', 'secondary_color', 'background_color', 'dark_or_light']:
            if selected_first.get(key):
                design[key] = selected_first[key]
        updated['design_system'] = design
        updated['style'] = selected_first.get('style', updated.get('style', 'business'))
    return updated


def _generate_pptx(plan: dict[str, Any], excel_df: pd.DataFrame | None = None) -> tuple[Path, bytes]:
    pptx_bytes = render_smart_deck(plan, excel_df)
    output_dir = Path('outputs') / 'generated_pptx'
    output_dir.mkdir(parents=True, exist_ok=True)
    title = safe_filename(f"{plan.get('title', '智能PPT')}_{date.today().isoformat()}", '.pptx')
    out_path = output_dir / title
    out_path.write_bytes(pptx_bytes)
    return out_path, pptx_bytes


def _generate_outline_json(plan: dict[str, Any]) -> bytes:
    return json.dumps(plan, ensure_ascii=False, indent=2).encode('utf-8')



# -----------------------------
# PPT 页面识别与候选规则
# -----------------------------

LAYOUT_HINT_RULES = [
    ('cover_slide', (r'封面', r'首页', r'打开页', r'cover', r'title page', r'hero')),
    ('agenda_slide', (r'目录', r'议程', r'大纲', r'agenda', r'contents')),
    ('section_slide', (r'章节', r'过渡', r'分节', r'section', r'divider')),
    ('kpi_slide', (r'KPI', r'指标', r'关键指标', r'看板', r'metric')),
    ('dashboard_slide', (r'仪表盘', r'数据看板', r'大屏', r'dashboard')),
    ('data_analysis_slide', (r'数据分析', r'洞察', r'分析报告', r'结论页', r'analysis')),
    ('chart_slide', (r'图表', r'柱状图', r'折线图', r'饼图', r'雷达图', r'散点图', r'热力图', r'趋势图', r'chart')),
    ('table_slide', (r'表格', r'明细', r'清单', r'汇总表', r'table', r'list')),
    ('timeline_slide', (r'时间线', r'时间轴', r'里程碑', r'roadmap', r'timeline', r'milestone')),
    ('process_slide', (r'流程', r'步骤', r'workflow', r'process', r'flow')),
    ('comparison_slide', (r'对比', r'差异', r'before after', r'contrast')),
    ('strategy_slide', (r'SWOT', r'PEST', r'策略', r'战略', r'strategy')),
    ('people_slide', (r'团队', r'人物', r'组织', r'profile', r'people')),
    ('product_slide', (r'产品', r'功能', r'方案', r'product')),
    ('gallery_slide', (r'图片', r'案例', r'作品集', r'gallery', r'portfolio')),
    ('map_slide', (r'地图', r'区域', r'分布', r'world', r'china')),
    ('ending_slide', (r'致谢', r'谢谢', r'结束', r'Q&A', r'thanks', r'end')),
]
def _resolve_preview_path(raw: Any) -> Path | None:
    if not raw:
        return None
    raw_path = Path(str(raw))
    if raw_path.exists():
        return raw_path
    alt = Path('ppt_template_library') / str(raw)
    if alt.exists():
        return alt
    return None


def _candidate_label(slide_id: str, slide_lookup: pd.DataFrame) -> str:
    if slide_lookup.empty or slide_id not in slide_lookup.index:
        return slide_id
    row = slide_lookup.loc[slide_id]
    if isinstance(row, pd.DataFrame):
        row = row.iloc[0]
    return ' | '.join([
        str(row.get('slide_id', slide_id)),
        str(row.get('layout_type', '')),
        str(row.get('style', '')),
        f"{int(float(row.get('overall_quality_score', 0)))}?",
    ])


def _candidate_ids_for_layout(layout_type: str, slide_lookup: pd.DataFrame, style: str = '') -> list[str]:
    if slide_lookup.empty or 'layout_type' not in slide_lookup.columns:
        return []
    result = slide_lookup.copy()
    result = result[result['layout_type'].astype(str).str.lower().eq(layout_type.lower())]
    if style and 'style' in result.columns:
        style_mask = result['style'].astype(str).str.lower().eq(style.lower())
        if style_mask.any():
            result = result[style_mask]
    if result.empty:
        return []
    ordered = result.sort_values(
        ['overall_quality_score', 'design_score', 'layout_score', 'slide_id'],
        ascending=[False, False, False, True],
    )
    return ordered['slide_id'].head(5).tolist()


def _infer_layout_from_instruction(instruction: str, fallback: str = 'content_slide') -> str:
    text = normalize_text(instruction)
    lowered = text.lower()
    for layout_type, keywords in LAYOUT_HINT_RULES:
        if any(keyword.lower() in lowered for keyword in keywords):
            return layout_type
    return fallback


def _extract_title_from_instruction(instruction: str) -> str:
    text = normalize_text(instruction)
    patterns = [
        r"(?:\u6807\u9898|\u4e3b\u9898|\u9898\u76ee|\u540d\u79f0|\u6539\u6210|\u6539\u4e3a|\u8bbe\u4e3a|\u547d\u540d\u4e3a)\s*[:?]?\s*[\"']?([^\"']{2,32})[\"']?",
        r"[\"']([^\"']{2,32})[\"']",
    ]
    for pattern in patterns:
        match = re.search(pattern, text)
        if match:
            candidate = match.group(1).strip()
            if candidate:
                return candidate
    return ''


def _extract_section_from_instruction(instruction: str) -> str:
    text = normalize_text(instruction)
    match = re.search(r"(?:\u7ae0\u8282|\u90e8\u5206|\u6a21\u5757|\u5c0f\u8282|\u5206\u8282|\u677f\u5757|section)\s*[:?]?\s*[\"']?([^\"']{2,20})[\"']?", text)
    return match.group(1).strip() if match else ''


def _apply_single_page_instruction(plan: dict[str, Any], page_no: int, instruction: str, slide_lookup: pd.DataFrame) -> dict[str, Any]:
    updated = deepcopy(plan)
    pages = list(updated.get('pages', []))
    if not (1 <= page_no <= len(pages)):
        return updated
    page = deepcopy(pages[page_no - 1])
    text = normalize_text(instruction)
    if not text:
        return updated

    title = _extract_title_from_instruction(text)
    if title:
        page['title'] = title

    section = _extract_section_from_instruction(text)
    if section:
        page['section'] = section

    layout_type = _infer_layout_from_instruction(text, page.get('layout_type', 'content_slide'))
    if layout_type:
        page['layout_type'] = layout_type
        page['recommended_slide_ids'] = _candidate_ids_for_layout(layout_type, slide_lookup, updated.get('style', ''))

    count_match = re.search(r'(\d+)\s*(?:个|条|点|项|页|段|要点|内容)?', text)
    if count_match:
        target_count = max(2, min(8, int(count_match.group(1))))
        bullets = list(page.get('bullets', []))[:target_count]
        while len(bullets) < target_count:
            bullets.append(f'要点 {len(bullets) + 1}')
        page['bullets'] = bullets
    else:
        fragments = [frag.strip() for frag in re.split(r'[\n?;|,?]+', text) if frag.strip()]
        fragments = [frag for frag in fragments if 3 <= len(frag) <= 32]
        if fragments:
            page['bullets'] = fragments[:6]

    pages[page_no - 1] = page
    updated['pages'] = pages
    updated['total_pages'] = len(pages)
    return updated
def _render_library_metrics() -> None:
    manifest = ensure_library_ready()
    summary = load_slide_summary()
    cols = st.columns(5)
    cols[0].metric('资源来源', int(manifest.get('source_count', 0)))
    cols[1].metric('模板数量', int(manifest.get('template_count', 0)))
    cols[2].metric('页面数量', int(manifest.get('slide_count', 0)))
    cols[3].metric('组件数量', int(manifest.get('component_count', 0)))
    cols[4].metric('优质模板', int(summary.get('premium_count', 0)))


def _render_uploaded_profile() -> None:
    profile = st.session_state.get(UPLOADED_PROFILE_KEY)
    if not profile:
        return
    st.markdown('**Excel 概览**')
    cols = st.columns(4)
    cols[0].metric('行数', profile.get('rows', 0))
    cols[1].metric('列数', profile.get('columns', 0))
    cols[2].metric('数值列', len(profile.get('numeric_columns', [])))
    cols[3].metric('可视化建议', len(profile.get('suggested_visuals', [])))
    st.caption(profile.get('summary', ''))


def _render_outline_editor(plan: dict[str, Any]) -> pd.DataFrame:
    st.subheader('页面大纲')
    outline_df = _page_editor_df(plan)
    if outline_df.empty:
        st.info('当前没有可编辑的大纲，请先生成方案。')
        return outline_df
    edited_df = st.data_editor(
        outline_df,
        use_container_width=True,
        hide_index=True,
        num_rows='fixed',
        disabled=['page_no', 'candidate_count'],
        key='ppt_outline_editor',
    )
    cols = st.columns(3)
    if cols[0].button('应用修改', type='primary', use_container_width=True):
        st.session_state[PLAN_KEY] = _apply_editor_changes(plan, edited_df)
        st.rerun()
    if cols[1].download_button(
        '导出 JSON',
        data=_generate_outline_json(plan),
        file_name='ppt_outline.json',
        mime='application/json',
        use_container_width=True,
    ):
        pass
    if cols[2].button('重置编辑', use_container_width=True):
        st.session_state.pop('ppt_outline_editor', None)
        st.rerun()
    return edited_df


STYLE_VARIANT_GROUPS: dict[str, list[str]] = {
    'business': ['\u7b80\u7ea6\u5546\u52a1', '\u6781\u7b80\u767d', '\u79d1\u6280\u84dd'],
    'academic': ['\u5b66\u672f\u6c47\u62a5', '\u6781\u7b80\u767d', '\u653f\u52a1\u7ea2'],
    'creative': ['\u521b\u610f\u6f14\u793a', '\u79d1\u6280\u84dd', '\u9ed1\u91d1\u9ad8\u7aef'],
    'technology': ['\u79d1\u6280\u84dd', '\u9ed1\u91d1\u9ad8\u7aef', '\u7b80\u7ea6\u5546\u52a1'],
    'finance': ['\u91d1\u878d\u7eff', '\u7b80\u7ea6\u5546\u52a1', '\u9ed1\u91d1\u9ad8\u7aef'],
    'government': ['\u653f\u52a1\u7ea2', '\u7b80\u7ea6\u5546\u52a1', '\u6781\u7b80\u767d'],
    'minimal': ['\u6781\u7b80\u767d', '\u7b80\u7ea6\u5546\u52a1', '\u5b66\u672f\u6c47\u62a5'],
    'dark': ['\u9ed1\u91d1\u9ad8\u7aef', '\u79d1\u6280\u84dd', '\u91d1\u878d\u7eff'],
    'eco_report': ['\u53ef\u6301\u7eed\u62a5\u544a', '\u62fc\u63a5\u4e2a\u4eba\u9875', '\u9ed1\u767d\u6587\u6848'],
    'split_portfolio': ['\u62fc\u63a5\u4e2a\u4eba\u9875', '\u9ed1\u767d\u6587\u6848', '\u53ef\u6301\u7eed\u62a5\u544a'],
    'editorial_copy': ['\u9ed1\u767d\u6587\u6848', '\u62fc\u63a5\u4e2a\u4eba\u9875', '\u53ef\u6301\u7eed\u62a5\u544a'],
}

LAYOUT_LABELS = {


    'cover_slide': '封面页',
    'agenda_slide': '目录页',
    'section_slide': '章节页',
    'content_slide': '内容页',
    'comparison_slide': '对比页',
    'chart_slide': '图表页',
    'dashboard_slide': '数据看板',
    'table_slide': '表格页',
    'timeline_slide': '时间线',
    'process_slide': '流程页',
    'data_analysis_slide': '数据分析页',
    'planning_slide': '规划页',
    'map_slide': '地图页',
    'kpi_slide': 'KPI 页',
    'gallery_slide': '图文展示页',
    'product_slide': '产品页',
    'people_slide': '团队页',
    'strategy_slide': '策略页',
    'ending_slide': '结束页',
}

STYLE_RHYTHM_LABELS = {
    'Low': '低密度',
    'Medium': '中密度',
    'High': '高密度',
}


def _dedupe_preserve(values: Iterable[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for value in values:
        if value and value not in seen:
            seen.add(value)
            result.append(value)
    return result


def _style_variant_labels(style_label: str) -> list[str]:
    base_style = resolve_style(style_label)
    base_code = style_code(base_style)
    labels = [base_style]
    labels.extend(STYLE_VARIANT_GROUPS.get(base_code, []))
    labels.append(style_label)
    return [label for label in _dedupe_preserve(labels) if label in PPT_STYLES][:3] or ['\u7b80\u7ea6\u5546\u52a1', '\u79d1\u6280\u84dd', '\u6781\u7b80\u767d']


def _humanize_layout(layout_type: str) -> str:
    return LAYOUT_LABELS.get(layout_type, layout_type or '内容页')


def _layout_density(page: dict[str, Any]) -> str:
    bullet_count = len(page.get('bullets', []) or [])
    if bullet_count <= 1:
        return 'Low'
    if bullet_count <= 3:
        return 'Medium'
    return 'High'


def _plan_insights(plan: dict[str, Any]) -> dict[str, Any]:
    pages = plan.get('pages', []) or []
    layouts = [str(page.get('layout_type', 'content_slide') or 'content_slide') for page in pages]
    layout_counts = Counter(layouts)
    total = max(1, len(layouts))
    dominant_layout, dominant_count = layout_counts.most_common(1)[0] if layout_counts else ('content_slide', 0)
    consecutive_same = sum(1 for idx in range(1, len(layouts)) if layouts[idx] == layouts[idx - 1])
    dominant_ratio = dominant_count / total
    consecutive_ratio = consecutive_same / max(1, len(layouts) - 1)
    similarity_score = round(min(100, dominant_ratio * 60 + consecutive_ratio * 40))
    density_series = [_layout_density(page) for page in pages]
    density_switches = sum(1 for idx in range(1, len(density_series)) if density_series[idx] != density_series[idx - 1])
    candidate_pages = sum(1 for page in pages if page.get('recommended_slide_ids'))
    return {
        'page_count': len(pages),
        'unique_layouts': len(layout_counts),
        'dominant_layout': dominant_layout,
        'dominant_layout_label': _humanize_layout(dominant_layout),
        'dominant_ratio': dominant_ratio,
        'layout_similarity_score': similarity_score,
        'layout_diversity_score': max(0, 100 - similarity_score),
        'density_series': density_series,
        'density_switches': density_switches,
        'candidate_pages': candidate_pages,
        'layout_counts': layout_counts,
    }


def _apply_style_override(plan: dict[str, Any], style_label: str) -> dict[str, Any]:
    updated = deepcopy(plan)
    palette = theme_hex(style_label)
    design = deepcopy(updated.get('design_system', {}))
    style_name = resolve_style(style_label)
    design['primary_color'] = palette['primary']
    design['secondary_color'] = palette['accent']
    design['accent_color'] = palette['accent']
    design['background_color'] = palette['bg']
    design['dark_or_light'] = 'dark' if style_code(style_name) == 'dark' else 'light'
    design['style_label'] = style_label
    updated['design_system'] = design
    updated['style'] = style_code(style_name)
    return updated


def _build_candidate_plans(prompt: str, style_label: str, page_range: tuple[int, int], excel_df: pd.DataFrame | None = None) -> list[dict[str, Any]]:
    variants: list[dict[str, Any]] = []
    for idx, variant_style in enumerate(_style_variant_labels(style_label)):
        plan = build_ppt_plan(prompt, variant_style, page_range, excel_df)
        plan = _apply_style_override(plan, variant_style)
        insights = _plan_insights(plan)
        variants.append({
            'label': f'方案{chr(65 + idx)}',
            'style_label': variant_style,
            'plan': plan,
            'insights': insights,
        })
    return variants


def _renumber_pages(plan: dict[str, Any]) -> dict[str, Any]:
    updated = deepcopy(plan)
    for idx, page in enumerate(updated.get('pages', []) or [], start=1):
        page['page_no'] = idx
    updated['total_pages'] = len(updated.get('pages', []) or [])
    return updated


def _apply_nl_edit(plan: dict[str, Any], command: str) -> dict[str, Any]:
    text = normalize_text(command)
    updated = deepcopy(plan)
    if not text:
        return updated

    style_rules = [
        (['黑金', '高端', '复古'], '黑金高端'),
        (['科技', 'AI', '数字', '智能'], '科技蓝'),
        (['金融', '银行', '证券', '保险'], '金融绿'),
        (['政务', '政府', '党建'], '政务红'),
        (['极简', '简洁', '留白'], '极简白'),
        (['学术', '论文', '答辩'], '学术汇报'),
        (['创意', '插画', '活泼'], '创意演示'),
        (['商务', '咨询', '麦肯锡'], '简约商务'),
    ]
    for keywords, style_name in style_rules:
        if any(keyword in text for keyword in keywords):
            updated = _apply_style_override(updated, style_name)
            break

    page_match = re.search(r'第(\d+)页', text)
    if page_match:
        page_no = int(page_match.group(1))
        if 1 <= page_no <= len(updated.get('pages', [])):
            page = updated['pages'][page_no - 1]
            if any(keyword in text for keyword in ['删除', '移除']):
                del updated['pages'][page_no - 1]
                return _renumber_pages(updated)
            if any(keyword in text for keyword in ['转成大数字', '冲击力', '大屏']):
                page['layout_type'] = 'dashboard_slide'
            elif any(keyword in text for keyword in ['更简洁', '简洁', '留白']):
                page['layout_type'] = 'section_slide'
            elif any(keyword in text for keyword in ['更科技', '科技感', 'AI']):
                page['layout_type'] = 'dashboard_slide'
            elif any(keyword in text for keyword in ['更高级', '高级感']):
                page['layout_type'] = 'gallery_slide'
            elif '表格' in text:
                page['layout_type'] = 'table_slide'
            elif '地图' in text:
                page['layout_type'] = 'map_slide'
            elif any(keyword in text for keyword in ['流程', '路线图']):
                page['layout_type'] = 'process_slide'
            elif any(keyword in text for keyword in ['时间线', '时间轴']):
                page['layout_type'] = 'timeline_slide'

    if '新增一页' in text or '添加一页' in text or '增加一页' in text:
        title = '新增页面'
        if '用户画像' in text:
            title = '用户画像'
        elif '行动计划' in text:
            title = '行动计划'
        elif '结论' in text:
            title = '结论'
        layout = 'content_slide'
        for keyword, layout_type in [
            ('数据', 'data_analysis_slide'),
            ('图表', 'chart_slide'),
            ('时间', 'timeline_slide'),
            ('地图', 'map_slide'),
            ('流程', 'process_slide'),
            ('对比', 'comparison_slide'),
        ]:
            if keyword in text:
                layout = layout_type
                break
        updated.setdefault('pages', []).append({
            'page_no': len(updated.get('pages', [])) + 1,
            'section': title,
            'title': title,
            'subtitle': '',
            'bullets': _sample_bullets(title, title),
            'layout_type': layout,
            'recommended_slide_ids': [],
        })
    return _renumber_pages(updated)


def _seed_demo_state() -> None:
    default_style = '简约商务'
    st.session_state.setdefault('ppt_prompt_input', DEFAULT_PPT_PROMPT)
    st.session_state.setdefault('ppt_style_input', default_style)
    st.session_state.setdefault('ppt_page_range', (8, 10))
    st.session_state.setdefault('table_prompt_input', DEFAULT_TABLE_PROMPT)
    st.session_state.setdefault('table_row_count', 8)
    st.session_state.setdefault('table_theme_input', default_style)
    if PLAN_KEY not in st.session_state:
        try:
            variants = _build_candidate_plans(DEFAULT_PPT_PROMPT, default_style, (8, 10), None)
            if variants:
                st.session_state['ppt_candidate_variants'] = variants
                st.session_state['ppt_variant_index'] = 0
                st.session_state[PLAN_KEY] = deepcopy(variants[0]['plan'])
            else:
                st.session_state[PLAN_KEY] = build_ppt_plan(DEFAULT_PPT_PROMPT, default_style, (8, 10), None)
        except Exception:
            try:
                st.session_state[PLAN_KEY] = build_ppt_plan(DEFAULT_PPT_PROMPT, default_style, (8, 10), None)
            except Exception:
                pass
    if TABLE_DF_KEY not in st.session_state:
        st.session_state[TABLE_DF_KEY] = generate_table_dataframe(DEFAULT_TABLE_PROMPT, 8)
        st.session_state[TABLE_SUMMARY_KEY] = None


def _render_badges(values: Iterable[str]) -> None:
    pills = ''.join(f"<span class='metric-pill'>{value}</span>" for value in values if value)
    if pills:
        st.markdown(pills, unsafe_allow_html=True)


def _render_plan_overview(plan: dict[str, Any]) -> dict[str, Any]:
    insights = _plan_insights(plan)
    cols = st.columns(4)
    cols[0].metric('页面数', insights['page_count'])
    cols[1].metric('布局重复度', f"{insights['layout_similarity_score']}%")
    cols[2].metric('路径切换', insights['density_switches'])
    cols[3].metric('候选覆盖', insights['candidate_pages'])
    if insights['layout_similarity_score'] >= 75:
        st.warning('检测到布局重复度较高，建议立即更换部分页面布局。')
    _render_badges([
        f"风格: {resolve_style(str(plan.get('style', 'business')))}",
        f"行业: {plan.get('industry', '')}",
        f"场景: {plan.get('scenario', '')}",
        f"主布局: {insights['dominant_layout_label']}",
    ])
    st.caption(' → '.join(STYLE_RHYTHM_LABELS.get(item, item) for item in insights['density_series']) or '暂无数据')
    return insights


def _render_candidate_rows(plan_page: dict[str, Any], slide_lookup: pd.DataFrame, page_no: int) -> str:
    rec_df = _recommendation_rows(plan_page, slide_lookup)
    default_choice = '自动推荐'
    if rec_df.empty:
        st.info('该页暂无对应候选布局，可以保持自动绘制或再刷新调整。')
        return default_choice
    preview_cols = [c for c in ['slide_id', 'slide_type', 'slide_subtype', 'layout_type', 'overall_quality_score', 'preview_image'] if c in rec_df.columns]
    st.dataframe(rec_df[preview_cols].head(3), use_container_width=True, hide_index=True)
    img_cols = st.columns(min(3, len(rec_df)))
    for idx, (_, row) in enumerate(rec_df.head(3).iterrows()):
        with img_cols[idx]:
            preview_image = str(row.get('preview_image', '')).strip()
            if preview_image and Path(preview_image).exists():
                st.image(preview_image, use_container_width=True)
            st.caption(f"{row.get('slide_id', '')} | {row.get('overall_quality_score', 0)}")
    options = ['自动推荐'] + rec_df['slide_id'].astype(str).tolist()
    preferred = str(plan_page.get('selected_slide_id', '')).strip()
    if preferred not in options:
        preferred = options[1] if len(options) > 1 else options[0]
    choice = st.selectbox('布局选择', options, index=options.index(preferred), key=f'ppt_candidate_choice_{page_no}')
    return choice


def render_ppt_workbench() -> None:
    st.header('PPT 生成工作台')
    st.caption('先出三套方案，再按页编辑，最后导出 .pptx。')
    left, right = st.columns([1.25, 1])
    with left:
        prompt = st.text_area('需求描述', value=st.session_state.get('ppt_prompt_input', DEFAULT_PPT_PROMPT), height=140, key='ppt_prompt_input')
        style_label = st.selectbox('主风格', list(STYLE_UI_TO_CODE.keys()), index=list(STYLE_UI_TO_CODE.keys()).index(st.session_state.get('ppt_style_input', '简约商务')) if st.session_state.get('ppt_style_input', '简约商务') in STYLE_UI_TO_CODE else 0, key='ppt_style_input')
        page_range = st.slider('页数范围', 6, 16, st.session_state.get('ppt_page_range', (8, 10)), key='ppt_page_range')
    with right:
        uploaded = st.file_uploader('上传 CSV / XLSX 作为数据参考', type=['csv', 'xlsx', 'xls'], key='ppt_data_upload')
        if uploaded is not None:
            try:
                uploaded_df = _read_dataframe_from_upload(uploaded)
                st.session_state[UPLOADED_DF_KEY] = uploaded_df
                st.session_state[UPLOADED_PROFILE_KEY] = analyze_uploaded_dataframe(uploaded_df)
                st.success(st.session_state[UPLOADED_PROFILE_KEY]['summary'])
                st.dataframe(uploaded_df.head(8), use_container_width=True, hide_index=True)
            except Exception as exc:  # noqa: BLE001
                st.error(f'数据读取失败：{exc}')
        elif st.session_state.get(UPLOADED_PROFILE_KEY):
            st.info(st.session_state[UPLOADED_PROFILE_KEY]['summary'])
    generate = st.button('生成 3 套方案', type='primary', use_container_width=True)
    if generate:
        try:
            with st.spinner('正在分析需求并生成候选方案...'):
                variants = _build_candidate_plans(prompt, style_label, page_range, st.session_state.get(UPLOADED_DF_KEY))
            if variants:
                st.session_state['ppt_candidate_variants'] = variants
                st.session_state['ppt_variant_index'] = 0
                st.session_state[PLAN_KEY] = deepcopy(variants[0]['plan'])
                st.session_state[PLAN_SOURCE_KEY] = 'variants'
                st.success('方案已生成')
                st.rerun()
        except Exception as exc:  # noqa: BLE001
            st.error(f'生成失败：{exc}')
    variants = st.session_state.get('ppt_candidate_variants', [])
    if not variants:
        st.info('点击?生成 3 套方案?，我们会先给出三套不同风格的布局方案。')
        return
    active_idx = int(st.session_state.get('ppt_variant_index', 0))
    active_idx = max(0, min(active_idx, len(variants) - 1))
    plan = deepcopy(st.session_state.get(PLAN_KEY, variants[active_idx]['plan']))
    st.subheader('方案选择')
    variant_cols = st.columns(len(variants))
    for idx, item in enumerate(variants):
        with variant_cols[idx]:
            st.markdown(f"**{item['label']}**")
            st.caption(item['style_label'])
            st.metric('页数', item['insights']['page_count'])
            st.metric('布局重复度', f"{item['insights']['layout_similarity_score']}%")
            st.write(f"主布局: {item['insights']['dominant_layout_label']}")
            st.write(' → '.join(STYLE_RHYTHM_LABELS.get(d, d) for d in item['insights']['density_series']) or '暂无数据')
            if st.button('采用此方案', key=f'choose_variant_{idx}', use_container_width=True):
                st.session_state['ppt_variant_index'] = idx
                st.session_state[PLAN_KEY] = deepcopy(item['plan'])
                st.session_state[PLAN_SOURCE_KEY] = item['label']
                st.rerun()
    active_idx = int(st.session_state.get('ppt_variant_index', 0))
    active_idx = max(0, min(active_idx, len(variants) - 1))
    plan = deepcopy(st.session_state.get(PLAN_KEY, variants[active_idx]['plan']))
    st.success(f"当前采用：{variants[active_idx]['label']} · {variants[active_idx]['style_label']}")
    _render_plan_overview(plan)
    st.subheader('页面大纲')
    editor_df = _page_editor_df(plan)
    edited_df = st.data_editor(editor_df, use_container_width=True, hide_index=True, disabled=['page_no', 'candidate_count'])
    editor_cols = st.columns(3)
    with editor_cols[0]:
        if st.button('应用页面修改', type='primary', use_container_width=True):
            st.session_state[PLAN_KEY] = _apply_editor_changes(plan, edited_df)
            st.rerun()
    with editor_cols[1]:
        if st.button('恢复当前方案', use_container_width=True):
            st.session_state[PLAN_KEY] = deepcopy(variants[active_idx]['plan'])
            st.rerun()
    with editor_cols[2]:
        if st.button('刷新候选方案', use_container_width=True):
            st.session_state.pop('ppt_candidate_variants', None)
            st.rerun()

    st.subheader('页面候选布局')
    slide_lookup = _slide_catalog_lookup()
    selection_map: dict[int, str] = {}
    for page in plan.get('pages', []) or []:
        page_no = int(page.get('page_no', 0))
        title = str(page.get('title', ''))
        section = str(page.get('section', ''))
        layout_label = _humanize_layout(str(page.get('layout_type', 'content_slide')))
        with st.expander(f"{page_no}. {section or title} - {title} - {layout_label}", expanded=page_no <= 2):
            bullets = page.get('bullets', []) or []
            if bullets:
                st.write(f"核心内容：{bullets[0]}")
            choice = _render_candidate_rows(page, slide_lookup, page_no)
            if choice != '自动推荐':
                selection_map[page_no] = choice
    if st.button('应用候选布局', use_container_width=True):
        st.session_state[PLAN_KEY] = _apply_candidate_selection(plan, selection_map, slide_lookup)
        st.rerun()

    st.subheader('自然语言修改')
    nl_command = st.text_area('修改指令', value=st.session_state.get('ppt_nl_command', ''), height=100, key='ppt_nl_command')
    nl_cols = st.columns(2)
    with nl_cols[0]:
        if st.button('应用修改', use_container_width=True):
            st.session_state[PLAN_KEY] = _apply_nl_edit(plan, nl_command)
            st.rerun()
    with nl_cols[1]:
        if st.button('清空修改指令', use_container_width=True):
            st.session_state['ppt_nl_command'] = ''
            st.rerun()

    st.subheader('导出')
    export_cols = st.columns(3)
    with export_cols[0]:
        if st.button('生成 PPTX', type='primary', use_container_width=True):
            try:
                with st.spinner('正在生成 PPTX...'):
                    out_path, pptx_bytes = _generate_pptx(st.session_state.get(PLAN_KEY, plan), st.session_state.get(UPLOADED_DF_KEY))
                st.session_state[PPT_PATH_KEY] = str(out_path)
                st.session_state[PPT_BYTES_KEY] = pptx_bytes
                st.success(f'已生成：{out_path.name}')
            except Exception as exc:  # noqa: BLE001
                st.error(f'PPTX 生成失败：{exc}')
    with export_cols[1]:
        if st.session_state.get(PPT_BYTES_KEY):
            st.download_button('下载 PPTX', data=st.session_state[PPT_BYTES_KEY], file_name=safe_filename(str(plan.get('title', '智能PPT')), '.pptx'), mime='application/vnd.openxmlformats-officedocument.presentationml.presentation', use_container_width=True)
    with export_cols[2]:
        st.download_button('下载大纲 JSON', data=_generate_outline_json(st.session_state.get(PLAN_KEY, plan)), file_name='ppt_outline.json', mime='application/json', use_container_width=True)
def _render_table_downloads(df: pd.DataFrame, summary_df: pd.DataFrame | None, theme_name: str, prefix: str) -> None:
    if df.empty:
        return
    xlsx_bytes = dataframe_to_bytes(df, 'xlsx', summary_df=summary_df, theme_name=theme_name)
    csv_bytes = dataframe_to_bytes(df, 'csv', summary_df=summary_df, theme_name=theme_name)
    cols = st.columns(2)
    with cols[0]:
        st.download_button(f'下载 {prefix}XLSX', data=xlsx_bytes, file_name=f'{prefix}.xlsx', mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', use_container_width=True)
    with cols[1]:
        st.download_button(f'下载 {prefix}CSV', data=csv_bytes, file_name=f'{prefix}.csv', mime='text/csv', use_container_width=True)


def render_table_workbench() -> None:
    st.header('表格工具')
    tabs = st.tabs(['文本生成', '文件处理'])
    with tabs[0]:
        st.subheader('文本生成表格')
        prompt = st.text_area('表格描述', value=st.session_state.get('table_prompt_input', DEFAULT_TABLE_PROMPT), height=120, key='table_prompt_input')
        row_count = st.slider('示例行数', 3, 20, int(st.session_state.get('table_row_count', 8)), key='table_row_count')
        theme_options = list(STYLE_UI_TO_CODE.keys())
        current_theme = st.session_state.get('table_theme_input', '简约商务')
        theme_name = st.selectbox('表格风格', theme_options, index=theme_options.index(current_theme) if current_theme in STYLE_UI_TO_CODE else 0, key='table_theme_input')
        cols = st.columns(2)
        with cols[0]:
            if st.button('生成示例表格', type='primary', use_container_width=True):
                generated_df = generate_table_dataframe(prompt, row_count)
                summary_df = None
                st.session_state[TABLE_DF_KEY] = generated_df
                st.session_state[TABLE_SUMMARY_KEY] = summary_df
                st.success(f'已生成 {len(generated_df)} 行数据')
        with cols[1]:
            if st.button('恢复系统示例', use_container_width=True):
                st.session_state[TABLE_DF_KEY] = generate_table_dataframe(DEFAULT_TABLE_PROMPT, 8)
                st.session_state[TABLE_SUMMARY_KEY] = None
                st.rerun()
        generated_df = st.session_state.get(TABLE_DF_KEY, pd.DataFrame())
        if not generated_df.empty:
            st.dataframe(generated_df, use_container_width=True, hide_index=True)
            _render_table_downloads(generated_df, st.session_state.get(TABLE_SUMMARY_KEY), theme_name, 'generated_table')
        else:
            st.info('暂无表格数据，请先生成。')
    with tabs[1]:
        st.subheader('文件处理')
        uploaded = st.file_uploader('上传 CSV / XLSX', type=['csv', 'xlsx', 'xls'], key='table_upload')
        raw_df = st.session_state.get(UPLOADED_DF_KEY, pd.DataFrame())
        if uploaded is not None:
            try:
                raw_df = _read_dataframe_from_upload(uploaded)
                st.session_state[UPLOADED_DF_KEY] = raw_df
                st.session_state[UPLOADED_PROFILE_KEY] = analyze_uploaded_dataframe(raw_df)
                st.success(st.session_state[UPLOADED_PROFILE_KEY]['summary'])
            except Exception as exc:  # noqa: BLE001
                st.error(f'文件读取失败：{exc}')
        if raw_df is not None and not raw_df.empty:
            profile = st.session_state.get(UPLOADED_PROFILE_KEY) or analyze_uploaded_dataframe(raw_df)
            st.json(profile)
            st.dataframe(raw_df.head(10), use_container_width=True, hide_index=True)
            controls = st.columns(3)
            with controls[0]:
                dedupe = st.checkbox('去重', value=True, key='table_dedupe')
                fill_missing = st.checkbox('填充空值', value=True, key='table_fill_missing')
                add_total_row = st.checkbox('添加汇总行', value=False, key='table_total_row')
            with controls[1]:
                sort_column = st.selectbox('排序列', [''] + raw_df.columns.tolist(), index=0, key='table_sort_column')
                ascending = st.checkbox('升序', value=False, key='table_sort_ascending')
            with controls[2]:
                group_column = st.selectbox('分组列', [''] + raw_df.columns.tolist(), index=0, key='table_group_column')
                agg_candidates = [''] + [col for col in raw_df.columns if pd.api.types.is_numeric_dtype(raw_df[col])]
                agg_column = st.selectbox('汇总列', agg_candidates, index=0 if len(agg_candidates) == 1 else 1, key='table_agg_column')
            agg_func = st.selectbox('汇总方式', ['sum', 'mean', 'max', 'min', 'count'], key='table_agg_func')
            if st.button('处理表格', type='primary', use_container_width=True):
                try:
                    processed_df, summary_df = process_dataframe(raw_df, dedupe=dedupe, fill_missing=fill_missing, sort_column=sort_column, ascending=ascending, group_column=group_column, agg_column=agg_column, agg_func=agg_func, add_total_row=add_total_row)
                    st.session_state[TABLE_DF_KEY] = processed_df
                    st.session_state[TABLE_SUMMARY_KEY] = summary_df
                    st.success(f'处理完成，当前 {len(processed_df)} 行')
                except Exception as exc:  # noqa: BLE001
                    st.error(f'处理失败：{exc}')
            processed_df = st.session_state.get(TABLE_DF_KEY)
            summary_df = st.session_state.get(TABLE_SUMMARY_KEY)
            if isinstance(processed_df, pd.DataFrame) and not processed_df.empty:
                st.dataframe(processed_df, use_container_width=True, hide_index=True)
                export_bytes = dataframe_to_bytes(processed_df, file_format='xlsx', summary_df=summary_df, theme_name=theme_name)
                st.download_button('下载 XLSX', data=export_bytes, file_name='processed_table.xlsx', mime='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet', use_container_width=True)
                csv_bytes = dataframe_to_bytes(processed_df, file_format='csv')
                st.download_button('下载 CSV', data=csv_bytes, file_name='processed_table.csv', mime='text/csv', use_container_width=True)
            if isinstance(summary_df, pd.DataFrame) and not summary_df.empty:
                st.markdown('**分组汇总**')
                st.dataframe(summary_df, use_container_width=True, hide_index=True)
        else:
            st.info('请先上传 CSV 或 XLSX 文件。')


def render_template_library_workbench() -> None:
    render_template_library_page()


def main() -> None:
    setup_page()
    _seed_demo_state()
    st.title(APP_TITLE)
    st.caption('一个页面里打通 PPT 、表格和模板库。')
    manifest = st.session_state.get('_office_library_manifest') or load_library_preview_counts()
    sidebar = st.sidebar
    sidebar.subheader('功能切换')
    sidebar.radio('', ['PPT 生成', '表格工具', '模板资源库'], key='nav_choice', label_visibility='collapsed')
    sidebar.caption(f"页面: {manifest.get('slide_count', 0)}  组件: {manifest.get('component_count', 0)}  资源: {manifest.get('template_count', 0)}")
    choice = st.session_state.get('nav_choice', 'PPT 生成')
    if choice == 'PPT 生成':
        render_ppt_workbench()
    elif choice == '表格工具':
        render_table_workbench()
    else:
        render_template_library_workbench()


if __name__ == '__main__':
    main()
