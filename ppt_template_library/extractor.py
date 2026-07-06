from __future__ import annotations

import re
from dataclasses import asdict
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .page_models import ComponentRecord, PageRecord
from .page_library import LAYOUT_BLUEPRINTS
from .preview import create_component_preview, create_page_preview

EMU_PER_INCH = 914400


def _emu_to_inches(value: int | float | None) -> float:
    if value is None:
        return 0.0
    return round(float(value) / EMU_PER_INCH, 2)


def _shape_text(shape: Any) -> str:
    if not getattr(shape, 'has_text_frame', False):
        return ''
    try:
        return '\n'.join(paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text)
    except Exception:
        return ''


def _slide_text(slide: Any) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        text = _shape_text(shape)
        if text:
            parts.append(text)
    return '\n'.join(parts)


def _shape_kind(shape: Any) -> str:
    if getattr(shape, 'has_chart', False):
        return 'chart'
    if getattr(shape, 'has_table', False):
        return 'table'
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return 'image'
    if getattr(shape, 'has_text_frame', False) and _shape_text(shape).strip():
        return 'text'
    return 'shape'


def analyze_slide(slide: Any, slide_number: int = 1, source_template_id: str = '', source_url: str = '') -> dict[str, Any]:
    text = _slide_text(slide)
    kinds: dict[str, int] = {'text': 0, 'image': 0, 'chart': 0, 'table': 0, 'shape': 0}
    for shape in slide.shapes:
        kinds[_shape_kind(shape)] = kinds.get(_shape_kind(shape), 0) + 1
    title = ''
    if slide.shapes and getattr(slide.shapes[0], 'has_text_frame', False):
        title = _shape_text(slide.shapes[0]).split('\n')[0].strip()
    return {
        'slide_number': slide_number,
        'title': title or f'Slide {slide_number}',
        'text': text,
        'shape_count': len(slide.shapes),
        'text_box_count': kinds.get('text', 0),
        'image_count': kinds.get('image', 0),
        'chart_count': kinds.get('chart', 0),
        'table_count': kinds.get('table', 0),
        'shape_types': kinds,
        'source_template_id': source_template_id,
        'source_url': source_url,
    }


def classify_slide(text: str, stats: dict[str, Any]) -> str:
    combined = f"{stats.get('title', '')}\n{text}".lower()
    if stats.get('slide_number', 1) == 1 or re.search(r'封面|cover|title', combined):
        return 'cover_slide'
    if re.search(r'目录|agenda|议程', combined):
        return 'agenda_slide'
    if re.search(r'结束|感谢|thank you|q&a|联系', combined):
        return 'ending_slide'
    if stats.get('chart_count', 0) or stats.get('table_count', 0):
        if re.search(r'仪表盘|dashboard|kpi|指标', combined):
            return 'dashboard_slide'
        if stats.get('table_count', 0) and not stats.get('chart_count', 0):
            return 'table_slide'
        return 'data_analysis_slide'
    if re.search(r'时间轴|里程碑|路线图|timeline|roadmap|历史|历程', combined):
        return 'timeline_slide'
    if re.search(r'流程|步骤|审批|工作流|workflow|process', combined):
        return 'process_slide'
    if re.search(r'对比|竞品|优劣|before|after', combined):
        return 'comparison_slide'
    if re.search(r'组织|架构|团队|成员|人员|关系|network|org', combined):
        return 'relation_slide'
    if re.search(r'swot|pest|波特|战略|定位|商业模式', combined):
        return 'strategy_slide'
    if re.search(r'计划|规划|gantt|任务|排期|roadmap', combined):
        return 'planning_slide'
    if re.search(r'产品|功能|发布|参数|feature', combined):
        return 'product_slide'
    if re.search(r'地图|区域|地区|分布|world|china', combined):
        return 'map_slide'
    if stats.get('image_count', 0) >= 2:
        return 'gallery_slide'
    if stats.get('text_box_count', 0) <= 3 and stats.get('image_count', 0) >= 1:
        return 'cover_slide'
    if stats.get('text_box_count', 0) >= 8 and stats.get('shape_count', 0) >= 10:
        return 'content_slide'
    return 'content_slide'


def extract_components_from_slide(slide: Any, slide_id: str, slide_type: str, style: str, industry: str) -> list[ComponentRecord]:
    components: list[ComponentRecord] = []
    for index, shape in enumerate(slide.shapes, start=1):
        kind = _shape_kind(shape)
        if kind == 'shape' and not getattr(shape, 'has_text_frame', False):
            continue
        left = _emu_to_inches(getattr(shape, 'left', 0))
        top = _emu_to_inches(getattr(shape, 'top', 0))
        width = _emu_to_inches(getattr(shape, 'width', 0))
        height = _emu_to_inches(getattr(shape, 'height', 0))
        text = _shape_text(shape)
        component_type = 'text_block' if kind == 'text' else kind if kind != 'shape' else 'shape_block'
        if kind == 'chart':
            component_type = 'chart_block'
        elif kind == 'table':
            component_type = 'table_block'
        elif kind == 'image':
            component_type = 'gallery_block'
        component_id = f'{slide_id}-{index:03d}-{component_type}'
        components.append(
            ComponentRecord(
                component_id=component_id,
                component_type=component_type.replace('_block', ''),
                component_subtype=f'{component_type}_{index}',
                source_slide_id=slide_id,
                layout_type=slide_type,
                bounding_box=f'x:{left},y:{top},w:{width},h:{height}',
                style_token=style,
                color_token='#1B3756',
                width_ratio=round(min(1.0, width / 13.33 if width else 0.0), 2),
                height_ratio=round(min(1.0, height / 7.5 if height else 0.0), 2),
                preview_image='',
                metadata_json={
                    'slide_type': slide_type,
                    'kind': kind,
                    'text': text[:300],
                    'industry': industry,
                },
                status='extracted',
                tags=(slide_type, kind, style, industry),
            )
        )
    return components


def build_page_record_from_slide(
    slide: Any,
    slide_number: int,
    source_template_id: str,
    source_url: str,
    style: str = 'business',
    industry: str = 'general',
) -> PageRecord:
    stats = analyze_slide(slide, slide_number=slide_number, source_template_id=source_template_id, source_url=source_url)
    slide_type = classify_slide(stats['text'], stats)
    blueprint = LAYOUT_BLUEPRINTS.get(slide_type, LAYOUT_BLUEPRINTS['content_slide'])
    preview_path = str(Path('ppt_template_library') / 'preview' / 'slides' / f"{source_template_id}-{slide_number:02d}.png")
    page = PageRecord(
        slide_id=f"{source_template_id}-{slide_number:02d}",
        source_template_id=source_template_id,
        source_file='',
        source_url=source_url,
        slide_number=slide_number,
        slide_type=blueprint['slide_type'],
        slide_subtype=slide_type,
        industry=industry,
        scenario='imported',
        style=style,
        layout_type=slide_type.replace('_slide', ''),
        primary_color='#1B3756',
        secondary_color='#245A8D',
        background_color='#F7F9FC',
        dark_or_light='light',
        aspect_ratio='16:9',
        text_density=min(1.0, round(stats['text_box_count'] / max(1, stats['shape_count']), 2)),
        image_density=min(1.0, round(stats['image_count'] / max(1, stats['shape_count']), 2)),
        chart_count=stats['chart_count'],
        table_count=stats['table_count'],
        shape_count=stats['shape_count'],
        icon_count=0,
        image_count=stats['image_count'],
        text_box_count=stats['text_box_count'],
        has_chart=stats['chart_count'] > 0,
        has_table=stats['table_count'] > 0,
        has_timeline=bool(re.search(r'时间轴|里程碑|路线图|timeline|roadmap', stats['text'], re.IGNORECASE)),
        has_process=bool(re.search(r'流程|步骤|审批|工作流', stats['text'], re.IGNORECASE)),
        has_map=bool(re.search(r'地图|区域|地区|分布', stats['text'], re.IGNORECASE)),
        has_people=bool(re.search(r'组织|团队|成员|人员|关系', stats['text'], re.IGNORECASE)),
        has_infographic=stats['shape_count'] >= 8,
        has_animation=False,
        editable_level='high' if stats['shape_count'] < 40 else 'medium',
        design_score=80,
        layout_score=82,
        color_score=80,
        usability_score=82,
        modern_score=78,
        overall_quality_score=82,
        preview_image=preview_path,
        thumbnail_path=preview_path,
        slide_file_path='',
        embedding_vector='',
        metadata_json={
            'title': stats['title'],
            'text': stats['text'][:1000],
            'shape_types': stats['shape_types'],
            'classification_reason': slide_type,
        },
        status='extracted',
        tags=(slide_type, style, industry),
    )
    return page


def analyze_pptx_file(
    pptx_path: str | Path,
    source_template_id: str | None = None,
    source_url: str = '',
    style: str = 'business',
    industry: str = 'general',
    generate_previews: bool = False,
) -> dict[str, Any]:
    path = Path(pptx_path)
    prs = Presentation(path)
    template_id = source_template_id or path.stem
    pages: list[PageRecord] = []
    components: list[ComponentRecord] = []
    for index, slide in enumerate(prs.slides, start=1):
        page = build_page_record_from_slide(slide, index, template_id, source_url or path.as_uri(), style=style, industry=industry)
        pages.append(page)
        components.extend(extract_components_from_slide(slide, page.slide_id, page.slide_subtype, page.style, page.industry))
    if generate_previews:
        for page in pages:
            create_page_preview(title=page.metadata_json.get('title', page.slide_type), subtitle=f'{page.style} ? {page.industry} ? {page.scenario}', layout_type=page.layout_type, score=page.overall_quality_score, tags=page.tags, out_path=page.preview_image)
        for component in components[:50]:
            create_component_preview(title=component.component_type.replace('_', ' ').title(), subtitle=f'{component.style_token} ? {component.layout_type}', component_type=component.component_type, score=80, tags=component.tags, out_path=component.preview_image)
    return {
        'template_id': template_id,
        'source_url': source_url or path.as_uri(),
        'slides': pages,
        'components': components,
        'summary': {
            'slide_count': len(pages),
            'component_count': len(components),
            'layout_types': len({page.layout_type for page in pages}),
            'chart_pages': sum(1 for page in pages if page.has_chart),
            'table_pages': sum(1 for page in pages if page.has_table),
            'timeline_pages': sum(1 for page in pages if page.has_timeline),
            'process_pages': sum(1 for page in pages if page.has_process),
            'map_pages': sum(1 for page in pages if page.has_map),
        },
    }
