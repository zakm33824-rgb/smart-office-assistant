from __future__ import annotations

import io
import math
import re
from dataclasses import dataclass
from datetime import date
from typing import Any, Iterable

import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


LIGHT_PALETTE = {
    'bg': '#F7F9FC',
    'paper': '#FFFFFF',
    'primary': '#1B3756',
    'secondary': '#245A8D',
    'accent': '#0F766E',
    'accent2': '#F59E0B',
    'text': '#102A43',
    'muted': '#5B6B7A',
    'line': '#D8E0EA',
    'soft': '#EEF5FF',
    'soft2': '#EAFBF9',
}

DARK_PALETTE = {
    'bg': '#0F172A',
    'paper': '#111827',
    'primary': '#E5E7EB',
    'secondary': '#93C5FD',
    'accent': '#22D3EE',
    'accent2': '#F59E0B',
    'text': '#F8FAFC',
    'muted': '#CBD5E1',
    'line': '#334155',
    'soft': '#1E293B',
    'soft2': '#14303C',
}


@dataclass(frozen=True)
class Palette:
    bg: RGBColor
    paper: RGBColor
    primary: RGBColor
    secondary: RGBColor
    accent: RGBColor
    accent2: RGBColor
    text: RGBColor
    muted: RGBColor
    line: RGBColor
    soft: RGBColor
    soft2: RGBColor



def _hex_to_rgb(value: str) -> RGBColor:
    value = value.strip().lstrip('#')
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _palette_from_design(design: dict[str, Any]) -> Palette:
    tone = (design.get('dark_or_light') or 'light').lower()
    base = DARK_PALETTE if tone == 'dark' else LIGHT_PALETTE
    return Palette(
        bg=_hex_to_rgb(design.get('background_color') or base['bg']),
        paper=_hex_to_rgb(base['paper']),
        primary=_hex_to_rgb(design.get('primary_color') or base['primary']),
        secondary=_hex_to_rgb(design.get('secondary_color') or base['secondary']),
        accent=_hex_to_rgb(design.get('accent_color') or base['accent']),
        accent2=_hex_to_rgb(base['accent2']),
        text=_hex_to_rgb(base['text']),
        muted=_hex_to_rgb(base['muted']),
        line=_hex_to_rgb(base['line']),
        soft=_hex_to_rgb(base['soft']),
        soft2=_hex_to_rgb(base['soft2']),
    )


def _font_size(design: dict[str, Any], key: str, default: int) -> int:
    value = design.get(key, default)
    try:
        return int(value)
    except Exception:
        return default


def _setup_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _set_fill(shape: Any, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _set_line(shape: Any, color: RGBColor | None = None, transparency: int = 100000) -> None:
    shape.line.fill.background()
    if color is not None:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = color
    shape.line.transparency = transparency


def _add_text(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = 'Microsoft YaHei',
) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_multiline(slide: Any, lines: Iterable[str], left: float, top: float, width: float, height: float, color: RGBColor, font_size: int = 15) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.space_after = Pt(6)
        run = paragraph.runs[0]
        run.font.name = 'Microsoft YaHei'
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def _add_card(slide: Any, left: float, top: float, width: float, height: float, fill: RGBColor, line: RGBColor) -> None:
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.03), Inches(top + 0.05), Inches(width), Inches(height))
    _set_fill(shadow, RGBColor(225, 230, 237))
    _set_line(shadow)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    _set_fill(card, fill)
    _set_line(card, line, transparency=0)


def _add_badge(slide: Any, text: str, left: float, top: float, width: float, fill: RGBColor, color: RGBColor) -> None:
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.34))
    _set_fill(badge, fill)
    _set_line(badge)
    _add_text(slide, text, left + 0.06, top + 0.03, width - 0.12, 0.14, 9, color, True, PP_ALIGN.CENTER)


def _footer(slide: Any, palette: Palette, page_no: int, total_pages: int) -> None:
    _add_text(slide, f'{page_no:02d} / {total_pages:02d}', 11.75, 6.95, 1.0, 0.2, 9, palette.muted, align=PP_ALIGN.RIGHT)


def _title_band(slide: Any, palette: Palette) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.16))
    _set_fill(bar, palette.accent)
    _set_line(bar)


def _metric_card(slide: Any, palette: Palette, metric: dict[str, str], left: float, top: float, width: float) -> None:
    _add_card(slide, left, top, width, 1.02, palette.paper, palette.line)
    _add_text(slide, metric['label'], left + 0.16, top + 0.12, width - 0.28, 0.2, 9, palette.muted)
    _add_text(slide, metric['value'], left + 0.16, top + 0.37, width - 0.28, 0.3, 18, palette.primary, True)
    _add_text(slide, metric.get('trend', ''), left + 0.16, top + 0.69, width - 0.28, 0.16, 8, palette.accent)


def _sample_line_values() -> list[float]:
    return [0.45, 0.62, 0.55, 0.73, 0.68, 0.84]


def _tiny_bar_chart(slide: Any, palette: Palette, left: float, top: float) -> None:
    values = _sample_line_values()
    for idx, value in enumerate(values):
        x = left + idx * 0.38
        height = 1.45 * value
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top + 1.5 - height), Inches(0.22), Inches(height))
        _set_fill(bar, palette.accent if idx in (2, 5) else palette.secondary)
        _set_line(bar)


def _simple_table(slide: Any, palette: Palette, rows: list[list[str]], left: float, top: float, width: float, height: float) -> None:
    if not rows:
        return
    row_count = len(rows)
    col_count = len(rows[0])
    _add_card(slide, left, top, width, height, palette.paper, palette.line)
    cell_w = width / col_count
    cell_h = height / row_count
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            fill_color = palette.soft if r == 0 else (palette.paper if r % 2 == 1 else palette.soft2)
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + c * cell_w), Inches(top + r * cell_h), Inches(cell_w), Inches(cell_h))
            _set_fill(cell, fill_color)
            _set_line(cell, palette.line, transparency=0)
            _add_text(slide, str(value), left + c * cell_w + 0.08, top + r * cell_h + 0.05, cell_w - 0.16, cell_h - 0.08, 10 if r == 0 else 9, palette.primary if r == 0 else palette.text, r == 0, PP_ALIGN.CENTER if c else PP_ALIGN.LEFT)


def _chart_from_df(df: pd.DataFrame | None, palette: Palette, slide: Any, left: float, top: float, width: float, height: float, chart_kind: str = 'bar') -> None:
    data = CategoryChartData()
    categories: list[str] = []
    series_data: dict[str, list[float]] = {}
    if df is not None and not df.empty:
        category_col = None
        numeric_cols = list(df.select_dtypes(include='number').columns)
        for candidate in df.columns:
            if candidate not in numeric_cols:
                category_col = candidate
                break
        if category_col is None:
            category_col = df.columns[0]
        categories = [str(v) for v in df[category_col].head(6).tolist()]
        if not numeric_cols:
            numeric_cols = [df.columns[-1]]
            df = df.copy()
            df[numeric_cols[0]] = range(1, len(df) + 1)
        for col in numeric_cols[:3]:
            series_data[str(col)] = [float(v) for v in df[col].head(6).fillna(0).tolist()]
    if not categories:
        categories = [f'Q{i}' for i in range(1, 7)]
        series_data = {'Series A': [12, 18, 15, 22, 20, 26], 'Series B': [8, 14, 11, 16, 19, 21]}
    data.categories = categories
    for name, values in series_data.items():
        data.add_series(name, values)
    chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
    if chart_kind == 'line':
        chart_type = XL_CHART_TYPE.LINE_MARKERS
    elif chart_kind == 'pie':
        chart_type = XL_CHART_TYPE.PIE
    elif chart_kind == 'area':
        chart_type = XL_CHART_TYPE.AREA
    chart = slide.shapes.add_chart(chart_type, Inches(left), Inches(top), Inches(width), Inches(height), data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.has_title = False
    if chart.series:
        try:
            chart.series[0].format.fill.solid()
            chart.series[0].format.fill.fore_color.rgb = palette.primary
            if len(chart.series) > 1:
                chart.series[1].format.fill.solid()
                chart.series[1].format.fill.fore_color.rgb = palette.accent
        except Exception:
            pass


def _render_cover(slide: Any, palette: Palette, plan: dict[str, Any], page: dict[str, Any], total_pages: int) -> None:
    hero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(2.05))
    _set_fill(hero, palette.primary)
    _set_line(hero)
    ribbon = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.05), Inches(13.333), Inches(0.18))
    _set_fill(ribbon, palette.accent)
    _set_line(ribbon)
    _add_text(slide, 'Q3 BUSINESS REVIEW', 0.82, 0.42, 3.2, 0.18, 9, RGBColor(220, 232, 245), True)
    _add_text(slide, plan['title'], 0.78, 0.83, 9.0, 0.76, 34, RGBColor(255, 255, 255), True)
    _add_text(slide, plan['subtitle'], 0.82, 1.58, 8.0, 0.28, 12, RGBColor(218, 229, 239))
    _add_badge(slide, plan.get('style', 'style'), 10.15, 0.4, 2.0, palette.accent2, RGBColor(255, 255, 255))
    _add_card(slide, 0.8, 2.9, 7.4, 2.9, palette.paper, palette.line)
    _add_text(slide, '\u672c\u6b21\u7ae0\u8282', 1.08, 3.18, 1.8, 0.24, 15, palette.primary, True)
    sections = plan.get('sections', [])[:5] or ['\u4e1a\u7ee9\u6982\u51b5', '\u7528\u6237\u589e\u957f', '\u95ee\u9898\u5206\u6790', '\u4e0b\u6708\u89c4\u5212']
    for idx, step in enumerate(sections, start=1):
        top = 3.66 + (idx - 1) * 0.38
        _add_text(slide, f'0{idx}', 1.08, top, 0.4, 0.18, 9, palette.accent, True)
        _add_text(slide, step, 1.65, top - 0.01, 5.6, 0.18, 12, palette.text)
    _add_card(slide, 8.5, 2.9, 3.95, 2.9, palette.soft, palette.line)
    _add_text(slide, '\u6c47\u62a5\u91cd\u70b9', 8.82, 3.18, 2.6, 0.24, 17, palette.primary, True)
    _add_multiline(slide, [
        '\u7ed3\u8bba\u5148\u884c',
        '\u6570\u636e\u652f\u6491',
        '\u95ee\u9898\u5f52\u56e0',
        '\u884c\u52a8\u843d\u5730',
    ], 8.84, 3.72, 2.95, 1.0, palette.text, 11)
    _tiny_bar_chart(slide, palette, 9.0, 4.28)
    _footer(slide, palette, 1, total_pages)


def _render_agenda(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _add_text(slide, '目录', 0.78, 0.5, 2.2, 0.52, 26, palette.primary, True)
    _add_text(slide, 'Agenda / Content Map', 0.82, 1.08, 3.0, 0.2, 10, palette.muted, True)
    timeline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.05), Inches(2.04), Inches(11.05), Inches(0.03))
    _set_fill(timeline, palette.line)
    _set_line(timeline)
    sections = page.get('bullets') or []
    if not sections:
        sections = ['背景', '核心内容', '结论']
    card_width = max(1.55, 10.7 / max(1, len(sections)))
    for idx, section in enumerate(sections, start=1):
        left = 0.92 + (idx - 1) * card_width
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.18), Inches(1.8), Inches(0.46), Inches(0.46))
        _set_fill(node, palette.accent if idx % 2 else palette.secondary)
        _set_line(node)
        _add_text(slide, str(idx), left + 0.18, 1.89, 0.46, 0.12, 9, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        _add_card(slide, left, 2.46, card_width - 0.15, 1.7, palette.paper, palette.line)
        _add_text(slide, section, left + 0.16, 2.72, card_width - 0.46, 0.3, 15, palette.primary, True)
        _add_text(slide, '关键结论 · 数据支撑 · 行动建议', left + 0.16, 3.42, card_width - 0.36, 0.26, 9, palette.muted)
    _footer(slide, palette, 2, total_pages)


def _render_section(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.72, 0.48, 9.5, 0.5, 24, palette.primary, True)
    _add_text(slide, page.get('subtitle', ''), 0.76, 1.03, 5.0, 0.22, 11, palette.muted)
    _add_card(slide, 0.82, 1.75, 11.6, 3.75, palette.paper, palette.line)
    _add_text(slide, '章节导语', 1.08, 2.0, 1.5, 0.22, 15, palette.primary, True)
    _add_multiline(slide, [
        '本页用于切换叙事节奏，承接上一章节结论，',
        '为下一页的分析、图表或行动计划建立语境。',
        '它不是纯装饰页，而是有明确信息层次的过渡页。',
    ], 1.08, 2.42, 5.0, 1.1, palette.text, 15)
    _add_card(slide, 7.0, 2.0, 4.95, 2.95, palette.soft, palette.line)
    _add_text(slide, '关键提示', 7.28, 2.32, 1.5, 0.2, 15, palette.primary, True)
    _add_multiline(slide, ['结论先行', '重点明确', '留白充足', '视觉简洁'], 7.32, 2.7, 2.3, 1.1, palette.text, 18)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_content(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.48, 9.1, 0.5, 24, palette.primary, True)
    _add_badge(slide, page.get('section', '\u5185\u5bb9\u9875'), 10.8, 0.56, 1.45, palette.soft, palette.primary)
    _add_card(slide, 0.78, 1.42, 8.15, 5.05, palette.paper, palette.line)
    _add_multiline(slide, page.get('bullets', []), 1.06, 1.78, 7.4, 3.8, palette.text, 16)
    _add_card(slide, 9.22, 1.42, 3.35, 5.05, palette.soft, palette.line)
    _add_text(slide, '\u6c47\u62a5\u8981\u70b9', 9.52, 1.78, 1.8, 0.24, 16, palette.primary, True)
    notes = [
        '\u7ed3\u8bba\u5148\u884c',
        '\u6570\u636e\u652f\u6491',
        '\u52a8\u4f5c\u95ed\u73af',
    ]
    _add_multiline(slide, notes, 9.58, 2.32, 2.0, 1.0, palette.text, 19)
    for idx, metric in enumerate(page.get('metrics', [])[:3]):
        top = 4.0 + idx * 0.72
        _metric_card(slide, palette, metric, 9.52, top, 2.58)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_data_page(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int, df: pd.DataFrame | None) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.46, 9.4, 0.5, 23, palette.primary, True)
    _add_text(slide, page.get('subtitle', ''), 10.15, 0.58, 2.2, 0.18, 10, palette.muted, True, PP_ALIGN.RIGHT)
    metrics = page.get('metrics', []) or [
        {'label': '核心指标', 'value': '92.4%', 'trend': '同比提升'},
        {'label': '关键增长', 'value': '18.7%', 'trend': '环比向上'},
        {'label': '关注风险', 'value': '3项', 'trend': '需持续跟踪'},
    ]
    for idx, metric in enumerate(metrics[:3]):
        _metric_card(slide, palette, metric, 0.78 + idx * 2.83, 1.3, 2.55)
    _add_card(slide, 0.78, 2.48, 7.4, 4.0, palette.paper, palette.line)
    _add_text(slide, '趋势图', 1.08, 2.72, 1.2, 0.2, 15, palette.primary, True)
    chart_kind = 'line' if page.get('layout_type') in {'dashboard', 'kpi'} else 'bar'
    _chart_from_df(df, palette, slide, 1.0, 3.02, 6.9, 2.9, chart_kind=chart_kind)
    _add_card(slide, 8.45, 2.48, 4.1, 4.0, palette.soft, palette.line)
    _add_text(slide, '数据解读', 8.75, 2.72, 1.5, 0.22, 15, palette.primary, True)
    narrative = page.get('bullets', [])[:4]
    if not narrative:
        narrative = ['围绕数据建立结论', '拆解变化原因', '明确下一步动作']
    _add_multiline(slide, narrative, 8.8, 3.12, 3.1, 1.8, palette.text, 14)
    if df is not None and not df.empty:
        preview = df.head(4)
        table_rows = [list(preview.columns[:4])] + preview.iloc[:, :4].astype(str).values.tolist()
    else:
        table_rows = [['指标', '当前', '目标', '变化'], ['A', '92', '95', '+3'], ['B', '76', '80', '+4']]
    _simple_table(slide, palette, table_rows, 8.7, 4.82, 3.2, 1.22)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_table(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int, df: pd.DataFrame | None) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.46, 9.4, 0.5, 24, palette.primary, True)
    preview = [['字段', '值', '字段', '值']]
    if df is not None and not df.empty:
        values = df.head(4).astype(str)
        cols = list(values.columns)
        for idx, row in values.iterrows():
            row_values = [str(row[col]) for col in cols[:4]]
            while len(row_values) < 4:
                row_values.append('')
            preview.append(row_values)
    else:
        preview += [
            ['门店', '销售额', '同比', '排名'],
            ['上海', '120万', '+18%', '1'],
            ['北京', '103万', '+14%', '2'],
            ['深圳', '98万', '+11%', '3'],
        ]
    _simple_table(slide, palette, preview, 0.84, 1.42, 8.0, 4.9)
    _add_card(slide, 9.18, 1.42, 3.36, 4.9, palette.soft, palette.line)
    _add_text(slide, '表格解读', 9.48, 1.76, 1.6, 0.2, 15, palette.primary, True)
    _add_multiline(slide, page.get('bullets', [])[:4] or ['适合呈现明细、清单、报价、排名和对比信息。', '结合筛选和汇总能直接用于汇报或审批。'], 9.5, 2.18, 2.6, 2.2, palette.text, 14)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_timeline(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.46, 9.4, 0.5, 24, palette.primary, True)
    _add_card(slide, 0.84, 1.75, 11.55, 3.8, palette.paper, palette.line)
    y = 3.15
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(y), Inches(10.2), Inches(0.04)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = palette.line
    slide.shapes[-1].line.transparency = 100000
    steps = page.get('bullets', [])[:5] or ['阶段一', '阶段二', '阶段三', '阶段四']
    for idx, item in enumerate(steps, start=1):
        x = 1.4 + (idx - 1) * 2.25
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.95), Inches(0.36), Inches(0.36))
        _set_fill(dot, palette.accent if idx % 2 else palette.secondary)
        _set_line(dot)
        _add_text(slide, str(idx), x + 0.06, 3.02, 0.22, 0.1, 8, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        _add_card(slide, x - 0.32, 1.95 if idx % 2 else 3.42, 1.55, 0.96, palette.soft, palette.line)
        _add_text(slide, item, x - 0.22, (2.22 if idx % 2 else 3.67), 1.3, 0.42, 11, palette.text, True)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_process(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.46, 9.4, 0.5, 24, palette.primary, True)
    _add_card(slide, 0.9, 1.5, 11.3, 4.8, palette.paper, palette.line)
    steps = page.get('bullets', [])[:4] or ['输入', '处理', '输出', '复盘']
    for idx, step in enumerate(steps, start=1):
        top = 1.88 + (idx - 1) * 1.0
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.25), Inches(top), Inches(1.0), Inches(0.42))
        _set_fill(pill, palette.primary if idx % 2 else palette.accent)
        _set_line(pill)
        _add_text(slide, f'Step {idx}', 1.32, top + 0.08, 0.86, 0.16, 10, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        _add_text(slide, step, 2.6, top + 0.02, 8.6, 0.2, 14, palette.text, True)
        if idx < len(steps):
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(1.63), Inches(top + 0.45), Inches(0.25), Inches(0.24))
            _set_fill(arrow, palette.line)
            _set_line(arrow)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_comparison(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.46, 9.4, 0.5, 24, palette.primary, True)
    _add_card(slide, 0.84, 1.55, 5.85, 4.8, palette.paper, palette.line)
    _add_card(slide, 6.64, 1.55, 5.85, 4.8, palette.soft, palette.line)
    _add_text(slide, '方案 A', 1.1, 1.84, 1.2, 0.2, 16, palette.primary, True)
    _add_text(slide, '方案 B', 6.92, 1.84, 1.2, 0.2, 16, palette.primary, True)
    left_bullets = page.get('bullets', [])[:3] or ['优点：高效', '缺点：成本较高', '适用：成熟业务']
    right_bullets = ['优势：灵活', '风险：需要协同', '适用：快速迭代']
    _add_multiline(slide, left_bullets, 1.1, 2.28, 4.8, 1.8, palette.text, 14)
    _add_multiline(slide, right_bullets, 6.94, 2.28, 4.8, 1.8, palette.text, 14)
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.16), Inches(1.85), Inches(0.04), Inches(3.95))
    _set_fill(divider, palette.line)
    _set_line(divider)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_strategy(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.46, 9.4, 0.5, 24, palette.primary, True)
    labels = ['Strength', 'Weakness', 'Opportunity', 'Threat']
    fills = [palette.soft, palette.soft2, palette.paper, palette.soft]
    positions = [(0.92, 1.62), (6.72, 1.62), (0.92, 4.0), (6.72, 4.0)]
    for (left, top), label, fill in zip(positions, labels, fills):
        _add_card(slide, left, top, 5.3, 1.95, fill, palette.line)
        _add_text(slide, label, left + 0.22, top + 0.18, 1.8, 0.18, 14, palette.primary, True)
        _add_multiline(slide, ['输入要点', '关键判断', '行动建议'], left + 0.24, top + 0.5, 3.8, 0.95, palette.text, 12)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_people(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.46, 9.4, 0.5, 24, palette.primary, True)
    for idx in range(4):
        left = 0.92 + idx * 3.1
        _add_card(slide, left, 1.72, 2.7, 3.65, palette.paper, palette.line)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.75), Inches(2.02), Inches(1.15), Inches(1.15))
        _set_fill(circle, palette.secondary if idx % 2 else palette.accent)
        _set_line(circle)
        _add_text(slide, f'P{idx + 1}', left + 0.92, 2.37, 0.6, 0.18, 12, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
        _add_text(slide, f'成员 {idx + 1}', left + 0.36, 3.42, 1.95, 0.18, 14, palette.primary, True, PP_ALIGN.CENTER)
        _add_text(slide, '角色 / 简介 / 联系方式', left + 0.2, 3.74, 2.28, 0.18, 9, palette.muted, align=PP_ALIGN.CENTER)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_product(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.46, 9.4, 0.5, 24, palette.primary, True)
    _add_card(slide, 0.84, 1.52, 6.0, 4.75, palette.paper, palette.line)
    _add_card(slide, 7.0, 1.52, 5.6, 1.44, palette.soft, palette.line)
    _add_card(slide, 7.0, 3.08, 5.6, 1.44, palette.soft2, palette.line)
    _add_card(slide, 7.0, 4.64, 5.6, 1.44, palette.paper, palette.line)
    _add_text(slide, '产品价值', 1.08, 1.82, 1.5, 0.2, 16, palette.primary, True)
    _add_multiline(slide, page.get('bullets', [])[:4] or ['功能亮点', '解决问题', '商业价值'], 1.12, 2.22, 4.8, 2.2, palette.text, 14)
    _add_text(slide, '版本路线图', 7.28, 1.86, 1.8, 0.2, 15, palette.primary, True)
    _tiny_bar_chart(slide, palette, 7.48, 2.22)
    _add_text(slide, '参数 / 规格 / 说明', 7.28, 3.38, 2.5, 0.2, 15, palette.primary, True)
    _add_multiline(slide, ['尺寸', '性能', '价格', '交付'], 7.34, 3.7, 2.0, 0.8, palette.text, 13)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_gallery(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.46, 9.4, 0.5, 24, palette.primary, True)
    for idx in range(6):
        row = idx // 3
        col = idx % 3
        left = 0.92 + col * 4.0
        top = 1.6 + row * 2.35
        _add_card(slide, left, top, 3.3, 2.0, palette.paper, palette.line)
        image = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.16), Inches(top + 0.16), Inches(3.0), Inches(1.15))
        _set_fill(image, palette.soft if idx % 2 == 0 else palette.soft2)
        _set_line(image, palette.line)
        _add_text(slide, f'Image {idx + 1}', left + 0.2, top + 1.42, 1.0, 0.18, 11, palette.primary, True)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_map(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    _title_band(slide, palette)
    _add_text(slide, page['title'], 0.76, 0.46, 9.4, 0.5, 24, palette.primary, True)
    _add_card(slide, 0.92, 1.52, 7.5, 4.8, palette.paper, palette.line)
    _add_card(slide, 8.62, 1.52, 3.9, 4.8, palette.soft, palette.line)
    map_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.86), Inches(6.95), Inches(3.96))
    _set_fill(map_panel, palette.soft)
    _set_line(map_panel, palette.line)
    for idx in range(8):
        x = 1.5 + (idx % 4) * 1.55
        y = 2.15 + (idx // 4) * 1.5
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.18), Inches(0.18))
        _set_fill(dot, palette.accent if idx % 2 else palette.secondary)
        _set_line(dot)
    _add_text(slide, '区域分布', 8.92, 1.82, 1.4, 0.2, 15, palette.primary, True)
    _add_multiline(slide, page.get('bullets', [])[:4] or ['覆盖范围', '区域热点', '资源分布'], 8.96, 2.24, 2.9, 2.0, palette.text, 14)
    _footer(slide, palette, page['page_no'], total_pages)


def _render_ending(slide: Any, palette: Palette, page: dict[str, Any], total_pages: int) -> None:
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.08), Inches(13.333), Inches(1.9))
    _set_fill(block, palette.primary)
    _set_line(block)
    _add_text(slide, page['title'], 0.0, 2.56, 13.333, 0.52, 32, RGBColor(255, 255, 255), True, PP_ALIGN.CENTER)
    _add_text(slide, page.get('subtitle', 'Thank You'), 0.0, 3.2, 13.333, 0.24, 12, RGBColor(219, 228, 240), True, PP_ALIGN.CENTER)
    _add_text(slide, '\u4e0b\u4e00\u6b65\u805a\u7126 Q4 \u884c\u52a8\u548c\u5173\u952e\u98ce\u9669\u8ddf\u8fdb\u3002', 0.0, 4.35, 13.333, 0.2, 11, palette.muted, align=PP_ALIGN.CENTER)
    _footer(slide, palette, total_pages, total_pages)


LAYOUT_RENDERERS = {
    'cover_slide': _render_cover,
    'agenda_slide': _render_agenda,
    'section_slide': _render_section,
    'content_slide': _render_content,
    'data_analysis_slide': _render_data_page,
    'chart_slide': _render_data_page,
    'table_slide': _render_table,
    'timeline_slide': _render_timeline,
    'process_slide': _render_process,
    'comparison_slide': _render_comparison,
    'relation_slide': _render_people,
    'strategy_slide': _render_strategy,
    'planning_slide': _render_timeline,
    'people_slide': _render_people,
    'product_slide': _render_product,
    'gallery_slide': _render_gallery,
    'map_slide': _render_map,
    'dashboard_slide': _render_data_page,
    'kpi_slide': _render_data_page,
    'ending_slide': _render_ending,
}


def render_smart_deck(plan: dict[str, Any], excel_df: pd.DataFrame | None = None) -> bytes:
    design = plan.get('design_system', {})
    prs = _setup_prs()
    blank = prs.slide_layouts[6]
    pages = plan.get('pages', [])
    total_pages = len(pages)
    if not total_pages:
        raise ValueError('plan must include pages')
    for index, page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(blank)
        page_design = dict(design)
        candidate = page.get('selected_candidate') or {}
        if isinstance(candidate, dict):
            for key in ('primary_color', 'secondary_color', 'background_color', 'dark_or_light', 'accent_color'):
                if candidate.get(key):
                    page_design[key] = candidate[key]
        palette = _palette_from_design(page_design)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = palette.bg
        layout_type = page.get('layout_type', 'content_slide')
        renderer = LAYOUT_RENDERERS.get(layout_type, _render_content)
        try:
            if renderer is _render_cover:
                renderer(slide, palette, plan, page, total_pages)
            elif renderer in {_render_data_page, _render_table}:
                renderer(slide, palette, page, total_pages, excel_df)
            else:
                renderer(slide, palette, page, total_pages)
        except TypeError:
            if renderer is _render_cover:
                renderer(slide, palette, plan, page, total_pages)
            elif renderer in {_render_data_page, _render_table}:
                renderer(slide, palette, page, total_pages, excel_df)
            else:
                renderer(slide, palette, page, total_pages)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
