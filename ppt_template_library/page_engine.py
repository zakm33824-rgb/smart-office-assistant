from __future__ import annotations

import json
import re
from datetime import date
from pathlib import Path
from typing import Any, Iterable

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .page_models import AssemblyRequest, ComponentRecord, DesignSystem, ExcelProfile, PageRecord, PageSearchFilters, normalize_tokens
from .paths import ensure_library_structure
from .preview import build_contact_sheet, create_component_preview, create_page_preview
from .slide_storage import initialize_slide_storage, load_component_catalog, load_slide_catalog, load_slide_summary, save_assembly_run, save_component_catalog, save_slide_catalog

INDUSTRY_MAP = [
    (r'新能源汽车|汽车', 'automotive'),
    (r'金融|银行|保险|投资|证券', 'finance'),
    (r'房地产|建筑|工程', 'construction'),
    (r'制造|工业', 'manufacturing'),
    (r'新能源|光伏|电力', 'energy'),
    (r'互联网|软件|人工智能|大数据|云计算|网络安全', 'technology'),
    (r'电子商务|零售|物流|供应链', 'commerce'),
    (r'医疗|医药|生物', 'healthcare'),
    (r'教育|学校|大学|科研', 'education'),
    (r'政府|公共事业', 'government'),
    (r'农业|食品|餐饮', 'consumer'),
    (r'旅游|酒店|航空|交通', 'travel'),
    (r'文化|艺术|音乐|影视|游戏|体育', 'creative'),
    (r'时尚|美容|环保|公益|法律|咨询|人力资源|市场营销|广告|媒体|创业|个人展示', 'business'),
]

SCENARIO_MAP = [
    (r'年度总结|年度报告|年报|复盘', 'annual_report'),
    (r'季度总结|季度报告|Q[1-4]', 'quarterly_report'),
    (r'月度总结|月报', 'monthly_report'),
    (r'周报', 'weekly_report'),
    (r'日报', 'daily_report'),
    (r'工作总结|总结', 'work_summary'),
    (r'项目汇报|项目报告', 'project_report'),
    (r'项目计划|路线图|里程碑', 'project_plan'),
    (r'商业计划|融资计划|创业计划', 'business_plan'),
    (r'产品发布|发布会|产品介绍', 'product_launch'),
    (r'市场分析|行业分析|竞争分析|SWOT', 'market_analysis'),
    (r'数据分析|财务报告|销售报告|运营报告', 'data_analysis'),
    (r'研究报告|调查报告|咨询报告|学术报告|科研汇报', 'research_report'),
    (r'培训|教学|答辩|会议演示|领导汇报|政府汇报|党建', 'presentation'),
    (r'简历|作品集|摄影|旅行|婚礼|生日|节日|纪念', 'personal'),
]

STYLE_MAP = [
    (r'简约|极简', 'minimal'),
    (r'商务|正式', 'business'),
    (r'现代|高级', 'modern'),
    (r'科技|AI|互联网|云', 'technology'),
    (r'学术|教育|科研', 'academic'),
    (r'创意|活泼|插画|卡通', 'creative'),
    (r'政府|红色政务', 'formal'),
    (r'绿色环保|环保', 'eco'),
    (r'中国风|国潮|日式|韩式|欧美|复古|怀旧', 'cultural'),
    (r'赛博朋克|蒸汽波|孟菲斯|玻璃拟态|毛玻璃|黑金|蓝色商务', 'trend'),
]

COLOR_MAP = [
    (r'蓝', ('blue', 'teal')),
    (r'红', ('red', 'gold')),
    (r'绿', ('green', 'teal')),
    (r'紫', ('purple', 'blue')),
    (r'橙', ('orange', 'gold')),
    (r'黑', ('dark', 'gold')),
]

COMPONENT_LABELS = {
    'title': '标题', 'subtitle': '副标题', 'text': '文本', 'card': '卡片', 'image': '图片', 'label': '标签', 'number': '数字', 'metric_card': 'KPI卡片', 'chart': '图表', 'table': '表格',
    'timeline': '时间节点', 'process': '流程步骤', 'node': '关系节点', 'matrix': '矩阵单元', 'title_block': '标题块', 'subtitle_block': '副标题块', 'hero_panel': '主视觉', 'accent_bar': '强调线', 'glow_panel': '高亮背景',
    'chip_row': '标签行', 'section_number': '目录数字', 'agenda_row': '目录项', 'agenda_label': '目录说明', 'chapter_number': '章节数字', 'chapter_title': '章节标题', 'chapter_rule': '章节线', 'text_block': '正文块',
    'image_panel': '图片面板', 'note_card': '说明卡片', 'kpi_card': 'KPI卡片', 'chart_panel': '图表面板', 'insight_card': '洞察卡片', 'headline_block': '标题区', 'chart_canvas': '图表画布', 'legend_block': '图例',
    'insight_block': '洞察说明', 'table_block': '表格块', 'header_block': '表头', 'summary_chip': '摘要标签', 'timeline_node': '时间节点', 'milestone_card': '里程碑卡片', 'connector': '连接线', 'process_step': '流程步骤',
    'arrow_connector': '箭头', 'compare_left': '左对比', 'compare_right': '右对比', 'divider': '分割线', 'center_node': '中心节点', 'child_node': '子节点', 'swot_cell': 'SWOT单元格', 'strategy_label': '战略标签', 'summary_box': '总结框',
    'gantt_row': '甘特行', 'task_bar': '任务条', 'milestone': '里程碑', 'profile_card': '人物卡', 'portrait': '头像', 'name_block': '姓名块', 'feature_card': '功能卡', 'product_image': '产品图', 'spec_block': '参数块',
    'image_frame': '图片框', 'caption_block': '说明文字', 'gallery_card': '画廊卡', 'map_panel': '地图面板', 'region_tag': '区域标签', 'metric_chip': '指标标签', 'thanks_block': '致谢标题', 'contact_block': '联系卡', 'footer_rule': '底部线',
}


def _slug(value: str) -> str:
    out: list[str] = []
    dash = False
    for ch in value.lower():
        if ch.isalnum():
            out.append(ch)
            dash = False
        elif not dash:
            out.append('-')
            dash = True
    return ''.join(out).strip('-') or 'item'


def _norm(text: str) -> str:
    return re.sub(r'\s+', ' ', (text or '').strip())


def _match(text: str, patterns: list[tuple[str, str]], default: str) -> str:
    for pattern, value in patterns:
        if re.search(pattern, text, re.IGNORECASE):
            return value
    return default


def _contains(text: str, pattern: str) -> bool:
    return bool(re.search(pattern, text, re.IGNORECASE))


def _infer_title(prompt: str) -> str:
    title = re.split(r'(?:包含|包括|围绕|结合|实现|并|，|。|；|;)', _norm(prompt), maxsplit=1)[0]
    title = re.sub(r'^(请|帮我|麻烦|生成|制作|输出|做一份|做个|帮忙)?', '', title).strip(' ，,。:：')
    if not title:
        return '智能PPT'
    if len(title) <= 6 and '报告' not in title and 'PPT' not in title.upper():
        title = f'{title}报告'
    return title


def _infer_sections(prompt: str) -> list[str]:
    text = _norm(prompt)
    if _contains(text, r'年度|年报|总结|复盘'):
        return ['业绩概览', '增长分析', '问题洞察', '行动计划']
    if _contains(text, r'项目|计划|路线图|里程碑'):
        return ['项目背景', '推进进度', '风险资源', '下一步计划']
    if _contains(text, r'产品|发布|介绍'):
        return ['产品背景', '核心亮点', '应用场景', '发布计划']
    if _contains(text, r'学术|研究|论文|答辩'):
        return ['研究背景', '方法与实验', '结果分析', '结论展望']
    if _contains(text, r'市场|营销|广告'):
        return ['市场背景', '渠道表现', '用户反馈', '营销计划']
    if _contains(text, r'财务|销售|运营|数据'):
        return ['核心指标', '趋势分析', '结构拆解', '管理建议']
    return ['核心内容', '关键发现', '问题分析', '行动建议']


def _infer_industry(text: str) -> str:
    return _match(text, INDUSTRY_MAP, 'business')


def _infer_scenario(text: str) -> str:
    return _match(text, SCENARIO_MAP, 'report')


def _infer_style(text: str, style_name: str) -> str:
    found = _match(text, STYLE_MAP, '')
    return found or _match(style_name, STYLE_MAP, 'business')


def _infer_colors(text: str, style: str) -> tuple[str, str]:
    for pattern, colors in COLOR_MAP:
        if _contains(text, pattern):
            return colors
    return ('dark', 'gold') if style == 'formal' else ('blue', 'teal')


def _design_system(style: str) -> DesignSystem:
    mapping = {
        'technology': ('#123C7A', '#2F80ED', '#14B8A6', '#0F172A'),
        'academic': ('#4A427C', '#996F2C', '#D7DEE8', '#F8F8F5'),
        'creative': ('#205375', '#D96C4D', '#F4A946', '#FAFBFF'),
        'minimal': ('#1B3756', '#94A3B8', '#00A89D', '#F7F9FC'),
        'formal': ('#205375', '#6B7280', '#F4A946', '#FCFCFD'),
        'eco': ('#1F5B4E', '#22C55E', '#A3E635', '#F6FBF7'),
    }
    primary, secondary, accent, background = mapping.get(style, ('#1B3756', '#00A89D', '#F4A946', '#F7F9FC'))
    return DesignSystem(primary, secondary, accent, background, 'Microsoft YaHei', 'Microsoft YaHei', 'Arial', 28, 16, 14, 10, 12, 'comfortable', (primary, accent, secondary), 'line', 'clean', 'dark' if style == 'technology' else 'light')


def _page_tags(bp: dict[str, Any]) -> tuple[str, ...]:
    values = [bp['slide_type'], bp['slide_subtype'], bp['industry'], bp['scenario'], bp['style'], bp['layout_type'], *bp['tags']]
    return tuple(dict.fromkeys(normalize_tokens(values)))


def _embedding(bp: dict[str, Any]) -> str:
    return json.dumps({'keywords': list(_page_tags(bp)), 'slide_type': bp['slide_type'], 'layout_type': bp['layout_type']}, ensure_ascii=False)


def _score(bp: dict[str, Any]) -> int:
    return max(0, min(100, round(bp['design_score'] * 0.28 + bp['layout_score'] * 0.22 + bp['color_score'] * 0.16 + bp['usability_score'] * 0.18 + bp['modern_score'] * 0.16)))


def _bp(slide_type: str, slide_subtype: str, title: str, subtitle: str, layout_type: str, tags: tuple[str, ...], components: tuple[tuple[str, str], ...], **kwargs: Any) -> dict[str, Any]:
    base = {
        'slide_type': slide_type, 'slide_subtype': slide_subtype, 'title': title, 'subtitle': subtitle, 'layout_type': layout_type,
        'industry': 'business', 'scenario': 'report', 'style': 'business', 'primary_color': '#1B3756', 'secondary_color': '#00A89D',
        'background_color': '#F7F9FC', 'dark_or_light': 'light', 'aspect_ratio': '16:9', 'text_density': 0.2, 'image_density': 0.1,
        'chart_count': 0, 'table_count': 0, 'shape_count': 12, 'icon_count': 1, 'image_count': 0, 'text_box_count': 5,
        'has_chart': False, 'has_table': False, 'has_timeline': False, 'has_process': False, 'has_map': False, 'has_people': False,
        'has_infographic': False, 'has_animation': False, 'editable_level': 'high', 'design_score': 90, 'layout_score': 90,
        'color_score': 88, 'usability_score': 91, 'modern_score': 88, 'tags': tags, 'components': components,
    }
    base.update(kwargs)
    return base


PAGE_BLUEPRINTS = [
    _bp('cover', 'business_cover', '商务封面', '适合年度总结、经营分析与企业汇报开场', 'hero_split', ('cover', 'business', 'annual_report', 'corporate'), (('title_block', 'title'), ('subtitle_block', 'subtitle'), ('hero_panel', 'image'), ('accent_bar', 'decorative')), scenario='annual_report', design_score=93, layout_score=91, color_score=90, usability_score=93, modern_score=90),
    _bp('cover', 'technology_cover', '科技封面', '面向 AI、互联网、云计算与技术发布场景', 'dark_hero', ('cover', 'technology', 'ai', 'modern'), (('title_block', 'title'), ('subtitle_block', 'subtitle'), ('glow_panel', 'image'), ('chip_row', 'label')), industry='technology', scenario='product_launch', style='technology', primary_color='#123C7A', secondary_color='#2F80ED', background_color='#0F172A', dark_or_light='dark', text_density=0.11, image_density=0.35, shape_count=16, icon_count=4, image_count=1, text_box_count=5, design_score=95, layout_score=93, color_score=94, usability_score=92, modern_score=96),
    _bp('agenda', 'number_agenda', '数字目录', '使用编号与分组概览整套结构', 'numbered_list', ('agenda', 'numbered', 'business', 'outline'), (('section_number', 'number'), ('agenda_row', 'list_item'), ('agenda_label', 'text')), industry='business', scenario='project_report', style='business', primary_color='#205375', secondary_color='#F4A946', text_density=0.2, shape_count=18, icon_count=3, text_box_count=7, design_score=90, layout_score=92, color_score=88, usability_score=94, modern_score=88),
    _bp('section', 'chapter_break', '章节页', '适合分段、过渡与强调重点章节', 'big_title', ('section', 'transition', 'minimal', 'chapter'), (('chapter_number', 'number'), ('chapter_title', 'title'), ('chapter_rule', 'decorative')), style='minimal', scenario='report_section', primary_color='#4A427C', secondary_color='#D96C4D', background_color='#F8F8F5', text_density=0.1, image_density=0.12, shape_count=8, design_score=91, layout_score=93, color_score=89, usability_score=92, modern_score=89),
    _bp('content', 'text_image_layout', '图文内容页', '左右结构，适合背景说明与结论展示', 'two_column', ('content', 'two_column', 'text', 'image'), (('title_block', 'title'), ('text_block', 'body'), ('image_panel', 'image'), ('note_card', 'card')), style='modern', scenario='work_summary', primary_color='#1F4E79', secondary_color='#00A89D', background_color='#FBFCFE', text_density=0.42, image_density=0.22, shape_count=14, text_box_count=6, design_score=89, layout_score=90, color_score=88, usability_score=92, modern_score=90),
    _bp('data_analysis', 'kpi_dashboard', 'KPI 数据页', '适合核心指标概览与经营数据展示', 'dashboard_cards', ('kpi', 'dashboard', 'data', 'finance'), (('kpi_card', 'metric_card'), ('chart_panel', 'chart'), ('insight_card', 'note_card'), ('headline_block', 'title')), industry='finance', scenario='annual_report', style='business', primary_color='#1B3756', secondary_color='#2F80ED', background_color='#F6FAFF', text_density=0.24, image_density=0.1, chart_count=2, shape_count=20, icon_count=3, text_box_count=8, has_chart=True, has_infographic=True, design_score=94, layout_score=93, color_score=92, usability_score=94, modern_score=93),
    _bp('chart', 'trend_chart', '趋势图页', '适合销售、增长与时间序列分析', 'line_chart', ('chart', 'trend', 'growth', 'technology'), (('chart_canvas', 'chart'), ('legend_block', 'legend'), ('insight_block', 'note_card')), industry='technology', scenario='sales_report', style='technology', primary_color='#123C7A', secondary_color='#14B8A6', background_color='#F6FAFF', text_density=0.18, image_density=0.08, chart_count=2, shape_count=18, icon_count=2, has_chart=True, has_infographic=True, design_score=93, layout_score=92, color_score=94, usability_score=93, modern_score=95),
    _bp('timeline', 'roadmap_timeline', '时间轴页', '适合历程、路线图、里程碑与节点规划', 'horizontal_timeline', ('timeline', 'roadmap', 'milestone', 'project'), (('timeline_node', 'timeline'), ('milestone_card', 'card'), ('connector', 'decorative')), scenario='project_plan', style='modern', primary_color='#4A427C', secondary_color='#00A89D', background_color='#F8FBFF', text_density=0.24, has_timeline=True, shape_count=18, icon_count=3, text_box_count=6, has_infographic=True, design_score=92, layout_score=94, color_score=90, usability_score=93, modern_score=91),
    _bp('process', 'workflow_process', '流程页', '适合业务流程、项目流程与工作流', 'arrow_flow', ('process', 'workflow', 'flow', 'operation'), (('process_step', 'process'), ('arrow_connector', 'decorative'), ('note_card', 'card')), industry='business', scenario='project_plan', style='business', primary_color='#3B6E8F', secondary_color='#F59E0B', background_color='#FBFAF7', text_density=0.2, has_process=True, shape_count=20, icon_count=4, text_box_count=6, has_infographic=True, design_score=91, layout_score=93, color_score=88, usability_score=94, modern_score=89),
    _bp('comparison', 'before_after', '对比页', '适合方案对比、前后对比与竞品分析', 'two_panel_compare', ('comparison', 'before_after', 'analysis', 'consulting'), (('compare_left', 'card'), ('compare_right', 'card'), ('divider', 'decorative')), industry='consulting', scenario='competitive_analysis', style='modern', primary_color='#205375', secondary_color='#EF4444', background_color='#FAFBFF', text_density=0.25, image_density=0.1, shape_count=14, icon_count=2, image_count=1, text_box_count=6, has_infographic=True, design_score=91, layout_score=91, color_score=89, usability_score=93, modern_score=90),
    _bp('team', 'team_profile', '团队介绍页', '适合组织成员、专家阵容与讲师介绍', 'profile_cards', ('team', 'people', 'profile', 'organization'), (('profile_card', 'card'), ('portrait', 'image'), ('name_block', 'text')), scenario='company_profile', style='modern', primary_color='#1F4E79', secondary_color='#F472B6', background_color='#FBFCFF', text_density=0.2, image_density=0.36, image_count=4, shape_count=18, icon_count=2, text_box_count=8, has_people=True, has_infographic=True, design_score=90, layout_score=91, color_score=89, usability_score=92, modern_score=92),
    _bp('product', 'product_highlight', '产品展示页', '适合产品功能、优势、参数与路线图', 'feature_grid', ('product', 'feature', 'showcase', 'technology'), (('feature_card', 'card'), ('product_image', 'image'), ('spec_block', 'text')), industry='technology', scenario='product_intro', style='creative', primary_color='#205375', secondary_color='#10B981', background_color='#FAFBFF', text_density=0.22, image_density=0.3, image_count=2, shape_count=18, icon_count=4, text_box_count=8, has_infographic=True, design_score=92, layout_score=92, color_score=91, usability_score=91, modern_score=93),
    _bp('map', 'business_map', '地图页', '适合区域分布、站点布局与市场覆盖', 'map_dashboard', ('map', 'distribution', 'region', 'logistics'), (('map_panel', 'map'), ('region_tag', 'label'), ('metric_chip', 'metric_card')), industry='logistics', scenario='market_analysis', style='business', primary_color='#205375', secondary_color='#22C55E', background_color='#F8FBFF', text_density=0.22, image_density=0.34, chart_count=1, has_chart=True, has_map=True, has_infographic=True, shape_count=18, icon_count=2, image_count=1, text_box_count=6, design_score=91, layout_score=92, color_score=90, usability_score=91, modern_score=89),
    _bp('ending', 'thanks_end', '结束页', '用于总结、行动计划、联系方式与致谢', 'closing_card', ('ending', 'thanks', 'closing', 'contact'), (('thanks_block', 'title'), ('contact_block', 'card'), ('footer_rule', 'decorative')), scenario='closing', style='business', primary_color='#1B3756', secondary_color='#94A3B8', background_color='#F7F9FC', text_density=0.1, shape_count=8, icon_count=1, text_box_count=4, design_score=90, layout_score=92, color_score=88, usability_score=94, modern_score=88),
]

def build_demo_page_records(root: str | Path = 'ppt_template_library') -> tuple[list[PageRecord], list[ComponentRecord]]:
    paths = ensure_library_structure(root)
    pages: list[PageRecord] = []
    components: list[ComponentRecord] = []
    for idx, bp in enumerate(PAGE_BLUEPRINTS, start=1):
        slide_id = f"demo-{idx:03d}-{_slug(bp['slide_subtype'])}"
        preview = paths.preview / 'slides' / f'{slide_id}.png'
        create_page_preview(bp['title'], bp['subtitle'], bp['slide_type'], bp['slide_subtype'], _score(bp), bp['tags'], preview)
        page = PageRecord(slide_id, 'demo-library', 'demo://synthetic', 'demo://page-library', idx, bp['slide_type'], bp['slide_subtype'], bp['industry'], bp['scenario'], bp['style'], bp['layout_type'], bp['primary_color'], bp['secondary_color'], bp['background_color'], bp['dark_or_light'], bp['aspect_ratio'], bp['text_density'], bp['image_density'], bp['chart_count'], bp['table_count'], bp['shape_count'], bp['icon_count'], bp['image_count'], bp['text_box_count'], bp['has_chart'], bp['has_table'], bp['has_timeline'], bp['has_process'], bp['has_map'], bp['has_people'], bp['has_infographic'], bp['has_animation'], bp['editable_level'], bp['design_score'], bp['layout_score'], bp['color_score'], bp['usability_score'], bp['modern_score'], _score(bp), str(preview), str(preview), '', _embedding(bp), {'title': bp['title'], 'subtitle': bp['subtitle'], 'bullets': list(bp['tags']), 'source': 'synthetic_demo'}, 'demo', _page_tags(bp))
        pages.append(page)
        for cidx, (ctype, csubtype) in enumerate(bp['components'], start=1):
            comp_id = f'{slide_id}-c{cidx:02d}'
            cpreview = paths.preview / 'components' / f'{comp_id}.png'
            create_component_preview(COMPONENT_LABELS.get(csubtype, csubtype), bp['slide_subtype'], ctype, bp['design_score'] - cidx, (bp['slide_type'], bp['style'], csubtype), cpreview)
            components.append(ComponentRecord(comp_id, ctype, csubtype, slide_id, bp['layout_type'], '0.00,0.00,0.80,0.18', bp['style'], bp['primary_color'], 0.80, 0.18, str(cpreview), {'description': COMPONENT_LABELS.get(csubtype, csubtype), 'source_slide_type': bp['slide_type']}, 'demo', (bp['slide_type'], bp['style'], csubtype)))
    return pages, components


def seed_demo_page_library(root: str | Path = 'ppt_template_library', refresh: bool = False) -> dict[str, Any]:
    initialize_slide_storage(root)
    if not refresh and not load_slide_catalog(root).empty and not load_component_catalog(root).empty:
        return {**load_slide_summary(root), 'seeded': False}
    pages, comps = build_demo_page_records(root)
    save_slide_catalog(pages, root)
    save_component_catalog(comps, root)
    return {**load_slide_summary(root), 'seeded': True, 'page_records': len(pages), 'component_records': len(comps)}


def _collect_slide_text(slide: Any) -> str:
    parts = []
    for shape in slide.shapes:
        try:
            if getattr(shape, 'has_text_frame', False) and getattr(shape, 'text', ''):
                parts.append(shape.text)
        except Exception:
            pass
    return _norm(' '.join(parts))


def _shape_counts(slide: Any) -> dict[str, int]:
    counts = {'chart': 0, 'table': 0, 'image': 0, 'text_box': 0, 'shape': 0}
    for shape in slide.shapes:
        counts['shape'] += 1
        stype = getattr(shape, 'shape_type', None)
        if stype == MSO_SHAPE_TYPE.PICTURE:
            counts['image'] += 1
        elif getattr(shape, 'has_chart', False):
            counts['chart'] += 1
        elif getattr(shape, 'has_table', False):
            counts['table'] += 1
        elif getattr(shape, 'has_text_frame', False):
            counts['text_box'] += 1
    return counts


def classify_slide(slide: Any, slide_number: int, source_template_id: str, source_file: str, source_url: str) -> tuple[str, str, str, dict[str, Any]]:
    text = _collect_slide_text(slide)
    counts = _shape_counts(slide)
    years = len(re.findall(r'20\d{2}', text))
    pct = len(re.findall(r'\d+(?:\.\d+)?%', text))
    num = len(re.findall(r'\d+(?:\.\d+)?', text))
    slide_type, slide_subtype, layout_type = 'content', 'generic_content', 'standard_text'
    if slide_number == 1 or _contains(text, r'封面|cover|title|首页'):
        slide_type, slide_subtype, layout_type = 'cover', 'business_cover', 'hero_split'
    elif _contains(text, r'目录|agenda|contents'):
        slide_type, slide_subtype, layout_type = 'agenda', 'number_agenda', 'numbered_list'
    elif _contains(text, r'章节|chapter|section'):
        slide_type, slide_subtype, layout_type = 'section', 'chapter_break', 'big_title'
    elif counts['table']:
        slide_type, slide_subtype, layout_type = 'table', 'ranking_table', 'table_grid'
    elif counts['chart'] or pct >= 2 or num >= 8:
        slide_type, slide_subtype, layout_type = 'data_analysis', 'kpi_dashboard', 'dashboard_cards'
    elif _contains(text, r'时间轴|里程碑|timeline|roadmap') or years >= 3:
        slide_type, slide_subtype, layout_type = 'timeline', 'roadmap_timeline', 'horizontal_timeline'
    elif _contains(text, r'流程|process|workflow|步骤|审批'):
        slide_type, slide_subtype, layout_type = 'process', 'workflow_process', 'arrow_flow'
    elif _contains(text, r'对比|comparison|before|after|竞品'):
        slide_type, slide_subtype, layout_type = 'comparison', 'before_after', 'two_panel_compare'
    elif _contains(text, r'组织|团队|人员|专家|成员'):
        slide_type, slide_subtype, layout_type = 'team', 'team_profile', 'profile_cards'
    elif _contains(text, r'产品|功能|参数|亮点'):
        slide_type, slide_subtype, layout_type = 'product', 'product_highlight', 'feature_grid'
    elif _contains(text, r'地图|区域|分布|world|china'):
        slide_type, slide_subtype, layout_type = 'map', 'business_map', 'map_dashboard'
    elif _contains(text, r'战略|swot|pest|五力|价值链'):
        slide_type, slide_subtype, layout_type = 'strategy', 'swot_matrix', 'matrix_grid'
    elif counts['image'] > counts['text_box'] and counts['image'] >= 2:
        slide_type, slide_subtype, layout_type = 'image', 'gallery_layout', 'gallery_grid'
    return slide_type, slide_subtype, layout_type, {'slide_number': slide_number, 'source_template_id': source_template_id, 'source_file': source_file, 'source_url': source_url, 'text': text, 'counts': counts, 'years': years, 'percentages': pct, 'numbers': num}


def _component_records_from_slide(slide: Any, slide_id: str, layout_type: str, style: str, primary_color: str) -> list[ComponentRecord]:
    records: list[ComponentRecord] = []
    prs = slide.part.slide_layout.part.presentation
    width = int(getattr(prs, 'slide_width', 1) or 1)
    height = int(getattr(prs, 'slide_height', 1) or 1)
    for idx, shape in enumerate(slide.shapes, start=1):
        try:
            if not (getattr(shape, 'has_text_frame', False) or getattr(shape, 'has_chart', False) or getattr(shape, 'has_table', False) or getattr(shape, 'shape_type', None) == MSO_SHAPE_TYPE.PICTURE):
                continue
            if getattr(shape, 'has_chart', False):
                ctype, csub = 'chart', 'chart_panel'
            elif getattr(shape, 'has_table', False):
                ctype, csub = 'table', 'table_block'
            elif getattr(shape, 'shape_type', None) == MSO_SHAPE_TYPE.PICTURE:
                ctype, csub = 'image', 'image_panel'
            elif getattr(shape, 'has_text_frame', False):
                ctype, csub = ('title', 'title_block') if len(_norm(getattr(shape, 'text', ''))) <= 18 else ('text', 'text_block')
            else:
                ctype, csub = 'card', 'note_card'
            w = float(shape.width) / width if width else 0.0
            h = float(shape.height) / height if height else 0.0
            records.append(ComponentRecord(f'{slide_id}-auto-{idx:02d}', ctype, csub, slide_id, layout_type, f'0.00,0.00,{w:.2f},{h:.2f}', style, primary_color, w, h, '', {'auto_extracted': True, 'shape_index': idx}, 'ingested', (ctype, csub, style)))
        except Exception:
            pass
    return records


def extract_pptx_to_slide_catalog(pptx_path: str | Path, source_template_id: str | None = None, source_url: str | None = None, root: str | Path = 'ppt_template_library') -> dict[str, Any]:
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f'找不到 PPTX 文件：{pptx_path}')
    initialize_slide_storage(root)
    prs = Presentation(str(pptx_path))
    source_template_id = source_template_id or _slug(pptx_path.stem)
    source_url = source_url or f'file://{pptx_path.as_posix()}'
    paths = ensure_library_structure(root)
    pages: list[PageRecord] = []
    comps: list[ComponentRecord] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_type, slide_subtype, layout_type, meta = classify_slide(slide, idx, source_template_id, pptx_path.name, source_url)
        tags = _page_tags({'slide_type': slide_type, 'slide_subtype': slide_subtype, 'industry': _infer_industry(meta['text']), 'scenario': _infer_scenario(meta['text']), 'style': _infer_style(meta['text'], 'business'), 'layout_type': layout_type, 'tags': (slide_type, slide_subtype, layout_type)})
        preview = paths.preview / 'slides' / f'{source_template_id}-{idx:03d}.png'
        create_page_preview((meta['text'].split(' ')[0] if meta['text'] else f'Slide {idx}')[:36], meta['text'][:90] or 'Auto extracted slide', slide_type, slide_subtype, 84, tags, preview)
        page = PageRecord(f'{source_template_id}-{idx:03d}', source_template_id, pptx_path.name, source_url, idx, slide_type, slide_subtype, _infer_industry(meta['text']), _infer_scenario(meta['text']), _infer_style(meta['text'], 'business'), layout_type, '#1B3756', '#00A89D', '#F7F9FC', 'light', '16:9', min(1.0, len(meta['text']) / 1000.0), min(1.0, meta['counts']['image'] / 5.0), meta['counts']['chart'], meta['counts']['table'], meta['counts']['shape'], 0, meta['counts']['image'], meta['counts']['text_box'], meta['counts']['chart'] > 0, meta['counts']['table'] > 0, slide_type == 'timeline', slide_type == 'process', slide_type == 'map', slide_type == 'team', slide_type in {'data_analysis', 'chart', 'strategy', 'timeline', 'process', 'comparison'}, False, 'high' if meta['counts']['chart'] == 0 else 'medium', 84, 83, 82, 86, 84, 84, str(preview), str(preview), str(pptx_path), json.dumps({'keywords': list(tags), 'summary': meta['text'][:200]}, ensure_ascii=False), meta, 'ingested', tags)
        pages.append(page)
        comps.extend(_component_records_from_slide(slide, page.slide_id, layout_type, page.style, page.primary_color))
    save_slide_catalog(pages, root)
    if comps:
        save_component_catalog(comps, root)
    summary = load_slide_summary(root)
    summary['ingested'] = len(pages)
    summary['source_file'] = pptx_path.name
    return summary

def search_pages(filters: PageSearchFilters, root: str | Path = 'ppt_template_library') -> pd.DataFrame:
    initialize_slide_storage(root)
    df = load_slide_catalog(root)
    if df.empty:
        return df
    result = df.copy()
    query = _norm(filters.query)
    if query:
        terms = [t for t in re.split(r'[\s,，。;；/|]+', query) if t]
        mask = pd.Series(True, index=result.index)
        for term in terms:
            term_mask = pd.Series(False, index=result.index)
            for col in ['slide_type', 'slide_subtype', 'industry', 'scenario', 'style', 'layout_type', 'metadata_json', 'tags_json']:
                term_mask |= result[col].astype(str).str.contains(term, case=False, na=False)
            mask &= term_mask
        result = result[mask]
    for attr, col in [('layout_types', 'layout_type'), ('slide_types', 'slide_type'), ('industries', 'industry'), ('scenarios', 'scenario'), ('styles', 'style')]:
        values = getattr(filters, attr)
        if values:
            lowered = [v.lower() for v in values]
            result = result[result[col].astype(str).str.lower().isin(lowered)]
    if filters.tags:
        pattern = '|'.join(filters.tags)
        result = result[result['tags_json'].astype(str).str.contains(pattern, case=False, na=False)]
    if filters.colors:
        pattern = '|'.join(filters.colors)
        result = result[result['primary_color'].astype(str).str.contains(pattern, case=False, na=False) | result['secondary_color'].astype(str).str.contains(pattern, case=False, na=False)]
    if filters.min_score:
        result = result[result['overall_quality_score'] >= filters.min_score]
    if filters.min_quality:
        result = result[result['overall_quality_score'] >= filters.min_quality]
    for attr in ['has_chart', 'has_table', 'has_timeline', 'has_process', 'has_map', 'has_people', 'has_infographic']:
        value = getattr(filters, attr)
        if value is not None:
            result = result[result[attr].astype(bool) == value]
    if filters.dark_or_light:
        result = result[result['dark_or_light'].astype(str).str.lower().eq(filters.dark_or_light.lower())]
    return result.sort_values(['overall_quality_score', 'design_score', 'layout_score'], ascending=[False, False, False])


def search_components(query: str = '', component_types: tuple[str, ...] = (), root: str | Path = 'ppt_template_library') -> pd.DataFrame:
    initialize_slide_storage(root)
    df = load_component_catalog(root)
    if df.empty:
        return df
    result = df.copy()
    if query:
        terms = [t for t in re.split(r'[\s,，。;；/|]+', _norm(query)) if t]
        mask = pd.Series(True, index=result.index)
        for term in terms:
            term_mask = result['component_type'].astype(str).str.contains(term, case=False, na=False)
            term_mask |= result['component_subtype'].astype(str).str.contains(term, case=False, na=False)
            term_mask |= result['tags_json'].astype(str).str.contains(term, case=False, na=False)
            mask &= term_mask
        result = result[mask]
    if component_types:
        lowered = [v.lower() for v in component_types]
        result = result[result['component_type'].astype(str).str.lower().isin(lowered) | result['component_subtype'].astype(str).str.lower().isin(lowered)]
    return result.sort_values(['component_type', 'component_subtype', 'component_id'])


def infer_request(prompt: str, style_name: str, page_range: tuple[int, int], data_profile: ExcelProfile | None = None) -> AssemblyRequest:
    text = _norm(prompt)
    topics = tuple(dict.fromkeys(normalize_tokens(re.split(r'[\s,，。;；/|]+', text))))
    preferred_layouts: list[str] = []
    if _contains(text, r'数据|销售|财务|增长|指标'):
        preferred_layouts.extend(['dashboard_cards', 'line_chart', 'table_grid'])
    if _contains(text, r'项目|计划|里程碑|路线图'):
        preferred_layouts.extend(['horizontal_timeline', 'arrow_flow', 'gantt_chart'])
    if _contains(text, r'产品|发布|功能'):
        preferred_layouts.extend(['feature_grid', 'two_panel_compare'])
    if _contains(text, r'团队|组织|专家'):
        preferred_layouts.extend(['profile_cards', 'hub_spoke'])
    if _contains(text, r'地图|区域|分布'):
        preferred_layouts.extend(['map_dashboard'])
    style = _infer_style(text, style_name)
    return AssemblyRequest(prompt, _infer_title(prompt), f'{style_name} · {date.today().strftime("%Y-%m-%d")}', style, page_range, _infer_industry(text), _infer_scenario(text), 'office', 'formal', _infer_colors(text, style), topics, tuple(dict.fromkeys(preferred_layouts)), data_profile)


def _sequence(prompt: str, request: AssemblyRequest, page_count: int) -> list[dict[str, Any]]:
    text = _norm(prompt)
    has_data = _contains(text, r'数据|经营|销售|财务|增长|指标')
    has_project = _contains(text, r'项目|计划|进度|排期|里程碑')
    has_product = _contains(text, r'产品|发布|功能|参数')
    has_people = _contains(text, r'团队|组织|专家|人员')
    has_map = _contains(text, r'地图|区域|分布')
    slots = [
        {'page_type': 'cover', 'page_subtype': request.style, 'title': request.title},
        {'page_type': 'agenda', 'page_subtype': 'number_agenda', 'title': '目录'},
    ]
    if has_data:
        slots += [{'page_type': 'data_analysis', 'page_subtype': 'kpi_dashboard', 'title': '核心数据'}, {'page_type': 'chart', 'page_subtype': 'trend_chart', 'title': '趋势变化'}]
    elif has_project:
        slots += [{'page_type': 'timeline', 'page_subtype': 'roadmap_timeline', 'title': '项目进度'}, {'page_type': 'process', 'page_subtype': 'workflow_process', 'title': '执行流程'}]
    elif has_product:
        slots += [{'page_type': 'product', 'page_subtype': 'product_highlight', 'title': '产品亮点'}, {'page_type': 'comparison', 'page_subtype': 'before_after', 'title': '功能对比'}]
    elif has_people:
        slots += [{'page_type': 'team', 'page_subtype': 'team_profile', 'title': '团队介绍'}, {'page_type': 'relation', 'page_subtype': 'org_relationship', 'title': '组织关系'}]
    else:
        slots += [{'page_type': 'content', 'page_subtype': 'text_image_layout', 'title': '背景与目标'}, {'page_type': 'strategy', 'page_subtype': 'swot_matrix', 'title': '分析与判断'}]
    if has_map:
        slots.append({'page_type': 'map', 'page_subtype': 'business_map', 'title': '区域分布'})
    slots.append({'page_type': 'ending', 'page_subtype': 'thanks_end', 'title': '总结与行动'})
    while len(slots) < page_count:
        slots.insert(-1, [{'page_type': 'content', 'page_subtype': 'text_image_layout', 'title': '核心内容'}, {'page_type': 'chart', 'page_subtype': 'trend_chart', 'title': '趋势补充'}, {'page_type': 'comparison', 'page_subtype': 'before_after', 'title': '对比分析'}][len(slots) % 3])
    return slots[:page_count]


def build_page_level_outline(prompt: str, style_name: str, page_range: tuple[int, int], data_profile: ExcelProfile | None = None, root: str | Path = 'ppt_template_library') -> dict[str, Any]:
    seed_demo_page_library(root)
    request = infer_request(prompt, style_name, page_range, data_profile)
    min_pages, max_pages = page_range
    total_pages = max(min_pages, min(max_pages, max(len(request.topics) + 3, (min_pages + max_pages) // 2)))
    slots = _sequence(prompt, request, total_pages)
    sections = _infer_sections(prompt)
    pages = []
    for idx, slot in enumerate(slots[2:-1], start=1):
        candidates = search_pages(PageSearchFilters(query=' '.join([slot['page_type'], slot['page_subtype'], request.industry, request.scenario, request.style]), slide_types=(slot['page_type'],), industries=(request.industry,), scenarios=(request.scenario,), styles=(request.style,), colors=request.colors, tags=request.topics, min_quality=70), root).head(5)
        cands = []
        for _, row in candidates.iterrows():
            cands.append({'slide_id': row['slide_id'], 'slide_type': row['slide_type'], 'slide_subtype': row['slide_subtype'], 'layout_type': row['layout_type'], 'quality_score': int(row['overall_quality_score']), 'preview_image': row.get('preview_image', ''), 'match_score': float(row['overall_quality_score'])})
        pages.append({'section': slot['title'], 'title': slot['title'], 'bullets': [f"围绕 {slot['title']} 组织内容", f"优先使用 {slot['page_type']} 页面", '保持设计系统统一'], 'tag': f"{slot['page_type']} / {slot['page_subtype']}", 'page_type': slot['page_type'], 'page_subtype': slot['page_subtype'], 'layout_type': cands[0]['layout_type'] if cands else slot['page_subtype'], 'candidate_count': len(cands), 'candidates': cands, 'slot': idx})
    design_system = _design_system(request.style)
    outline = {'title': request.title, 'subtitle': request.subtitle, 'style': request.style, 'industry': request.industry, 'scenario': request.scenario, 'sections': sections, 'pages': pages, 'page_plan': slots, 'design_system': design_system.to_row(), 'request': request.to_row(), 'library_summary': load_slide_summary(root), 'total_pages': len(slots)}
    save_assembly_run({'request_json': outline['request'], 'plan_json': {'page_plan': slots, 'design_system': outline['design_system']}, 'output_path': '', 'slide_count': len(slots), 'layout_summary': {'sections': sections, 'page_types': [slot['page_type'] for slot in slots]}}, root)
    return outline


def build_library_overview(root: str | Path = 'ppt_template_library') -> dict[str, Any]:
    seed_demo_page_library(root)
    slide_df = load_slide_catalog(root)
    component_df = load_component_catalog(root)
    return {'slides': int(len(slide_df)), 'components': int(len(component_df)), 'page_types': int(slide_df['slide_type'].nunique()) if not slide_df.empty else 0, 'industries': int(slide_df['industry'].nunique()) if not slide_df.empty else 0, 'scenarios': int(slide_df['scenario'].nunique()) if not slide_df.empty else 0, 'average_score': round(float(slide_df['overall_quality_score'].mean()), 1) if not slide_df.empty else 0.0}


def build_page_preview_sheet(df: pd.DataFrame, out_path: str | Path, columns: int = 3, limit: int = 9) -> Path:
    image_paths = [row['preview_image'] for _, row in df.head(limit).iterrows() if str(row.get('preview_image', '')).strip()]
    if not image_paths:
        raise ValueError('No preview images available')
    return build_contact_sheet(image_paths, out_path, columns=columns)
