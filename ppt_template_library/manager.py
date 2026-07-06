from __future__ import annotations

import hashlib
import json
import re
import shutil
import subprocess
import tempfile
import zipfile
from datetime import datetime
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Iterable

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .catalog import build_source_catalog, filter_source_dataframe, save_source_catalog_json
from .download import discover_template_links, directory_hash, download_file, file_hash, is_supported_template_url
from .models import SearchFilters, SourceEntry, default_template_record
from .page_library import (
    build_component_summary,
    build_layout_summary,
    build_seed_component_catalog,
    build_seed_slide_catalog,
    generate_component_previews,
    generate_slide_previews,
)
from .page_models import ComponentRecord, PageRecord, PageSearchFilters
from .paths import ensure_library_structure, library_paths
from .preview import build_contact_sheet, create_page_preview, create_source_preview
from .scoring import score_template
from .slide_storage import (
    initialize_slide_storage,
    load_component_catalog,
    load_slide_catalog,
    load_slide_summary,
    save_component_catalog,
    save_slide_catalog,
)
from .storage import (
    initialize_storage,
    load_catalog_summary,
    load_source_catalog,
    load_template_catalog,
    save_source_catalog,
    save_template_catalog,
)


INDUSTRY_KEYWORDS: dict[str, tuple[str, ...]] = {
    "finance": ("finance", "financial", "bank", "banking", "investment", "investor", "insurance", "securities", "stock", "capital"),
    "technology": ("technology", "tech", "software", "saas", "ai", "cloud", "data", "internet", "digital", "cybersecurity"),
    "education": ("education", "school", "university", "academic", "lecture", "course", "training", "thesis", "research"),
    "government": ("government", "policy", "public", "official", "政务", "党建", "government"),
    "manufacturing": ("manufacturing", "industrial", "factory", "production", "supplier", "supply chain"),
    "healthcare": ("healthcare", "medical", "pharma", "pharmaceutical", "biotech", "clinical", "hospital"),
    "real_estate": ("real estate", "property", "construction", "architecture", "building"),
    "logistics": ("logistics", "delivery", "shipping", "supply chain", "warehouse", "transport"),
    "energy": ("energy", "power", "electric", "renewable", "solar", "wind", "photovoltaic", "utility"),
    "automotive": ("automotive", "car", "vehicle", "mobility", "EV", "new energy", "新能源"),
    "retail": ("retail", "ecommerce", "e-commerce", "shop", "consumer", "sales"),
    "media": ("media", "film", "tv", "music", "content", "broadcast"),
    "consulting": ("consulting", "strategy", "advisor", "analysis", "business model"),
    "marketing": ("marketing", "brand", "campaign", "advertising", "advertisement"),
    "startup": ("startup", "founder", "pitch deck", "venture", "fundraising"),
    "nonprofit": ("nonprofit", "public welfare", "公益", "charity", "foundation"),
    "sports": ("sports", "athletic", "fitness", "game", "team"),
    "travel": ("travel", "tourism", "hotel", "hospitality", "airline", "flight"),
    "art": ("art", "creative", "design", "illustration", "photography", "music"),
    "law": ("law", "legal", "attorney", "court", "compliance"),
    "hr": ("hr", "human resources", "talent", "recruiting", "people ops"),
}

SCENARIO_KEYWORDS: dict[str, tuple[str, ...]] = {
    "annual_report": ("annual report", "yearly report", "year end", "年度总结", "年报"),
    "quarterly_report": ("quarterly report", "q1", "q2", "q3", "q4", "季度总结"),
    "monthly_report": ("monthly report", "月报"),
    "weekly_report": ("weekly report", "周报"),
    "daily_report": ("daily report", "日报"),
    "project_report": ("project report", "项目汇报", "项目复盘"),
    "project_plan": ("project plan", "项目计划", "roadmap"),
    "business_plan": ("business plan", "商业计划书", "融资计划书", "创业计划书"),
    "pitch_deck": ("pitch deck", "融资", "投资人", "investor pitch"),
    "product_launch": ("product launch", "产品发布", "新品发布", "launch"),
    "product_intro": ("product intro", "产品介绍", "product overview"),
    "company_profile": ("company profile", "公司介绍", "企业宣传"),
    "market_analysis": ("market analysis", "市场分析", "行业分析", "竞争分析"),
    "data_analysis": ("data analysis", "数据分析", "dashboard", "kpi"),
    "financial_report": ("financial report", "财务报告", "budget", "profit", "loss"),
    "sales_report": ("sales report", "销售报告", "业绩"),
    "operations_report": ("operations report", "运营报告", "ops"),
    "research_report": ("research report", "research", "study", "科研汇报", "学术报告"),
    "training": ("training", "course", "培训课件", "教学课件"),
    "defense": ("defense", "毕业答辩", "论文答辩"),
    "government_briefing": ("government briefing", "政府汇报", "党建汇报"),
    "event_plan": ("event plan", "活动策划", "节日", "婚礼", "生日"),
    "resume": ("resume", "cv", "个人简历"),
    "portfolio": ("portfolio", "作品集", "摄影展示", "旅行记录"),
}

STYLE_KEYWORDS: dict[str, tuple[str, ...]] = {
    "business": ("business", "corporate", "commercial", "executive"),
    "minimal": ("minimal", "simple", "clean", "简约", "极简"),
    "modern": ("modern", "contemporary", "today", "current"),
    "technology": ("technology", "tech", "ai", "digital", "internet", "future", "科技", "未来"),
    "academic": ("academic", "research", "scholarly", "学术"),
    "government": ("government", "official", "formal", "政务", "党建"),
    "creative": ("creative", "editorial", "art", "magazine", "creative"),
    "industrial": ("industrial", "mechanical", "machinery", "manufacturing"),
    "finance": ("finance", "financial", "bank", "blue", "green"),
    "dark": ("dark", "black", "night", "黑金", "深色"),
    "colorful": ("colorful", "memphis", "playful", "children", "卡通"),
    "chinese": ("chinese", "国风", "国潮", "东方"),
    "japanese": ("japanese", "日式", "wabi", "minimal"),
    "korean": ("korean", "韩式", "soft", "clean"),
    "luxury": ("luxury", "premium", "elegant", "high-end", "高端", "奢华"),
}

COLOR_KEYWORDS: dict[str, tuple[str, ...]] = {
    "blue": ("blue", "navy", "azure", "blue商务", "蓝", "蓝色"),
    "green": ("green", "emerald", "mint", "环保", "绿色"),
    "red": ("red", "crimson", "burgundy", "政务", "红色"),
    "purple": ("purple", "violet", "lavender", "紫"),
    "orange": ("orange", "amber", "tangerine", "橙"),
    "black_gold": ("black gold", "black-gold", "黑金", "gold"),
    "gray": ("gray", "grey", "slate", "neutral"),
    "white": ("white", "ivory", "cream"),
}

CONTENT_KEYWORDS: dict[str, tuple[str, ...]] = {
    "contains_chart": ("chart", "graph", "diagram", "图表", "analytics"),
    "contains_table": ("table", "grid", "表格", "matrix"),
    "contains_timeline": ("timeline", "roadmap", "milestone", "时间轴"),
    "contains_process": ("process", "workflow", "flow", "流程"),
    "contains_map": ("map", "地图", "world", "china"),
    "contains_infographic": ("infographic", "information", "diagram", "data viz"),
    "contains_images": ("image", "photo", "gallery", "摄影", "图片"),
    "contains_icons": ("icon", "icons", "symbol", "pictogram"),
    "contains_animations": ("animation", "animated", "motion", "transition"),
}


class _LinkCollector(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.links: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag.lower() != "a":
            return
        href = dict(attrs).get("href")
        if href:
            self.links.append(href)


def _normalize_text(value: str | Iterable[str]) -> str:
    if isinstance(value, str):
        raw = value
    else:
        raw = " ".join(str(item) for item in value)
    return re.sub(r"\s+", " ", raw).strip().lower()


def _match_keywords(text: str, mapping: dict[str, tuple[str, ...]]) -> list[str]:
    hits: list[str] = []
    for label, keywords in mapping.items():
        if any(keyword.lower() in text for keyword in keywords):
            hits.append(label)
    return hits


def _first_or_empty(values: list[str], fallback: str = "") -> str:
    return values[0] if values else fallback


def _as_bool_flag(text: str, keywords: tuple[str, ...]) -> bool:
    return any(keyword.lower() in text for keyword in keywords)


def infer_classification(text: str, tags: Iterable[str] = ()) -> dict[str, Any]:
    merged = _normalize_text([text, *tags])
    industry_hits = _match_keywords(merged, INDUSTRY_KEYWORDS)
    scenario_hits = _match_keywords(merged, SCENARIO_KEYWORDS)
    style_hits = _match_keywords(merged, STYLE_KEYWORDS)
    color_hits = _match_keywords(merged, COLOR_KEYWORDS)

    result = {
        "industries": industry_hits,
        "scenarios": scenario_hits,
        "styles": style_hits,
        "colors": color_hits,
    }
    return result


def infer_content_flags(text: str, tags: Iterable[str] = ()) -> dict[str, bool]:
    merged = _normalize_text([text, *tags])
    return {key: _as_bool_flag(merged, keywords) for key, keywords in CONTENT_KEYWORDS.items()}


def infer_preview_layouts(text: str, tags: Iterable[str] = ()) -> list[str]:
    merged = _normalize_text([text, *tags])
    layout_candidates = [
        ("dashboard_slide", ("dashboard", "kpi", "data", "finance", "sales", "operations", "report")),
        ("timeline_slide", ("timeline", "roadmap", "history", "milestone")),
        ("process_slide", ("process", "workflow", "flow", "step")),
        ("map_slide", ("map", "location", "logistics", "distribution")),
        ("people_slide", ("team", "people", "profile", "staff", "speaker")),
        ("product_slide", ("product", "feature", "launch", "spec")),
        ("strategy_slide", ("strategy", "swot", "matrix", "pest", "consulting")),
        ("comparison_slide", ("comparison", "contrast", "before after", "pros cons")),
        ("chart_slide", ("chart", "graph", "analytics", "insight")),
        ("table_slide", ("table", "grid", "rank", "quotation")),
    ]
    selected = "content_slide"
    for layout_type, keywords in layout_candidates:
        if any(keyword in merged for keyword in keywords):
            selected = layout_type
            break
    layouts = ["cover_slide", "agenda_slide", selected]
    if selected in {"dashboard_slide", "chart_slide"}:
        layouts.insert(3, "chart_slide")
    if selected in {"table_slide", "comparison_slide"}:
        layouts.insert(3, "table_slide")
    layouts.append("ending_slide")
    ordered: list[str] = []
    for layout in layouts:
        if layout not in ordered:
            ordered.append(layout)
    return ordered


def _license_bucket(license_hint: str, commercial_use_hint: str) -> str:
    license_lower = license_hint.lower()
    if commercial_use_hint == "allowed":
        return "commercial_allowed"
    if commercial_use_hint == "restricted":
        return "personal_use_only"
    if license_lower in {"mit", "apache-2.0", "cc0-1.0", "public-domain", "cc-by", "cc-by-sa", "lgpl-2.1", "gpl-3.0"}:
        return "commercial_allowed"
    return "license_uncertain"


def _color_pair(color_hits: list[str]) -> tuple[str, str, str]:
    primary = _first_or_empty(color_hits, "blue")
    secondary_map = {
        "blue": "gray",
        "green": "blue",
        "red": "gold",
        "purple": "blue",
        "orange": "blue",
        "black_gold": "gold",
        "gray": "blue",
        "white": "blue",
    }
    secondary = secondary_map.get(primary, "blue")
    dark_or_light = "dark" if primary in {"black_gold"} else "light"
    return primary, secondary, dark_or_light


def _build_preview_bundle(
    template_id: str,
    template_name: str,
    subtitle: str,
    score: int,
    tags: Iterable[str],
    preview_dir: Path,
    layout_types: list[str],
) -> str:
    preview_dir.mkdir(parents=True, exist_ok=True)
    page_dir = preview_dir / "pages" / template_id
    page_dir.mkdir(parents=True, exist_ok=True)
    page_paths: list[Path] = []
    for index, layout_type in enumerate(layout_types, start=1):
        page_path = page_dir / f"{index:02d}-{layout_type}.png"
        create_page_preview(
            title=template_name if index == 1 else f"{template_name} · Page {index}",
            subtitle=subtitle,
            layout_type=layout_type,
            score=max(60, min(100, score)),
            tags=tags,
            out_path=page_path,
        )
        page_paths.append(page_path)
    contact_sheet = preview_dir / f"{template_id}.png"
    build_contact_sheet(page_paths, contact_sheet, columns=2)
    return str(contact_sheet)


def build_template_scaffold_catalog(entries: Iterable[SourceEntry] | None = None, root: str | Path = "ppt_template_library") -> list[dict[str, Any]]:
    paths = ensure_library_structure(root)
    source_entries = list(entries) if entries is not None else build_source_catalog()
    records: list[dict[str, Any]] = []
    for index, entry in enumerate(source_entries, start=1):
        text = " ".join([entry.name, entry.notes, " ".join(entry.tags), entry.source_website, entry.source_type])
        classification = infer_classification(text, entry.tags)
        content_flags = infer_content_flags(text, entry.tags)
        primary_color, secondary_color, dark_or_light = _color_pair(classification["colors"])
        quality_score = int(entry.score)
        design_score = max(55, min(100, quality_score - 4))
        usability_score = max(55, min(100, quality_score - 2))
        richness_bonus = sum(1 for value in content_flags.values() if value)
        overall_score = score_template(design_score, usability_score, quality_score, richness_bonus)
        template_id = f"scaffold-{index:04d}"
        layout_types = infer_preview_layouts(text, entry.tags)
        preview_image = _build_preview_bundle(
            template_id=template_id,
            template_name=entry.name,
            subtitle=f"{entry.source_website} · {entry.language} · {entry.source_type}",
            score=overall_score,
            tags=entry.tags,
            preview_dir=paths.preview / "templates",
            layout_types=layout_types,
        )
        record = default_template_record(
            template_id=template_id,
            template_name=entry.name,
            source_url=entry.source_url,
            source_website=entry.source_website,
        )
        record.update(
            {
                "author": entry.source_website,
                "license": entry.license_hint,
                "commercial_use": entry.commercial_use_hint,
                "modification_allowed": entry.modification_allowed_hint,
                "download_date": datetime.utcnow().date().isoformat(),
                "file_format": "scaffold",
                "file_size": None,
                "slide_count": len(layout_types),
                "aspect_ratio": "16:9",
                "language": entry.language,
                "industry": _first_or_empty(classification["industries"], _first_or_empty(list(entry.tags), "")),
                "scenario": _first_or_empty(classification["scenarios"], "presentation_template"),
                "style": _first_or_empty(classification["styles"], "modern"),
                "primary_color": primary_color,
                "secondary_color": secondary_color,
                "dark_or_light": dark_or_light,
                "contains_chart": content_flags["contains_chart"],
                "contains_table": content_flags["contains_table"],
                "contains_timeline": content_flags["contains_timeline"],
                "contains_process": content_flags["contains_process"],
                "contains_map": content_flags["contains_map"],
                "contains_infographic": content_flags["contains_infographic"],
                "contains_images": content_flags["contains_images"],
                "contains_icons": content_flags["contains_icons"],
                "contains_animations": content_flags["contains_animations"],
                "quality_score": quality_score,
                "design_score": design_score,
                "usability_score": usability_score,
                "duplicate_hash": entry.source_id,
                "preview_image": preview_image,
                "local_path": "",
                "metadata_json": {
                    "scaffold": True,
                    "source_type": entry.source_type,
                    "download_method": entry.download_method,
                    "programmatic_access": entry.programmatic_access,
                    "tags": list(entry.tags),
                    "notes": entry.notes,
                    "classification": classification,
                    "content_flags": content_flags,
                    "layout_types": layout_types,
                    "license_bucket": _license_bucket(entry.license_hint, entry.commercial_use_hint),
                    "overall_score": overall_score,
                },
                "status": "scaffold",
            }
        )
        records.append(record)
    return records


def write_index_json(path: Path, records: Iterable[dict[str, Any]]) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(list(records), ensure_ascii=False, indent=2), encoding="utf-8")
    return path


def export_bucket_indexes(root: str | Path = "ppt_template_library") -> dict[str, Path]:
    paths = ensure_library_structure(root)
    template_df = load_template_catalog(root)
    slide_df = load_slide_catalog(root)
    component_df = load_component_catalog(root)

    outputs: dict[str, Path] = {}
    if not template_df.empty:
        for bucket_name, subset in {
            "commercial_allowed": template_df[template_df["commercial_use"].astype(str).str.lower().eq("allowed")],
            "personal_use_only": template_df[template_df["commercial_use"].astype(str).str.lower().eq("restricted")],
            "license_uncertain": template_df[template_df["commercial_use"].astype(str).str.lower().isin(["uncertain", "n/a", ""])],
            "premium_quality": template_df[template_df["quality_score"] >= 90],
        }.items():
            payload = subset[["template_id", "template_name", "quality_score", "license", "commercial_use", "preview_image", "local_path", "status"]].to_dict(orient="records")
            outputs[bucket_name] = write_index_json(getattr(paths, bucket_name) / "index.json", payload)

        for field_name, base_dir in [
            ("industry", paths.industry),
            ("scenario", paths.scenario),
            ("style", paths.style),
            ("primary_color", paths.color),
        ]:
            if field_name not in template_df.columns:
                continue
            for value, subset in template_df.groupby(template_df[field_name].fillna("unknown").astype(str)):
                payload = subset[["template_id", "template_name", "quality_score", "preview_image", "status"]].to_dict(orient="records")
                outputs[f"{field_name}:{value}"] = write_index_json(base_dir / f"{value or 'unknown'}.json", payload)

    if not slide_df.empty:
        for value, subset in slide_df.groupby(slide_df["layout_type"].fillna("unknown").astype(str)):
            payload = subset[["slide_id", "slide_type", "slide_subtype", "overall_quality_score", "preview_image", "status"]].to_dict(orient="records")
            outputs[f"layout:{value}"] = write_index_json(paths.layout / f"{value or 'unknown'}.json", payload)

    if not component_df.empty:
        for value, subset in component_df.groupby(component_df["component_type"].fillna("unknown").astype(str)):
            payload = subset[["component_id", "component_type", "component_subtype", "layout_type", "preview_image", "status"]].to_dict(orient="records")
            outputs[f"component:{value}"] = write_index_json(paths.components / f"{value or 'unknown'}.json", payload)

    return outputs


def bootstrap_source_library(root: str | Path = "ppt_template_library") -> dict[str, Any]:
    paths = ensure_library_structure(root)
    initialize_storage(root)
    entries = build_source_catalog()
    save_source_catalog(entries, root)
    save_source_catalog_json(paths.metadata / "source_catalog.json")
    write_index_json(paths.source_records / "index.json", [entry.to_row() for entry in entries])
    source_df = load_source_catalog(root)
    return {
        "root": str(paths.root),
        "source_count": int(len(entries)),
        "commercial_allowed": int((source_df["commercial_use_hint"] == "allowed").sum()) if not source_df.empty else 0,
        "uncertain": int((source_df["commercial_use_hint"] == "uncertain").sum()) if not source_df.empty else 0,
        "average_score": round(float(source_df["score"].mean()), 1) if not source_df.empty else 0.0,
    }


def bootstrap_template_scaffolds(root: str | Path = "ppt_template_library", generate_previews: bool = True) -> dict[str, Any]:
    paths = ensure_library_structure(root)
    initialize_storage(root)
    entries = build_source_catalog()
    scaffolds = build_template_scaffold_catalog(entries, root=root)
    save_template_catalog(scaffolds, root)
    if generate_previews:
        # _build_preview_bundle already writes the preview assets.
        pass
    write_index_json(paths.metadata / "template_scaffold_catalog.json", scaffolds)
    return {
        "template_count": len(scaffolds),
        "premium_count": sum(1 for record in scaffolds if int(record.get("quality_score", 0)) >= 90),
    }


def bootstrap_page_components(root: str | Path = "ppt_template_library", generate_previews: bool = True) -> dict[str, Any]:
    initialize_slide_storage(root)
    slide_records = build_seed_slide_catalog()
    component_records = build_seed_component_catalog(slide_records)
    save_slide_catalog(slide_records, root)
    save_component_catalog(component_records, root)
    if generate_previews:
        generate_slide_previews(slide_records)
        generate_component_previews(component_records)
    paths = ensure_library_structure(root)
    write_index_json(paths.metadata / "slide_catalog.json", [record.to_row() for record in slide_records])
    write_index_json(paths.metadata / "component_catalog.json", [record.to_row() for record in component_records])
    return {
        "slide_count": len(slide_records),
        "component_count": len(component_records),
        "premium_count": sum(1 for record in slide_records if record.overall_quality_score >= 90),
    }


def bootstrap_library(root: str | Path = "ppt_template_library", generate_previews: bool = True) -> dict[str, Any]:
    bootstrap_source_library(root)
    bootstrap_template_scaffolds(root, generate_previews=generate_previews)
    bootstrap_page_components(root, generate_previews=generate_previews)
    export_bucket_indexes(root)
    manifest = build_library_manifest(root)
    payload = {
        **manifest,
        "source_count": manifest.get("sources", 0),
        "template_count": manifest.get("template_scaffolds", 0),
        "slide_count": manifest.get("slides", 0),
        "component_count": manifest.get("components", 0),
    }
    write_index_json(ensure_library_structure(root).metadata / "library_manifest.json", [payload])
    return payload


def _stringify_metadata(value: Any) -> str:
    if isinstance(value, (dict, list)):
        return json.dumps(value, ensure_ascii=False)
    return str(value or "")


def _frame_text(df: pd.DataFrame, columns: Iterable[str]) -> pd.Series:
    parts = []
    for column in columns:
        if column in df.columns:
            parts.append(df[column].astype(str))
    if not parts:
        return pd.Series("", index=df.index)
    result = parts[0]
    for part in parts[1:]:
        result = result + " " + part
    return result.str.lower()


def filter_template_dataframe(
    df: pd.DataFrame,
    query: str = "",
    industries: tuple[str, ...] = (),
    scenarios: tuple[str, ...] = (),
    styles: tuple[str, ...] = (),
    colors: tuple[str, ...] = (),
    licenses: tuple[str, ...] = (),
    statuses: tuple[str, ...] = (),
    min_score: int = 0,
    commercial_only: bool = False,
    scaffold_only: bool = False,
) -> pd.DataFrame:
    result = df.copy()
    if result.empty:
        return result

    result["metadata_text"] = result.get("metadata_json", "{}").astype(str).str.lower()
    if query:
        terms = [term for term in _normalize_text(query).split() if term]
        if terms:
            searchable = _frame_text(
                result,
                [
                    "template_name",
                    "source_url",
                    "source_website",
                    "author",
                    "license",
                    "commercial_use",
                    "modification_allowed",
                    "industry",
                    "scenario",
                    "style",
                    "primary_color",
                    "secondary_color",
                    "status",
                    "metadata_text",
                ],
            )
            mask = pd.Series(True, index=result.index)
            for term in terms:
                mask &= searchable.str.contains(re.escape(term), na=False)
            result = result[mask]

    def _contains_any(column: str, values: tuple[str, ...]) -> pd.Series:
        if not values or column not in result.columns:
            return pd.Series(True, index=result.index)
        needle = "|".join(re.escape(value.lower()) for value in values)
        return result[column].astype(str).str.lower().str.contains(needle, na=False)

    if industries:
        result = result[_contains_any("industry", industries)]
    if scenarios:
        result = result[_contains_any("scenario", scenarios)]
    if styles:
        result = result[_contains_any("style", styles)]
    if colors:
        color_mask = _contains_any("primary_color", colors)
        if "secondary_color" in result.columns:
            color_mask |= _contains_any("secondary_color", colors)
        result = result[color_mask]
    if licenses:
        result = result[_contains_any("license", licenses)]
    if statuses:
        result = result[_contains_any("status", statuses)]
    if min_score and "quality_score" in result.columns:
        result = result[result["quality_score"].fillna(0).astype(int) >= int(min_score)]
    if commercial_only and "commercial_use" in result.columns:
        result = result[result["commercial_use"].astype(str).str.lower().eq("allowed")]
    if scaffold_only and "status" in result.columns:
        result = result[result["status"].astype(str).str.lower().eq("scaffold")]

    sort_columns = [column for column in ["quality_score", "design_score", "usability_score", "template_name"] if column in result.columns]
    ascending = [False, False, False, True][: len(sort_columns)]
    if sort_columns:
        result = result.sort_values(sort_columns, ascending=ascending)
    return result.drop(columns=["metadata_text"], errors="ignore")


def filter_slide_dataframe(df: pd.DataFrame, filters: PageSearchFilters) -> pd.DataFrame:
    result = df.copy()
    if result.empty:
        return result

    if filters.query:
        terms = [part for part in _normalize_text(filters.query).split() if part]
        searchable = _frame_text(
            result,
            [
                "slide_id",
                "source_template_id",
                "slide_type",
                "slide_subtype",
                "industry",
                "scenario",
                "style",
                "layout_type",
                "metadata_json",
            ],
        )
        mask = pd.Series(True, index=result.index)
        for term in terms:
            mask &= searchable.str.contains(re.escape(term), na=False)
        result = result[mask]

    for column_name, values in [
        ("layout_type", filters.layout_types),
        ("slide_type", filters.slide_types),
        ("industry", filters.industries),
        ("scenario", filters.scenarios),
        ("style", filters.styles),
        ("primary_color", filters.colors),
        ("tags_json", filters.tags),
    ]:
        if values and column_name in result.columns:
            needle = "|".join(re.escape(value.lower()) for value in values)
            result = result[result[column_name].astype(str).str.lower().str.contains(needle, na=False)]

    if filters.min_score and "overall_quality_score" in result.columns:
        result = result[result["overall_quality_score"].fillna(0).astype(int) >= int(filters.min_score)]
    if filters.has_chart is not None and "has_chart" in result.columns:
        result = result[result["has_chart"].astype(bool).eq(filters.has_chart)]
    if filters.has_table is not None and "has_table" in result.columns:
        result = result[result["has_table"].astype(bool).eq(filters.has_table)]
    if filters.has_timeline is not None and "has_timeline" in result.columns:
        result = result[result["has_timeline"].astype(bool).eq(filters.has_timeline)]
    if filters.has_process is not None and "has_process" in result.columns:
        result = result[result["has_process"].astype(bool).eq(filters.has_process)]
    if filters.has_map is not None and "has_map" in result.columns:
        result = result[result["has_map"].astype(bool).eq(filters.has_map)]
    if filters.has_people is not None and "has_people" in result.columns:
        result = result[result["has_people"].astype(bool).eq(filters.has_people)]
    if filters.has_infographic is not None and "has_infographic" in result.columns:
        result = result[result["has_infographic"].astype(bool).eq(filters.has_infographic)]
    if filters.dark_or_light and "dark_or_light" in result.columns:
        result = result[result["dark_or_light"].astype(str).str.lower().eq(filters.dark_or_light.lower())]
    if filters.min_quality and "overall_quality_score" in result.columns:
        result = result[result["overall_quality_score"].fillna(0).astype(int) >= int(filters.min_quality)]

    sort_columns = [column for column in ["overall_quality_score", "design_score", "layout_score", "slide_id"] if column in result.columns]
    ascending = [False, False, False, True][: len(sort_columns)]
    if sort_columns:
        result = result.sort_values(sort_columns, ascending=ascending)
    return result


def filter_component_dataframe(
    df: pd.DataFrame,
    query: str = "",
    component_types: tuple[str, ...] = (),
    layout_types: tuple[str, ...] = (),
    styles: tuple[str, ...] = (),
    min_score: int = 0,
) -> pd.DataFrame:
    result = df.copy()
    if result.empty:
        return result
    if query:
        terms = [term for term in _normalize_text(query).split() if term]
        searchable = _frame_text(
            result,
            ["component_id", "component_type", "component_subtype", "layout_type", "style_token", "color_token", "metadata_json"],
        )
        mask = pd.Series(True, index=result.index)
        for term in terms:
            mask &= searchable.str.contains(re.escape(term), na=False)
        result = result[mask]
    if component_types and "component_type" in result.columns:
        needle = "|".join(re.escape(value.lower()) for value in component_types)
        result = result[result["component_type"].astype(str).str.lower().str.contains(needle, na=False)]
    if layout_types and "layout_type" in result.columns:
        needle = "|".join(re.escape(value.lower()) for value in layout_types)
        result = result[result["layout_type"].astype(str).str.lower().str.contains(needle, na=False)]
    if styles and "style_token" in result.columns:
        needle = "|".join(re.escape(value.lower()) for value in styles)
        result = result[result["style_token"].astype(str).str.lower().str.contains(needle, na=False)]
    if min_score and "overall_quality_score" in result.columns:
        result = result[result["overall_quality_score"].fillna(0).astype(int) >= int(min_score)]
    sort_columns = [column for column in ["overall_quality_score", "component_type", "component_id"] if column in result.columns]
    ascending = [False, True, True][: len(sort_columns)]
    if sort_columns:
        result = result.sort_values(sort_columns, ascending=ascending)
    return result


def build_library_manifest(root: str | Path = "ppt_template_library") -> dict[str, Any]:
    source_summary = load_catalog_summary(root)
    slide_summary = load_slide_summary(root)
    template_df = load_template_catalog(root)
    component_df = load_component_catalog(root)
    return {
        "sources": source_summary.get("sources", 0),
        "templates": source_summary.get("templates", 0),
        "template_scaffolds": int(len(template_df)),
        "slides": slide_summary.get("slide_count", 0),
        "components": slide_summary.get("component_count", 0),
        "premium_templates": int((template_df["quality_score"] >= 90).sum()) if not template_df.empty else 0,
        "premium_slides": slide_summary.get("premium_count", 0),
        "average_slide_score": slide_summary.get("average_score", 0.0),
        "component_types": int(component_df["component_type"].nunique()) if not component_df.empty else 0,
    }


def refresh_library(root: str | Path = "ppt_template_library", generate_previews: bool = True) -> dict[str, Any]:
    manifest = bootstrap_library(root=root, generate_previews=generate_previews)
    return {**manifest, **build_library_manifest(root)}


def validate_template_asset(path: str | Path) -> dict[str, Any]:
    file_path = Path(path)
    result: dict[str, Any] = {
        "path": str(file_path),
        "exists": file_path.exists(),
        "is_directory": file_path.is_dir(),
        "size": file_path.stat().st_size if file_path.exists() and file_path.is_file() else 0,
        "file_format": file_path.suffix.lower().lstrip("."),
        "valid": False,
        "reason": "",
    }
    if not file_path.exists():
        result["reason"] = "not_found"
        return result
    if file_path.is_dir():
        candidates = [candidate for candidate in file_path.rglob("*") if candidate.is_file()]
        result["valid"] = any(candidate.suffix.lower() in {".pptx", ".potx", ".ppt", ".odp", ".md", ".qmd", ".html", ".htm"} for candidate in candidates)
        result["reason"] = "directory_scan"
        return result
    if result["size"] <= 0:
        result["reason"] = "empty_file"
        return result
    suffix = file_path.suffix.lower()
    try:
        if suffix in {".pptx", ".potx"}:
            if not zipfile.is_zipfile(file_path):
                raise ValueError("invalid_zip")
            Presentation(str(file_path))
        elif suffix == ".odp":
            if not zipfile.is_zipfile(file_path):
                raise ValueError("invalid_odp_zip")
        elif suffix == ".zip":
            if not zipfile.is_zipfile(file_path):
                raise ValueError("invalid_zip")
        elif suffix in {".ppt", ".md", ".qmd", ".html", ".htm", ".txt"}:
            pass
        else:
            pass
        result["valid"] = True
        result["reason"] = "ok"
    except Exception as exc:  # pragma: no cover - defensive
        result["reason"] = str(exc)
    return result


def _aspect_ratio(width: int, height: int) -> str:
    if not width or not height:
        return ""
    ratio = width / height
    if abs(ratio - (16 / 9)) < 0.08:
        return "16:9"
    if abs(ratio - (4 / 3)) < 0.08:
        return "4:3"
    return f"{ratio:.2f}:1"


def _has_picture(shape: Any) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return hasattr(shape, "image")


def extract_pptx_metadata(path: str | Path) -> dict[str, Any]:
    file_path = Path(path)
    prs = Presentation(str(file_path))
    slide_count = len(prs.slides)
    charts = tables = pictures = text_boxes = shapes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            shapes += 1
            if getattr(shape, "has_chart", False):
                charts += 1
            if getattr(shape, "has_table", False):
                tables += 1
            if _has_picture(shape):
                pictures += 1
            if getattr(shape, "has_text_frame", False):
                text_boxes += 1
    width = int(prs.slide_width)
    height = int(prs.slide_height)
    return {
        "slide_count": slide_count,
        "aspect_ratio": _aspect_ratio(width, height),
        "width_emu": width,
        "height_emu": height,
        "contains_chart": charts > 0,
        "contains_table": tables > 0,
        "contains_images": pictures > 0,
        "text_box_count": text_boxes,
        "shape_count": shapes,
        "contains_timeline": False,
        "contains_process": False,
        "contains_map": False,
        "contains_icons": False,
        "contains_animations": False,
    }


def _directory_fingerprint(path: Path) -> str:
    sha = hashlib.sha256()
    for file_path in sorted(candidate for candidate in path.rglob("*") if candidate.is_file() and ".git" not in candidate.parts):
        sha.update(str(file_path.relative_to(path)).encode("utf-8"))
        sha.update(str(file_path.stat().st_size).encode("utf-8"))
    return sha.hexdigest()


def file_signature(path: str | Path) -> dict[str, str]:
    file_path = Path(path)
    if file_path.is_dir():
        return {
            "file_hash": _directory_fingerprint(file_path),
            "image_hash": "",
            "structure_hash": _directory_fingerprint(file_path),
        }
    file_hash_value = file_hash(file_path)
    image_hash_value = ""
    if file_path.suffix.lower() in {".png", ".jpg", ".jpeg", ".webp", ".bmp"}:
        try:
            image = Image.open(file_path).convert("L").resize((8, 8))
            pixels = list(image.getdata())
            average = sum(pixels) / len(pixels)
            bits = "".join("1" if pixel >= average else "0" for pixel in pixels)
            image_hash_value = f"{int(bits, 2):016x}"
        except Exception:
            image_hash_value = ""
    return {
        "file_hash": file_hash_value,
        "image_hash": image_hash_value,
        "structure_hash": file_hash_value,
    }


def _clone_git_repository(source_url: str, target_dir: Path) -> Path:
    git_executable = shutil.which("git") or shutil.which("git.exe")
    if not git_executable:
        raise RuntimeError("git_not_available")
    target_dir.parent.mkdir(parents=True, exist_ok=True)
    if target_dir.exists():
        shutil.rmtree(target_dir)
    subprocess.run([git_executable, "clone", "--depth", "1", source_url, str(target_dir)], check=True, capture_output=True, text=True)
    return target_dir


def download_template_source(source: SourceEntry, root: str | Path = "ppt_template_library", timeout: int = 45) -> dict[str, Any]:
    paths = ensure_library_structure(root)
    template_root = paths.templates / source.source_id
    result: dict[str, Any] = {
        "source_id": source.source_id,
        "source_url": source.source_url,
        "download_method": source.download_method,
        "target_path": str(template_root),
        "valid": False,
        "reason": "",
    }
    try:
        if source.download_method == "git" or source.source_url.endswith(".git") or "github.com" in source.source_url.lower():
            downloaded = _clone_git_repository(source.source_url, template_root)
            result["target_path"] = str(downloaded)
            result["valid"] = True
            result["reason"] = "git_clone"
        elif is_supported_template_url(source.source_url):
            downloaded = download_file(source.source_url, template_root.parent, timeout=timeout)
            result["target_path"] = str(downloaded)
            result["valid"] = validate_template_asset(downloaded)["valid"]
            result["reason"] = "file_download"
        else:
            response_links = discover_template_links(source.source_url, timeout=timeout)
            result["discovered_links"] = response_links[:20]
            if response_links:
                candidate = response_links[0]
                downloaded = download_file(candidate, template_root.parent, timeout=timeout)
                result["target_path"] = str(downloaded)
                result["valid"] = validate_template_asset(downloaded)["valid"]
                result["reason"] = "discovered_file"
            else:
                result["reason"] = "no_downloadable_link_found"
    except Exception as exc:  # pragma: no cover - network / environment dependent
        result["reason"] = str(exc)
    return result


def build_template_record_from_asset(
    asset_path: str | Path,
    source: SourceEntry,
    root: str | Path = "ppt_template_library",
) -> dict[str, Any]:
    paths = ensure_library_structure(root)
    asset = Path(asset_path)
    validation = validate_template_asset(asset)
    metadata: dict[str, Any] = {}
    if asset.is_file() and asset.suffix.lower() in {".pptx", ".potx"} and validation["valid"]:
        metadata = extract_pptx_metadata(asset)
    elif asset.is_dir():
        metadata = {
            "slide_count": None,
            "aspect_ratio": "",
            "contains_chart": False,
            "contains_table": False,
            "contains_timeline": False,
            "contains_process": False,
            "contains_map": False,
            "contains_infographic": False,
            "contains_images": False,
            "contains_icons": False,
            "contains_animations": False,
        }
    else:
        metadata = {
            "slide_count": None,
            "aspect_ratio": "",
            "contains_chart": False,
            "contains_table": False,
            "contains_timeline": False,
            "contains_process": False,
            "contains_map": False,
            "contains_infographic": False,
            "contains_images": False,
            "contains_icons": False,
            "contains_animations": False,
        }

    text = " ".join([source.name, source.notes, " ".join(source.tags), source.source_website])
    classification = infer_classification(text, source.tags)
    content_flags = infer_content_flags(text, source.tags)
    primary_color, secondary_color, dark_or_light = _color_pair(classification["colors"])
    signature = file_signature(asset)
    quality_score = max(60, min(100, source.score))
    design_score = min(100, quality_score - 5)
    usability_score = min(100, quality_score - 3)
    richness_bonus = sum(1 for value in content_flags.values() if value)
    overall_score = score_template(design_score, usability_score, quality_score, richness_bonus)
    template_id = f"ingested-{source.source_id}"
    preview_image = _build_preview_bundle(
        template_id=template_id,
        template_name=source.name,
        subtitle=f"{source.source_website} · {source.language}",
        score=overall_score,
        tags=source.tags,
        preview_dir=paths.preview / "templates",
        layout_types=infer_preview_layouts(text, source.tags),
    )
    record = default_template_record(
        template_id=template_id,
        template_name=source.name,
        source_url=source.source_url,
        source_website=source.source_website,
    )
    record.update(
        {
            "author": source.source_website,
            "license": source.license_hint,
            "commercial_use": source.commercial_use_hint,
            "modification_allowed": source.modification_allowed_hint,
            "download_date": datetime.utcnow().date().isoformat(),
            "file_format": asset.suffix.lower().lstrip(".") if asset.is_file() else "directory",
            "file_size": asset.stat().st_size if asset.is_file() else None,
            "slide_count": metadata.get("slide_count"),
            "aspect_ratio": metadata.get("aspect_ratio", ""),
            "language": source.language,
            "industry": _first_or_empty(classification["industries"], _first_or_empty(list(source.tags), "")),
            "scenario": _first_or_empty(classification["scenarios"], "presentation_template"),
            "style": _first_or_empty(classification["styles"], "modern"),
            "primary_color": primary_color,
            "secondary_color": secondary_color,
            "dark_or_light": dark_or_light,
            "contains_chart": metadata.get("contains_chart", False) or content_flags["contains_chart"],
            "contains_table": metadata.get("contains_table", False) or content_flags["contains_table"],
            "contains_timeline": metadata.get("contains_timeline", False) or content_flags["contains_timeline"],
            "contains_process": metadata.get("contains_process", False) or content_flags["contains_process"],
            "contains_map": metadata.get("contains_map", False) or content_flags["contains_map"],
            "contains_infographic": metadata.get("contains_infographic", False) or content_flags["contains_infographic"],
            "contains_images": metadata.get("contains_images", False) or content_flags["contains_images"],
            "contains_icons": metadata.get("contains_icons", False) or content_flags["contains_icons"],
            "contains_animations": metadata.get("contains_animations", False) or content_flags["contains_animations"],
            "quality_score": overall_score,
            "design_score": design_score,
            "usability_score": usability_score,
            "duplicate_hash": "|".join(filter(None, [signature["file_hash"], signature["image_hash"], signature["structure_hash"]])),
            "preview_image": preview_image,
            "local_path": str(asset),
            "metadata_json": {
                "source_type": source.source_type,
                "download_method": source.download_method,
                "programmatic_access": source.programmatic_access,
                "validation": validation,
                "signature": signature,
                "classification": classification,
                "content_flags": content_flags,
                "scaffold": False,
                "tags": list(source.tags),
            },
            "status": "downloaded" if validation["valid"] else "needs_review",
        }
    )
    return record


def ingest_downloaded_asset(
    asset_path: str | Path,
    source: SourceEntry,
    root: str | Path = "ppt_template_library",
) -> dict[str, Any]:
    record = build_template_record_from_asset(asset_path, source, root=root)
    save_template_catalog([record], root)
    return record


def deduplicate_template_catalog(df: pd.DataFrame) -> pd.DataFrame:
    if df.empty:
        return df
    work = df.copy()
    if "metadata_json" in work.columns:
        work["metadata_text"] = work["metadata_json"].astype(str)
    else:
        work["metadata_text"] = ""
    sort_columns = [column for column in ["quality_score", "design_score", "usability_score", "template_name"] if column in work.columns]
    ascending = [False, False, False, True][: len(sort_columns)]
    if sort_columns:
        work = work.sort_values(sort_columns, ascending=ascending)
    return work.drop_duplicates(subset=["duplicate_hash"], keep="first").drop(columns=["metadata_text"], errors="ignore")


def search_sources(root: str | Path = "ppt_template_library", filters: SearchFilters | None = None) -> pd.DataFrame:
    filters = filters or SearchFilters()
    df = load_source_catalog(root)
    return filter_source_dataframe(df, filters)


def search_templates(
    root: str | Path = "ppt_template_library",
    query: str = "",
    industries: tuple[str, ...] = (),
    scenarios: tuple[str, ...] = (),
    styles: tuple[str, ...] = (),
    colors: tuple[str, ...] = (),
    licenses: tuple[str, ...] = (),
    statuses: tuple[str, ...] = (),
    min_score: int = 0,
    commercial_only: bool = False,
    scaffold_only: bool = False,
) -> pd.DataFrame:
    df = load_template_catalog(root)
    return filter_template_dataframe(
        df,
        query=query,
        industries=industries,
        scenarios=scenarios,
        styles=styles,
        colors=colors,
        licenses=licenses,
        statuses=statuses,
        min_score=min_score,
        commercial_only=commercial_only,
        scaffold_only=scaffold_only,
    )


def search_slides(root: str | Path = "ppt_template_library", filters: PageSearchFilters | None = None) -> pd.DataFrame:
    df = load_slide_catalog(root)
    filters = filters or PageSearchFilters()
    return filter_slide_dataframe(df, filters)


def search_components(
    root: str | Path = "ppt_template_library",
    query: str = "",
    component_types: tuple[str, ...] = (),
    layout_types: tuple[str, ...] = (),
    styles: tuple[str, ...] = (),
    min_score: int = 0,
) -> pd.DataFrame:
    df = load_component_catalog(root)
    return filter_component_dataframe(
        df,
        query=query,
        component_types=component_types,
        layout_types=layout_types,
        styles=styles,
        min_score=min_score,
    )


def load_library_preview_counts(root: str | Path = "ppt_template_library") -> dict[str, Any]:
    manifest = build_library_manifest(root)
    source_df = load_source_catalog(root)
    template_df = load_template_catalog(root)
    slide_df = load_slide_catalog(root)
    component_df = load_component_catalog(root)
    return {
        **manifest,
        "source_count": int(len(source_df)),
        "template_count": int(len(template_df)),
        "slide_count": int(len(slide_df)),
        "component_count": int(len(component_df)),
    }

